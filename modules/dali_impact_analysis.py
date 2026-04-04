"""DALI impact analysis and Data4Sec inventory enrichment pipeline.

This module orchestrates monitored-UID extraction from DALI, inventory enrichment
for Gen2 servers, beneficiary-based discovery, and report artifact generation.
"""

import argparse
import base64
import csv
import gzip
import ipaddress
import json
import logging
import math
import os
import random
import re
import time
import uuid
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import urljoin
from xml.etree import ElementTree as ET
from xml.sax.saxutils import escape

from config import QUERY_CONFIG
from d4s_client import Data4secClient

log = logging.getLogger(__name__)

HEADER_ALIASES = {
    "kear": "kear",
    "kear_id": "kear",
    "kearid": "kear",
    "uid": "uid",
    "program": "program",
    "programme": "program",
    "network": "network",
    "net": "network",
    "taken": "taken",
    "is_taken": "taken",
}


IMPACT_DEFAULT_PARAMS = {
    "ciLabel": "Application",
    "attributeName": "uid",
    "matchType": "equals",
    "direction": "to",
    "relationship": [
        "CHANGES",
        "IS_ASSIGNED_TO",
        "IS_CONTAINED_BY",
        "IS_GRANTED_TO",
        "IS_HOSTED_BY",
        "IS_LOCATED_BY",
        "IS_MANAGED_BY",
        "IS_MEMBER_OF",
        "IS_USED_BY",
        "USE",
        "USE_STORAGE",
        "MANAGE_RESOURCE",
        "IS_PROVIDED_BY",
        "IS_CONNECTED_TO",
        "COMPOSED_BY",
        "CLUSTER_CONTAINS",
    ],
    "impactedCis": "Server",
    "status": "In use",
    "reliability": "false",
    "criticality": ["Critical", "High", "Medium", "Low", "Unknown"],
    "includeLiveSources": "true",
    "zones": ["EUR", "ASIA", "AMER", "BCO", "UK", "Unknown"],
    "environments": ["Production", "Not in production"],
    "excludeDuplicates": "true",
    "boost": "false",
    "includeGTSInfra": "true",
    "includeCount": "true",
    "skip": "0",
}


RETRY_STATUSES = {429, 500, 502, 503, 504}

RAW_FILTER_COLUMN_PAIRS: List[Tuple[Optional[str], str, str]] = [
    ("FILTER_VALUE_environment", "F_FILTER_PRD_ENV", "FILTER_PRD_ENV"),
    ("FILTER_VALUE_os_name", "F_FILTER_OS_NAME", "FILTER_OS_NAME"),
    ("FILTER_VALUE_server.status", "F_FILTER_SERVER_STATUS", "FILTER_SERVER_STATUS"),
    ("FILTER_VALUE_cloud_type", "F_FILTER_CLOUD_TYPE_NOT_TAKEN", "FILTER_CLOUD_TYPE_NOT_TAKEN"),
    ("FILTER_VALUE_main_application", "F_FILTER_MAIN_APP_NOT_TAKEN", "FILTER_MAIN_APP_NOT_TAKEN"),
    ("FILTER_VALUE_domain", "F_FILTER_DOMAIN", "FILTER_DOMAIN_NOT_TAKEN"),
    ("FILTER_VALUE_typology", "F_FILTER_TYPOLOGY_NOT_TAKEN", "FILTER_TYPOLOGY_NOT_TAKEN"),
]

MARLEY_ENRICHMENT_MAPPING_TABLE: List[Tuple[str, str, str, str, str]] = [
    ("get_marley_gen2_by_uuid", "app_info.kear_uuid", "UID REL", "uid", "keep"),
    ("get_marley_gen2_by_uuid", "app_info.app_name", "NAME REL", "name", "keep"),
    ("get_marley_gen2_by_uuid", "app_info.app_name", "SHORT LABEL REL", "short_label", "keep"),
    ("get_marley_gen2_by_uuid", "", "ASA REL", "asa", "keep"),
    ("get_marley_gen2_by_uuid", "app_info.app_id", "IRT CODE REL", "irt_code", "keep"),
    ("get_marley_gen2_by_uuid", "", "IAPPLI CODE REL", "iappli_code", "keep"),
    ("get_marley_gen2_by_uuid", "", "TRIGRAM REL", "trigram", "keep"),
    ("get_marley_gen2_by_uuid", "app_info.kear_library", "DSI REL", "dsi", "keep"),
    ("get_marley_gen2_by_uuid", "app_info.service_line_name", "APPLICATION MANAGEMENT RC REL", "application_management_rc", "keep"),
    (
        "get_marley_gen2_by_uuid",
        "",
        "APPLICATION DEVELOPMENT MANAGER REL",
        "application_development_manager",
        "keep",
    ),
    ("get_marley_gen2_by_uuid", "app_info.app_id", "MAIN APPLICATION", "main_application", "keep"),
    ("get_marley_gen2_by_uuid", "app_info.env", "ENVIRONMENT", "environment", "keep"),
    ("get_marley_gen2_by_uuid", "uuid", "SERVER UID", "server_uid", "keep"),
    ("get_marley_gen2_by_uuid", "ocs_name", "HOSTNAME", "hostname", "keep"),
    ("get_marley_gen2_by_uuid", "status", "DALI STATUS", "usage", "keep"),
    ("get_marley_gen2_by_uuid", "status", "STATUS", "status", "keep"),
    ("get_marley_gen2_by_uuid", "ocs_name", "USUAL NAME", "usual_name", "keep"),
    ("get_inv_by_account", "hostname", "FRIENDLY NAME", "friendly_name", "keep"),
    ("get_marley_gen2_by_uuid", "", "DNS NAME", "dns_name", "keep"),
    ("get_marley_gen2_by_uuid", "", "TYPOLOGY", "typology", "keep"),
    ("get_marley_gen2_by_uuid", "Gen 2", "CLOUD TYPE", "cloud_type", "keep"),
    ("get_marley_gen2_by_uuid", "", "SERVICE OFFER", "service_offer", "keep"),
    ("get_marley_gen2_by_uuid", "os_name", "OS NAME", "os_name", "keep"),
    ("get_marley_gen2_by_uuid", "", "OS RELEASE", "os_release", "keep"),
    ("get_marley_gen2_by_uuid", "", "VRF NAME", "vrf_name", "keep"),
    ("get_marley_gen2_by_uuid", "", "SILO", "silo", "keep"),
    ("get_marley_gen2_by_uuid", "", "UPDATED BY", "updated_by", "keep"),
    ("get_marley_gen2_by_uuid", "beneficiary", "INV_Beneficiary_Account", "INV_Beneficiary_Account", "keep"),
    ("get_marley_gen2_by_uuid", "owner_app_name", "INV_Owner_Account", "INV_Owner_Account", "keep"),
    ("get_inv_by_account", "ocs_name", "INV_ocs_name", "INV_ocs_name", "keep"),
    ("get_inv_by_account", "status", "INV_status", "INV_status", "keep"),
    ("get_inv_by_account", "hostname", "INV_hostname", "INV_hostname", "keep"),
    ("FILTRED", "IPLIST", "network", "", "put same value for network from IPLIST column"),
]


def _response_error_details(status_code: int, body: str, max_len: int = 500) -> str:
    compact = " ".join(str(body or "").split())
    if len(compact) > max_len:
        compact = compact[:max_len] + "..."
    return f"HTTP {status_code} | response={compact or '<empty>'}"


def normalize_header_name(name: Optional[str]) -> str:
    if name is None:
        return ""
    normalized = str(name).replace("\ufeff", "").strip().lower()
    normalized = normalized.replace(" ", "_").replace("-", "_")
    while "__" in normalized:
        normalized = normalized.replace("__", "_")
    return HEADER_ALIASES.get(normalized, normalized)


def load_env_file(env_file: str = ".env") -> None:
    path = Path(env_file)
    if not path.is_file():
        return
    with open(path, "r", encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            key, value = line.split("=", 1)
            key = key.strip()
            value = value.strip().strip('"').strip("'")
            if key and key not in os.environ:
                os.environ[key] = value


def resolve_verify_ca() -> Any:
    """Resolve TLS verification strategy: VERIFY_CA > sg_cacert_file > default True."""
    value = os.getenv("VERIFY_CA")
    if value is not None and str(value).strip() != "":
        lowered = str(value).strip().lower()
        if lowered in {"true", "1", "yes", "on"}:
            return True
        if lowered in {"false", "0", "no", "off"}:
            return False
        return str(value).strip()

    try:
        from sg_cacert_file import get_cacert_path

        ca_path = get_cacert_path()
        log.info("Using CA bundle from sg_cacert_file: %s", ca_path)
        return ca_path
    except Exception:
        return True


def parse_positive_int(name: str, value: Optional[str]) -> Optional[int]:
    if value is None or str(value).strip() == "":
        return None
    parsed = int(value)
    if parsed <= 0:
        raise ValueError(f"{name} must be > 0")
    return parsed


def detect_csv_delimiter(csv_file: str, default: str = ",") -> str:
    with open(csv_file, "r", encoding="utf-8", newline="") as handle:
        sample = handle.read(4096)
    if not sample.strip():
        return default
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;")
        if dialect.delimiter in {",", ";"}:
            return dialect.delimiter
    except csv.Error:
        pass
    return ";" if sample.count(";") > sample.count(",") else default


def read_csv_with_original_headers(csv_file: str) -> Tuple[List[str], List[Dict[str, str]]]:
    delimiter = detect_csv_delimiter(csv_file)
    log.info("Detected generic CSV delimiter '%s' for %s", delimiter, csv_file)
    with open(csv_file, "r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle, delimiter=delimiter)
        headers = list(reader.fieldnames or [])
        rows: List[Dict[str, str]] = []
        for raw in reader:
            row: Dict[str, str] = {}
            for header in headers:
                value = raw.get(header, "") if raw is not None else ""
                row[header] = "" if value is None else str(value)
            if any(str(v).strip() for v in row.values()):
                rows.append(row)
    return headers, rows


class DaliImpactAnalysisClient:
    def __init__(self) -> None:
        self.base_url = (os.getenv("DALI_BASE_URL") or "").rstrip("/")
        self.token_url = (os.getenv("SGMARKET_TOKEN_URL") or "").strip()
        self.client_id = (os.getenv("SGCONNECT_CLIENT_ID") or "").strip()
        self.client_secret = (os.getenv("SGCONNECT_CLIENT_SECRET") or "").strip()
        self.scopes = (os.getenv("SGCONNECT_SCOPES") or "").strip()
        self.dali_client_id = (os.getenv("DALI_CLIENT_ID") or "").strip()
        self.dali_client_id_header = (os.getenv("DALI_CLIENT_ID_HEADER") or "x-client-id").strip()
        self.verify = resolve_verify_ca()
        self._token: Optional[str] = None
        self._token_expiry_epoch: float = 0.0

    def _validate_settings(self) -> None:
        missing = []
        for key, value in {
            "DALI_BASE_URL": self.base_url,
            "SGMARKET_TOKEN_URL": self.token_url,
            "SGCONNECT_CLIENT_ID": self.client_id,
            "SGCONNECT_CLIENT_SECRET": self.client_secret,
            "SGCONNECT_SCOPES": self.scopes,
        }.items():
            if not value:
                missing.append(key)
        if missing:
            raise ValueError(f"Missing DALI settings in .env: {', '.join(missing)}")

    def _basic_auth_header(self) -> str:
        raw = f"{self.client_id}:{self.client_secret}".encode("utf-8")
        return "Basic " + base64.b64encode(raw).decode("ascii")

    def fetch_sg_token(self) -> Tuple[str, int]:
        self._validate_settings()
        import requests

        headers = {
            "Authorization": self._basic_auth_header(),
            "Content-Type": "application/x-www-form-urlencoded",
            "Accept": "application/json",
        }
        payload = {
            "grant_type": "client_credentials",
            "scope": self.scopes,
        }
        response = requests.post(self.token_url, data=payload, headers=headers, timeout=30, verify=self.verify)
        response.raise_for_status()
        body = response.json()
        token = body.get("access_token")
        expires_in = int(body.get("expires_in", 3600))
        if not token:
            raise RuntimeError("No access_token found in OAuth2 response")
        return token, expires_in

    def get_bearer_token(self, force_refresh: bool = False) -> str:
        now = time.time()
        if not force_refresh and self._token and now < (self._token_expiry_epoch - 30):
            return self._token
        token, expires_in = self.fetch_sg_token()
        self._token = token
        self._token_expiry_epoch = now + max(60, expires_in)
        return token

    def dali_headers(self, force_refresh: bool = False) -> Dict[str, str]:
        headers = {
            "Accept": "application/json",
            "Authorization": f"Bearer {self.get_bearer_token(force_refresh=force_refresh)}",
        }
        if self.dali_client_id:
            headers[self.dali_client_id_header] = self.dali_client_id
        return headers

    def get_json(self, endpoint: str, params: Dict[str, Any], timeout_s: int = 60, retries: int = 4) -> Dict[str, Any]:
        import requests

        url = urljoin(f"{self.base_url}/", endpoint.lstrip("/"))
        last_exc: Optional[Exception] = None
        uid = params.get("attributeValue")

        for attempt in range(retries + 1):
            force_refresh = attempt > 0
            try:
                log.info(
                    "DALI GET request attempt=%s/%s url=%s params=%s",
                    attempt + 1,
                    retries + 1,
                    url,
                    params,
                )
                response = requests.get(
                    url,
                    params=params,
                    headers=self.dali_headers(force_refresh=force_refresh),
                    timeout=timeout_s,
                    verify=self.verify,
                )
                status_code = int(response.status_code)

                if status_code in {401, 403}:
                    self._token = None
                    self._token_expiry_epoch = 0
                    if attempt < retries:
                        continue

                if status_code in RETRY_STATUSES and attempt < retries:
                    delay = (2**attempt) + random.uniform(0, 0.5)
                    log.warning("DALI transient status=%s for uid=%s, retry in %.2fs", status_code, uid, delay)
                    time.sleep(delay)
                    continue

                if status_code >= 400:
                    details = _response_error_details(status_code=status_code, body=response.text)
                    raise RuntimeError(f"DALI request failed for uid={params.get('attributeValue')}: {details}")

                response.raise_for_status()
                payload = response.json()
                result_count = 0
                if isinstance(payload, dict):
                    result = payload.get("result")
                    if isinstance(result, list):
                        result_count = len(result)
                log.info(
                    "DALI GET response status=%s url=%s result_edges=%s count=%s",
                    status_code,
                    url,
                    result_count,
                    payload.get("count") if isinstance(payload, dict) else "n/a",
                )
                return payload
            except requests.HTTPError as exc:
                status_code = int(exc.response.status_code) if exc.response is not None else -1
                details = _response_error_details(
                    status_code=status_code,
                    body=exc.response.text if exc.response is not None else str(exc),
                )
                raise RuntimeError(f"DALI request failed for uid={uid}: {details}") from exc
            except requests.RequestException as exc:
                last_exc = exc
                if attempt < retries:
                    delay = (2**attempt) + random.uniform(0, 0.5)
                    log.warning(
                        "DALI request error for uid=%s on attempt %s/%s: %s; retry in %.2fs",
                        uid,
                        attempt + 1,
                        retries + 1,
                        exc,
                        delay,
                    )
                    time.sleep(delay)
                    continue
                break

        raise RuntimeError(f"DALI request failed after retries for uid={uid}: {last_exc}")


def read_headers_mapping(headers_file: str) -> List[Tuple[str, str]]:
    mappings: List[Tuple[str, str]] = []
    delimiter = detect_csv_delimiter(headers_file)
    log.info("Detected headers delimiter '%s' for %s", delimiter, headers_file)
    with open(headers_file, "r", encoding="utf-8", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for row in reader:
            display_name = str(row[0]).strip() if len(row) > 0 and row[0] else ""
            dali_attr = str(row[1]).strip() if len(row) > 1 and row[1] else ""
            if display_name and dali_attr:
                mappings.append((display_name, dali_attr))
    if not mappings:
        raise ValueError(f"No valid mappings found in {headers_file}")
    return mappings


def read_monitored_kears(monitored_file: str) -> List[Dict[str, str]]:
    delimiter = detect_csv_delimiter(monitored_file)
    log.info("Detected monitored_kears delimiter '%s' for %s", delimiter, monitored_file)
    with open(monitored_file, "r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle, delimiter=delimiter)
        raw_headers = reader.fieldnames or []
        headers = [normalize_header_name(h) for h in raw_headers]
        required = ["program", "network", "taken"]
        missing = [col for col in required if col not in headers]
        if "kear" not in headers and "uid" not in headers:
            missing = ["kear_or_uid"] + missing
        if missing:
            raise ValueError(
                f"Missing required columns in {monitored_file}: {', '.join(missing)} | detected_headers={headers}"
            )

        rows: List[Dict[str, str]] = []
        for raw in reader:
            normalized: Dict[str, str] = {}
            for key, value in raw.items():
                normalized_key = normalize_header_name(key)
                normalized[normalized_key] = str(value).strip() if value is not None else ""
            uid = normalized.get("uid") or normalized.get("kear")
            if not uid:
                continue
            rows.append(
                {
                    "uid": uid,
                    "kear": normalized.get("kear", uid),
                    "program": normalized.get("program", ""),
                    "network": normalized.get("network", ""),
                    "taken": normalized.get("taken", ""),
                    "short_label": normalized.get("short_label", ""),
                    "slide": normalized.get("slide", ""),
                }
            )
    if not rows:
        raise ValueError(f"No monitored KEAR rows found in {monitored_file}")
    return rows


def read_filters_conf(filters_file: str) -> Dict[str, str]:
    filters: Dict[str, str] = {}
    with open(filters_file, "r", encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.strip()
            if not line or line.startswith("#"):
                continue

            if "=" in line:
                key, value = line.split("=", 1)
            elif "," in line:
                key, value = line.split(",", 1)
            else:
                log.warning("Ignoring invalid filter line (expected key=value or key,value): %s", line)
                continue

            key = key.strip()
            value = value.strip()
            if key:
                filters[key] = value
    return filters


def node_properties_to_dict(node: Any) -> Dict[str, Any]:
    if not isinstance(node, dict):
        return {}
    props = node.get("properties")
    if isinstance(props, dict):
        return props
    if not isinstance(props, list):
        return {}
    out: Dict[str, Any] = {}
    for p in props:
        if not isinstance(p, dict):
            continue
        name = p.get("name")
        if not name:
            continue
        out[str(name)] = p.get("value")
    return out


def _extract_server_uid_from_edge(edge: Dict[str, Any]) -> str:
    """Extract Server UID from nodes labeled `Server` in a DALI edge."""
    for node_key in ("leading_node", "trailing_node"):
        node = edge.get(node_key)
        if not isinstance(node, dict):
            continue
        labels = node.get("labels")
        normalized_labels = {str(label).strip().lower() for label in labels} if isinstance(labels, list) else set()
        if "server" not in normalized_labels:
            continue
        props = node_properties_to_dict(node)
        uid = props.get("uid")
        if uid is None:
            continue
        value = str(uid).strip()
        if value:
            return value
    return ""


def _node_has_label(node: Any, expected_label: str) -> bool:
    if not isinstance(node, dict):
        return False
    labels = node.get("labels")
    normalized_labels = {str(label).strip().lower() for label in labels} if isinstance(labels, list) else set()
    return str(expected_label).strip().lower() in normalized_labels


def _resolve_edge_mapping_value(edge: Dict[str, Any], dali_attr: str, base_row: Dict[str, Any]) -> Any:
    attr = str(dali_attr or "").strip()
    if not attr:
        return ""

    lower_attr = attr.lower()
    leading_node = edge.get("leading_node")
    trailing_node = edge.get("trailing_node")
    lead = node_properties_to_dict(leading_node)
    trail = node_properties_to_dict(trailing_node)

    scoped_value: Any = None
    scoped_attr = attr

    if "." in attr:
        scope, scoped_attr = attr.split(".", 1)
        scope = scope.strip().lower()
        scoped_attr = scoped_attr.strip()

        if scope == "leading":
            scoped_value = lead.get(scoped_attr)
        elif scope == "trailing":
            scoped_value = trail.get(scoped_attr)
        elif scope == "server":
            for node, props in ((leading_node, lead), (trailing_node, trail)):
                if _node_has_label(node, "server"):
                    scoped_value = props.get(scoped_attr)
                    if scoped_value is not None and (not isinstance(scoped_value, str) or scoped_value.strip()):
                        break
        elif scope == "application":
            for node, props in ((leading_node, lead), (trailing_node, trail)):
                if _node_has_label(node, "application"):
                    scoped_value = props.get(scoped_attr)
                    if scoped_value is not None and (not isinstance(scoped_value, str) or scoped_value.strip()):
                        break

        if scoped_value is not None and (not isinstance(scoped_value, str) or scoped_value.strip()):
            return scoped_value
        if scoped_value is not None:
            return scoped_value
        lower_attr = scoped_attr.lower()

    raw_value = lead.get(scoped_attr)
    if raw_value is None or (isinstance(raw_value, str) and not raw_value.strip()):
        raw_value = trail.get(scoped_attr)
    if raw_value is None and lower_attr in {"uid", "application_uid", "app_uid"}:
        raw_value = base_row.get("uid", "")
    return raw_value


def _normalize_cell_value(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, list):
        return ", ".join(str(x) for x in value if str(x).strip())
    if isinstance(value, dict):
        return json.dumps(value, ensure_ascii=False)
    return str(value)


def _parse_filter_tokens(filters: Optional[Dict[str, str]], key: str) -> List[str]:
    if not filters:
        return []
    raw = filters.get(key)
    if raw is None:
        raw = filters.get(key.lower())
    if not raw:
        return []
    return [chunk.strip().upper() for chunk in raw.split(",") if chunk.strip()]


def _parse_filter_bool(filters: Optional[Dict[str, str]], key: str, default: bool) -> bool:
    if not filters:
        return default
    raw = filters.get(key)
    if raw is None:
        raw = filters.get(key.lower())
    if raw is None or str(raw).strip() == "":
        return default
    lowered = str(raw).strip().lower()
    if lowered in {"true", "1", "yes", "y", "on"}:
        return True
    if lowered in {"false", "0", "no", "n", "off"}:
        return False
    return default


def _property_value_from_nodes(
    lead: Dict[str, Any],
    trail: Dict[str, Any],
    property_name: str,
    leading_node: Optional[Dict[str, Any]] = None,
    trailing_node: Optional[Dict[str, Any]] = None,
) -> str:
    attr = str(property_name or "").strip()
    if not attr:
        return ""

    if "." in attr:
        scope, scoped_attr = attr.split(".", 1)
        scope = scope.strip().lower()
        scoped_attr = scoped_attr.strip()

        if scope == "leading":
            return _normalize_cell_value(lead.get(scoped_attr))
        if scope == "trailing":
            return _normalize_cell_value(trail.get(scoped_attr))
        if scope in {"server", "application"}:
            expected_label = scope
            for node, props in ((leading_node, lead), (trailing_node, trail)):
                if _node_has_label(node, expected_label):
                    value = props.get(scoped_attr)
                    if value is not None and (not isinstance(value, str) or value.strip()):
                        return _normalize_cell_value(value)
            # Fallback when labels are unavailable in provided nodes.
            value = lead.get(scoped_attr)
            if value is None or (isinstance(value, str) and not value.strip()):
                value = trail.get(scoped_attr)
            return _normalize_cell_value(value)

        # Unknown scope: fallback to raw attribute name without scope.
        attr = scoped_attr

    value = lead.get(attr)
    if value is None or (isinstance(value, str) and not value.strip()):
        value = trail.get(attr)
    return _normalize_cell_value(value)


def _contains_any_token(value: str, tokens: List[str]) -> bool:
    if not tokens:
        return False
    normalized = str(value or "").upper()
    return any(token in normalized for token in tokens)


def _matches_exact_token(value: str, tokens: List[str]) -> bool:
    if not tokens:
        return False
    normalized = str(value or "").upper().strip()
    if normalized in tokens:
        return True
    parts = [part.strip() for part in normalized.split(",") if part.strip()]
    return any(part in tokens for part in parts)


def _edge_matches_filters(
    lead: Dict[str, Any],
    trail: Dict[str, Any],
    filters: Optional[Dict[str, str]],
    leading_node: Optional[Dict[str, Any]] = None,
    trailing_node: Optional[Dict[str, Any]] = None,
    row: Optional[Dict[str, Any]] = None,
) -> bool:
    env_tokens = _parse_filter_tokens(filters, "FILTER_PRD_ENV")
    if env_tokens:
        environment = _property_value_from_nodes(lead, trail, "environment", leading_node=leading_node, trailing_node=trailing_node)
        if not _contains_any_token(environment, env_tokens):
            return False

    os_tokens = _parse_filter_tokens(filters, "FILTER_OS_NAME")
    if os_tokens:
        os_name = _property_value_from_nodes(lead, trail, "os_name", leading_node=leading_node, trailing_node=trailing_node)
        if not _matches_exact_token(os_name, os_tokens):
            return False

    server_status_tokens = _parse_filter_tokens(filters, "FILTER_SERVER_STATUS")
    if server_status_tokens:
        server_status = _normalize_cell_value((row or {}).get("Server Status", ""))
        if not server_status:
            server_status = _property_value_from_nodes(lead, trail, "server.status", leading_node=leading_node, trailing_node=trailing_node)
        if not _matches_exact_token(server_status, server_status_tokens):
            return False

    cloud_type_not_taken = _parse_filter_tokens(filters, "FILTER_CLOUD_TYPE_NOT_TAKEN")
    if cloud_type_not_taken:
        cloud_type = _property_value_from_nodes(lead, trail, "cloud_type", leading_node=leading_node, trailing_node=trailing_node)
        if _contains_any_token(cloud_type, cloud_type_not_taken):
            return False

    main_app_not_taken = _parse_filter_tokens(filters, "FILTER_MAIN_APP_NOT_TAKEN")
    if main_app_not_taken:
        main_application = _property_value_from_nodes(lead, trail, "main_application", leading_node=leading_node, trailing_node=trailing_node)
        if _contains_any_token(main_application, main_app_not_taken):
            return False

    typology_not_taken = _parse_filter_tokens(filters, "FILTER_TYPOLOGY_NOT_TAKEN")
    if typology_not_taken:
        typology = _property_value_from_nodes(lead, trail, "typology", leading_node=leading_node, trailing_node=trailing_node)
        if _contains_any_token(typology, typology_not_taken):
            return False

    domain_not_taken = _parse_filter_tokens(filters, "FILTER_DOMAIN_NOT_TAKEN")
    if domain_not_taken:
        domain = _property_value_from_nodes(lead, trail, "dns_name", leading_node=leading_node, trailing_node=trailing_node)
        if _contains_any_token(domain, domain_not_taken):
            return False

    return True


def _raw_filter_debug_columns(
    lead: Dict[str, Any],
    trail: Dict[str, Any],
    filters: Optional[Dict[str, str]],
    leading_node: Optional[Dict[str, Any]] = None,
    trailing_node: Optional[Dict[str, Any]] = None,
    row: Optional[Dict[str, Any]] = None,
) -> Dict[str, str]:
    env_value = _property_value_from_nodes(lead, trail, "environment", leading_node=leading_node, trailing_node=trailing_node)
    os_value = _property_value_from_nodes(lead, trail, "os_name", leading_node=leading_node, trailing_node=trailing_node)
    server_status_value = _normalize_cell_value((row or {}).get("Server Status", ""))
    if not server_status_value:
        server_status_value = _property_value_from_nodes(lead, trail, "server.status", leading_node=leading_node, trailing_node=trailing_node)
    cloud_value = _property_value_from_nodes(lead, trail, "cloud_type", leading_node=leading_node, trailing_node=trailing_node)
    main_app_value = _property_value_from_nodes(lead, trail, "main_application", leading_node=leading_node, trailing_node=trailing_node)
    domain_value = _property_value_from_nodes(lead, trail, "dns_name", leading_node=leading_node, trailing_node=trailing_node)
    typology_value = _property_value_from_nodes(lead, trail, "typology", leading_node=leading_node, trailing_node=trailing_node)

    env_tokens = _parse_filter_tokens(filters, "FILTER_PRD_ENV")
    os_tokens = _parse_filter_tokens(filters, "FILTER_OS_NAME")
    server_status_tokens = _parse_filter_tokens(filters, "FILTER_SERVER_STATUS")
    cloud_tokens = _parse_filter_tokens(filters, "FILTER_CLOUD_TYPE_NOT_TAKEN")
    main_app_tokens = _parse_filter_tokens(filters, "FILTER_MAIN_APP_NOT_TAKEN")
    domain_tokens = _parse_filter_tokens(filters, "FILTER_DOMAIN_NOT_TAKEN")
    typology_tokens = _parse_filter_tokens(filters, "FILTER_TYPOLOGY_NOT_TAKEN")

    env_ok = True if not env_tokens else _contains_any_token(env_value, env_tokens)
    os_ok = True if not os_tokens else _matches_exact_token(os_value, os_tokens)
    server_status_ok = True if not server_status_tokens else _matches_exact_token(server_status_value, server_status_tokens)
    cloud_ok = not _contains_any_token(cloud_value, cloud_tokens) if cloud_tokens else True
    main_app_ok = not _contains_any_token(main_app_value, main_app_tokens) if main_app_tokens else True
    domain_ok = not _contains_any_token(domain_value, domain_tokens) if domain_tokens else True
    typology_ok = not _contains_any_token(typology_value, typology_tokens) if typology_tokens else True
    excluded_ok = _normalize_lookup_value((row or {}).get("F_Excluded", "N")) != "Y"

    return {
        "FILTER_VALUE_environment": env_value,
        "F_FILTER_PRD_ENV": "Y" if env_ok else "N",
        "FILTER_VALUE_os_name": os_value,
        "F_FILTER_OS_NAME": "Y" if os_ok else "N",
        "FILTER_VALUE_server.status": server_status_value,
        "F_FILTER_SERVER_STATUS": "Y" if server_status_ok else "N",
        "FILTER_VALUE_cloud_type": cloud_value,
        "F_FILTER_CLOUD_TYPE_NOT_TAKEN": "Y" if cloud_ok else "N",
        "FILTER_VALUE_main_application": main_app_value,
        "F_FILTER_MAIN_APP_NOT_TAKEN": "Y" if main_app_ok else "N",
        "FILTER_VALUE_domain": domain_value,
        "F_FILTER_DOMAIN": "Y" if domain_ok else "N",
        "FILTER_VALUE_typology": typology_value,
        "F_FILTER_TYPOLOGY_NOT_TAKEN": "Y" if typology_ok else "N",
        "F_FILTER_ALL": "Y" if all([env_ok, os_ok, server_status_ok, cloud_ok, main_app_ok, domain_ok, typology_ok, excluded_ok]) else "N",
    }


def _enrich_filter_columns_from_enrich_row(
    row: Dict[str, Any],
    filters: Optional[Dict[str, str]],
    servers_to_exclude: Optional[List[str]] = None,
) -> Dict[str, str]:
    server_status_value = _normalize_cell_value(_get_row_value_by_candidates(row, ["DALI [CI] Server Status"]))
    os_value = _normalize_cell_value(_get_row_value_by_candidates(row, ["DALI [CI] OS NAME"]))
    main_app_value = _normalize_cell_value(_get_row_value_by_candidates(row, ["DALI [CI] MAIN APPLICATION"]))
    env_value = _normalize_cell_value(_get_row_value_by_candidates(row, ["INV_Beneficiary_Account_ENV"]))
    cloud_value = _normalize_cell_value(_get_row_value_by_candidates(row, ["DALI [CI] CLOUD TYPE"]))
    domain_value = _normalize_cell_value(_get_row_value_by_candidates(row, ["DALI [CI] DNS NAME"]))
    typology_value = _normalize_cell_value(_get_row_value_by_candidates(row, ["DALI [CI] TYPOLOGY"]))

    excluded_lookup = {_normalize_hostname_for_compare(value) for value in (servers_to_exclude or []) if _normalize_hostname_for_compare(value)}
    hostname_candidates = [
        _get_row_value_by_candidates(row, ["DALI [CI] HOSTNAME", "INV_ocs_name", "INV_hostname"]),
        _get_row_value_by_candidates(row, ["DALI [CI] USUAL NAME"]),
        _get_row_value_by_candidates(row, ["DALI [CI] FRIENDLY NAME"]),
    ]
    excluded_hit = any(_normalize_hostname_for_compare(candidate) in excluded_lookup for candidate in hostname_candidates if _normalize_hostname_for_compare(candidate))

    env_tokens = _parse_filter_tokens(filters, "FILTER_PRD_ENV")
    os_tokens = _parse_filter_tokens(filters, "FILTER_OS_NAME")
    server_status_tokens = _parse_filter_tokens(filters, "FILTER_SERVER_STATUS")
    cloud_tokens = _parse_filter_tokens(filters, "FILTER_CLOUD_TYPE_NOT_TAKEN")
    main_app_tokens = _parse_filter_tokens(filters, "FILTER_MAIN_APP_NOT_TAKEN")
    domain_tokens = _parse_filter_tokens(filters, "FILTER_DOMAIN_NOT_TAKEN")
    typology_tokens = _parse_filter_tokens(filters, "FILTER_TYPOLOGY_NOT_TAKEN")

    env_ok = True if not env_tokens else _contains_any_token(env_value, env_tokens)
    os_ok = True if not os_tokens else _matches_exact_token(os_value, os_tokens)
    server_status_ok = True if not server_status_tokens else _matches_exact_token(server_status_value, server_status_tokens)
    cloud_ok = not _contains_any_token(cloud_value, cloud_tokens) if cloud_tokens else True
    main_app_ok = not _contains_any_token(main_app_value, main_app_tokens) if main_app_tokens else True
    domain_ok = not _contains_any_token(domain_value, domain_tokens) if domain_tokens else True
    typology_ok = not _contains_any_token(typology_value, typology_tokens) if typology_tokens else True
    excluded_ok = not excluded_hit

    return {
        "FILTER_VALUE_server.status": server_status_value,
        "F_FILTER_SERVER_STATUS": "Y" if server_status_ok else "N",
        "FILTER_VALUE_os_name": os_value,
        "F_FILTER_OS_NAME": "Y" if os_ok else "N",
        "FILTER_VALUE_main_application": main_app_value,
        "F_FILTER_MAIN_APP_NOT_TAKEN": "Y" if main_app_ok else "N",
        "FILTER_VALUE_environment": env_value,
        "F_FILTER_PRD_ENV": "Y" if env_ok else "N",
        "FILTER_VALUE_cloud_type": cloud_value,
        "F_FILTER_CLOUD_TYPE_NOT_TAKEN": "Y" if cloud_ok else "N",
        "FILTER_VALUE_domain": domain_value,
        "F_FILTER_DOMAIN": "Y" if domain_ok else "N",
        "FILTER_VALUE_typology": typology_value,
        "F_FILTER_TYPOLOGY_NOT_TAKEN": "Y" if typology_ok else "N",
        "F_Excluded": "Y" if excluded_hit else "N",
        "F_FILTER_ALL": "Y" if all([env_ok, os_ok, server_status_ok, cloud_ok, main_app_ok, domain_ok, typology_ok, excluded_ok]) else "N",
    }


INVENTORY_HEADERS = [
    "INV_ocs_name",
    "INV_status",
    "INV_hostname",
    "Retrived from",
    "INV_Owner_Account",
    "INV_Beneficiary_Account",
    "INV_Beneficiary_Account_ENV",
]

WORKLOAD_MATCH_HEADERS = [
    "ILU_managed",
    "ILU_IPLIST",
    "ILU_SUBNET",
    "ILU_enforcement",
    "ILU_role",
    "ILU_app",
    "ILU_env",
    "ILU_loc",
    "F_Excluded",
    "In scope",
]

WORKLOAD_RAW_ADDITIONAL_HEADERS = [
    "ILU_hostname",
    "ILU_short_hostname",
    "ILU_ip_with_default_gw",
    "ILU_OS",
    "ILU_ocs_name_from_IP",
]

RAW_SCOPE_TRACE_HEADERS = [
    "INV_ocs_name",
    "INV_status",
    "INV_hostname",
    "INV_Owner_Account",
    "INV_Beneficiary_Account",
    "INV_Beneficiary_Account_ENV",
    "ILU_managed",
    "ILU_IPLIST",
    "ILU_SUBNET",
    "ILU_enforcement",
    "ILU_role",
    "ILU_app",
    "ILU_env",
    "ILU_loc",
    "F_Excluded",
    "ILU_hostname",
    "ILU_short_hostname",
    "ILU_ip_with_default_gw",
    "ILU_OS",
    "ILU_ocs_name_from_IP",
]

RAW_SCOPE_PROGRAM_HEADERS = [
    "In Scope(s)",
    "Program(s)",
]

SCOPE_WORKSHEET_PREFERRED_COLUMNS = [
    "uid",
    "program",
    "network",
    "taken",
    "Server UID",
    "UID REL",
    "NAME REL",
    "SHORT LABEL REL",
    "IRT CODE REL",
    "IAPPLI CODE REL",
    "TRIGRAM REL",
    "ENVIRONMENT",
    "HOSTNAME",
    "STATUS",
    "Server Status",
    "USUAL NAME",
    "FRIENDLY NAME",
    "TYPOLOGY",
    "CLOUD TYPE",
    "INV_ocs_name",
    "INV_status",
    "INV_hostname",
    "INV_Owner_Account",
    "INV_Beneficiary_Account",
    "INV_Beneficiary_Account_ENV",
    "ILU_managed",
    "ILU_IPLIST",
    "ILU_SUBNET",
    "ILU_enforcement",
    "ILU_role",
    "ILU_app",
    "ILU_env",
    "ILU_loc",
]

RAW_FILTER_TAIL_HEADERS = [
    "FILTER_VALUE_server.status",
    "F_FILTER_SERVER_STATUS",
    "FILTER_VALUE_os_name",
    "F_FILTER_OS_NAME",
    "FILTER_VALUE_main_application",
    "F_FILTER_MAIN_APP_NOT_TAKEN",
    "FILTER_VALUE_environment",
    "F_FILTER_PRD_ENV",
    "FILTER_VALUE_cloud_type",
    "F_FILTER_CLOUD_TYPE_NOT_TAKEN",
    "FILTER_VALUE_domain",
    "F_FILTER_DOMAIN",
    "FILTER_VALUE_typology",
    "F_FILTER_TYPOLOGY_NOT_TAKEN",
    "F_Excluded",
    "F_FILTER_ALL",
    "In Scope(s)",
    "Program(s)",
]

FILTERED_FIXED_EXTRA_HEADERS = [
    "INV_ocs_name",
    "INV_status",
    "INV_hostname",
    "INV_Owner_Account",
    "INV_Beneficiary_Account",
    "INV_Beneficiary_Account_ENV",
    "ILU_managed",
    "ILU_IPLIST",
    "ILU_SUBNET",
    "ILU_enforcement",
    "ILU_role",
    "ILU_app",
    "ILU_env",
    "ILU_loc",
]

EXCLUDED_SHEET_HEADERS = [
    "Server to exclude",
    "Retrived by",
    "uid",
    "short_label",
    "DSI REL",
    "DALI STATUS",
    "USUAL NAME",
    "cloud_type",
    "OS NAME",
    "INV_Ocs_Name",
    "managed",
    "IPLIST",
    "enforcement",
    "role",
    "app",
    "env",
    "loc",
]

DEFAULT_PROD_BENEFICIARY_TOKENS = ["PRD", "DRP", "BCK"]

HIDDEN_WORKSHEET_NAMES = {"get_inv_by_account", "get_marley_gen2_by_uuid", "DictKearAccount", "GLOBAL", "OUT_OF_SCOPE"}

CANDIDATE_SHEET_HEADERS = [
    "short_hostname",
    "status",
    "app_info.kear_uuid",
    "uuid",
    "IPLIST",
    "scope_alread_monitored?",
    "Program",
    "owner_app_name",
    "beneficiary",
    "beneficiary_account_env",
    "hostname",
    "ocs_name_from_IP",
    "SUBNET",
    "managed",
    "enforcement",
    "created_at",
    "ip_with_default_gw",
    "app",
    "env",
    "loc",
    "role",
    "app_info.factor",
    "app_info.app_id",
    "app_info.app_name",
    "app_info.env",
    "typologie",
    "os_name",
    "dns",
    "lookup_status",
    "FILTER_VALUE_environment",
    "F_FILTER_PRD_ENV",
    "FILTER_VALUE_os_name",
    "F_FILTER_OS_NAME",
    "FILTER_VALUE_server.status",
    "F_FILTER_SERVER_STATUS",
    "FILTER_VALUE_domain",
    "F_FILTER_DOMAIN",
    "FILTER_VALUE_typology",
    "F_FILTER_TYPOLOGY_NOT_TAKEN",
    "F_FILTER_ALL",
]

PPTX_STATS_COLUMNS = [
    "Index",
    "Program",
    "Entity",
    "Sub-Entity",
    "Application Short Label",
    "Total Assets in Dali (Enriched)",
    "Assets in Dali (Enriched) not in illumio",
    "% servers with illumio installed (Enriched)",
    "% servers with illumio installed (Enriched) Indicator Icon",
    "% servers with illumio installed (Enriched) Trend Icon",
    "% servers with illumio agent in blocking mode (Enriched)",
    "% servers with illumio agent in blocking mode (Enriched) Indicator Icon",
    "% servers with illumio agent in blocking mode (Enriched) Trend Icon",
]


def _normalize_lookup_value(value: Any) -> str:
    return str(value or "").strip().upper()


def _short_hostname(value: Any) -> str:
    raw = str(value or "").strip()
    if not raw:
        return ""
    return raw.split(".", 1)[0].strip()


def _normalize_hostname_for_compare(value: Any) -> str:
    return _normalize_lookup_value(_short_hostname(value))


def _parse_ipv4_strings(value: Any) -> List[str]:
    out: List[str] = []
    seen: set[str] = set()
    for token in re.findall(r"(?:\d{1,3}\.){3}\d{1,3}(?:/\d{1,2})?", str(value or "")):
        try:
            ip_obj = ipaddress.ip_interface(token).ip if "/" in token else ipaddress.ip_address(token)
        except ValueError:
            continue
        if isinstance(ip_obj, ipaddress.IPv4Address):
            rendered = str(ip_obj)
            if rendered not in seen:
                seen.add(rendered)
                out.append(rendered)
    return out


def _pick_main_ips_for_subnet(ipv4_list: List[str], subnet_value: Any) -> List[str]:
    subnet_raw = str(subnet_value or "").strip()
    if not subnet_raw:
        return ipv4_list
    try:
        subnet = ipaddress.ip_network(subnet_raw, strict=False)
    except ValueError:
        return ipv4_list
    if not isinstance(subnet, ipaddress.IPv4Network):
        return ipv4_list
    matched = [ip for ip in ipv4_list if ipaddress.ip_address(ip) in subnet]
    return matched or ipv4_list


def _normalize_status(value: Any) -> str:
    return str(value or "").strip().upper()


def _inventory_hostid_from_server_uid(server_uid: Any) -> str:
    """Build inventory hostid key as VM_<SERVER_UID>."""
    normalized_uid = _normalize_lookup_value(server_uid)
    if not normalized_uid:
        return ""
    return f"VM_{normalized_uid}"


def _inventory_srn_from_server_uid(server_uid: Any) -> str:
    """Build canonical SRN prefix pattern from Server UID."""
    normalized_uid = _normalize_lookup_value(server_uid)
    if not normalized_uid:
        return ""
    return f"SRN:SGCP:VCS.EU-FR-PARIS:SERVER:{normalized_uid}"


def _normalize_uuid_from_hostid(hostid: Any) -> str:
    """Normalize inventory hostid by removing VM_ prefix and lowercasing."""
    raw = str(hostid or "").strip()
    if not raw:
        return ""
    if raw.upper().startswith("VM_"):
        raw = raw[3:]
    return raw.lower()


def _normalize_uuid_from_srn(srn: Any) -> str:
    """Normalize inventory SRN by taking suffix after server: and lowercasing."""
    raw = str(srn or "").strip()
    if not raw:
        return ""
    lower_raw = raw.lower()
    marker = "server:"
    idx = lower_raw.find(marker)
    if idx >= 0:
        raw = raw[idx + len(marker):]
    return raw.lower()


def _normalize_column_key(value: str) -> str:
    return "".join(ch for ch in str(value or "").lower() if ch.isalnum())


def _nested_get(data: Dict[str, Any], dotted_key: str, default: Any = "") -> Any:
    if dotted_key in data:
        return data.get(dotted_key, default)
    current: Any = data
    for part in str(dotted_key or "").split("."):
        if isinstance(current, list):
            current = current[0] if current else default
        if not isinstance(current, dict):
            return default
        if part not in current:
            return default
        current = current.get(part)
    return current


def _get_row_value_by_candidates(row: Dict[str, Any], candidates: List[str]) -> str:
    normalized_candidates = {_normalize_column_key(name) for name in candidates}
    for key, value in row.items():
        if _normalize_column_key(key) in normalized_candidates:
            return str(value or "")
    return ""


def _get_first_non_empty_by_candidates(row: Dict[str, Any], candidates: List[str]) -> str:
    for candidate in candidates:
        value = str(_get_row_value_by_candidates(row, [candidate]) or "").strip()
        if value:
            return value
    return ""


def _get_prod_beneficiary_tokens(filters: Optional[Dict[str, str]]) -> List[str]:
    tokens = _parse_filter_tokens(filters, "FILTER_PRD_ENV")
    return tokens or DEFAULT_PROD_BENEFICIARY_TOKENS


def _get_beneficiary_not_taken_tokens(filters: Optional[Dict[str, str]]) -> List[str]:
    return _parse_filter_tokens(filters, "FILTER_BENEFICIARY_ACCOUNT_NOT_TAKEN")


def _is_prod_beneficiary(value: Any, prod_tokens: List[str]) -> bool:
    normalized = _normalize_lookup_value(value)
    if not normalized:
        return False
    return any(token in normalized for token in prod_tokens)


def _effective_gen2_prd_env_value(row: Dict[str, Any]) -> str:
    cloud_type = _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type", "CLOUD TYPE"]))
    if cloud_type == "GEN 2":
        candidate = _get_row_value_by_candidates(row, ["INV_Beneficiary_Account_ENV"])
        if str(candidate or "").strip() and _normalize_lookup_value(candidate) not in {"NOT_FOUND", "NOT_GEN2"}:
            return str(candidate or "")
        return _get_row_value_by_candidates(row, ["INV_Beneficiary_Account"])
    return _get_row_value_by_candidates(row, ["FILTER_VALUE_environment", "ENVIRONMENT", "environment"])


def _recompute_prd_env_flags(rows: List[Dict[str, Any]], filters: Optional[Dict[str, str]]) -> None:
    env_tokens = _parse_filter_tokens(filters, "FILTER_PRD_ENV")
    for row in rows:
        effective_env_value = _effective_gen2_prd_env_value(row)
        env_ok = True if not env_tokens else _contains_any_token(effective_env_value, env_tokens)
        row["FILTER_VALUE_environment"] = effective_env_value
        row["F_FILTER_PRD_ENV"] = "Y" if env_ok else "N"
        flags = [
            _normalize_lookup_value(row.get("F_FILTER_PRD_ENV", "Y")),
            _normalize_lookup_value(row.get("F_FILTER_OS_NAME", "Y")),
            _normalize_lookup_value(row.get("F_FILTER_SERVER_STATUS", "Y")),
            _normalize_lookup_value(row.get("F_FILTER_CLOUD_TYPE_NOT_TAKEN", "Y")),
            _normalize_lookup_value(row.get("F_FILTER_MAIN_APP_NOT_TAKEN", "Y")),
            _normalize_lookup_value(row.get("F_FILTER_DOMAIN", "Y")),
            _normalize_lookup_value(row.get("F_FILTER_TYPOLOGY_NOT_TAKEN", "Y")),
            "N" if _normalize_lookup_value(row.get("F_Excluded", "N")) == "Y" else "Y",
        ]
        row["F_FILTER_ALL"] = "Y" if all(flag == "Y" for flag in flags) else "N"


def _parse_managed_flag(value: Any) -> bool:
    return _normalize_lookup_value(value) in {"TRUE", "1", "YES", "Y"}


def _read_workload_derived_rows(workload_csv: Path) -> List[Dict[str, str]]:
    if not workload_csv.is_file():
        log.warning("Workload derived CSV not found: %s", workload_csv)
        return []

    with workload_csv.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle)
        fieldnames = reader.fieldnames or []
        required = {"short_hostname", "managed", "IPLIST", "SUBNET", "enforcement", "role", "app", "env", "loc"}
        if not required.issubset(set(fieldnames)):
            missing = sorted(required.difference(set(fieldnames)))
            log.warning("Workload derived CSV missing required columns %s in %s", ",".join(missing), workload_csv)
            return []

        rows: List[Dict[str, str]] = []
        for row in reader:
            rows.append({key: str(value or "").strip() for key, value in row.items()})
        return rows


def _find_workload_match(workload_rows: List[Dict[str, str]], lookup_value: str) -> Optional[Dict[str, str]]:
    lookup_values = [str(lookup_value or "")]
    managed_true_short_idx, managed_true_ocs_idx, managed_false_short_idx, managed_false_ocs_idx = _build_workload_lookup_indexes(workload_rows)
    return _find_workload_match_from_candidates(
        lookup_values,
        managed_true_short_idx,
        managed_true_ocs_idx,
        managed_false_short_idx,
        managed_false_ocs_idx,
    )


def _normalize_lookup_candidates(values: List[str]) -> List[str]:
    out: List[str] = []
    seen: set[str] = set()
    for value in values:
        raw = str(value or "").strip()
        if not raw:
            continue

        normalized_full = _normalize_lookup_value(raw)
        if normalized_full and normalized_full not in seen:
            seen.add(normalized_full)
            out.append(normalized_full)

        normalized_short = _normalize_lookup_value(_short_hostname(raw))
        if normalized_short and normalized_short not in seen:
            seen.add(normalized_short)
            out.append(normalized_short)
    return out


def _build_workload_lookup_indexes(workload_rows: List[Dict[str, str]]) -> Tuple[Dict[str, Dict[str, str]], Dict[str, Dict[str, str]], Dict[str, Dict[str, str]], Dict[str, Dict[str, str]]]:
    managed_true = [row for row in workload_rows if _parse_managed_flag(row.get("managed", ""))]
    managed_false = [row for row in workload_rows if not _parse_managed_flag(row.get("managed", ""))]

    def _index(rows: List[Dict[str, str]], field_name: str, short_hostname_field: bool = False) -> Dict[str, Dict[str, str]]:
        indexed: Dict[str, Dict[str, str]] = {}
        for row in rows:
            field_value = row.get(field_name, "")
            normalized = _normalize_lookup_value(_short_hostname(field_value) if short_hostname_field else field_value)
            if normalized and normalized not in indexed:
                indexed[normalized] = row
        return indexed

    return (
        _index(managed_true, "short_hostname", short_hostname_field=True),
        _index(managed_true, "ocs_name_from_IP"),
        _index(managed_false, "short_hostname", short_hostname_field=True),
        _index(managed_false, "ocs_name_from_IP"),
    )


def _find_workload_match_from_candidates(
    lookup_candidates: List[str],
    managed_true_short_idx: Dict[str, Dict[str, str]],
    managed_true_ocs_idx: Dict[str, Dict[str, str]],
    managed_false_short_idx: Dict[str, Dict[str, str]],
    managed_false_ocs_idx: Dict[str, Dict[str, str]],
) -> Optional[Dict[str, str]]:
    normalized_candidates = _normalize_lookup_candidates(lookup_candidates)
    if not normalized_candidates:
        return None

    for candidate in normalized_candidates:
        match = managed_true_short_idx.get(candidate)
        if match:
            return match
    for candidate in normalized_candidates:
        match = managed_true_ocs_idx.get(candidate)
        if match:
            return match
    for candidate in normalized_candidates:
        match = managed_false_short_idx.get(candidate)
        if match:
            return match
    for candidate in normalized_candidates:
        match = managed_false_ocs_idx.get(candidate)
        if match:
            return match
    return None


def _build_filtered_workload_lookup_candidates(row: Dict[str, Any]) -> List[str]:
    candidates: List[str] = []

    hostname_value = _get_row_value_by_candidates(row, ["HOSTNAME", "hostname", "server_hostname", "host_name"])
    if hostname_value:
        candidates.append(hostname_value)

    usual_name_value = _get_row_value_by_candidates(row, ["USUAL NAME", "usual_name", "server_usual_name"])
    if usual_name_value:
        candidates.append(usual_name_value)

    friendly_name_value = _get_row_value_by_candidates(row, ["FRIENDLY NAME", "friendly_name", "server_friendly_name"])
    if friendly_name_value and not any(ch.isspace() for ch in friendly_name_value):
        candidates.append(friendly_name_value)

    inv_ocs_name = _get_row_value_by_candidates(row, ["INV_ocs_name"])
    if inv_ocs_name:
        candidates.append(inv_ocs_name)

    inv_hostname = _get_row_value_by_candidates(row, ["INV_hostname"])
    if inv_hostname:
        candidates.append(inv_hostname)

    return candidates


def _workload_value(match: Dict[str, str], field_name: str) -> str:
    if field_name == "OS":
        return str(match.get("OS") or match.get("os") or match.get("os_name") or "")
    if field_name == "ocs_name_from_IP":
        return str(match.get("ocs_name_from_IP") or match.get("ocs_name_from_ip") or "")
    if field_name == "ip_with_default_gw":
        return str(match.get("ip_with_default_gw") or match.get("ip_with_default_gateway") or "")
    return str(match.get(field_name, ""))


def _workload_source_field(field_name: str) -> str:
    return str(field_name or "").replace("ILU_", "", 1) if str(field_name or "").startswith("ILU_") else str(field_name or "")


def enrich_filtered_rows_with_workload_matches(filtered_rows: List[Dict[str, Any]], workload_csv: Path) -> None:
    workload_rows = _read_workload_derived_rows(workload_csv)
    workload_headers_to_copy = WORKLOAD_MATCH_HEADERS + [name for name in WORKLOAD_RAW_ADDITIONAL_HEADERS if name not in WORKLOAD_MATCH_HEADERS]
    if not workload_rows:
        log.info("Workload match enrichment skipped: no workload rows loaded")
        for row in filtered_rows:
            for header in workload_headers_to_copy:
                row[header] = row.get(header, "")
        return

    managed_true_short_idx, managed_true_ocs_idx, managed_false_short_idx, managed_false_ocs_idx = _build_workload_lookup_indexes(workload_rows)

    matched_rows = 0
    for row in filtered_rows:
        lookup_candidates = _build_filtered_workload_lookup_candidates(row)
        match = _find_workload_match_from_candidates(
            lookup_candidates,
            managed_true_short_idx,
            managed_true_ocs_idx,
            managed_false_short_idx,
            managed_false_ocs_idx,
        )

        if match:
            matched_rows += 1
            for header in workload_headers_to_copy:
                source_field = _workload_source_field(header)
                value = _workload_value(match, source_field)
                row[header] = value
                row[source_field] = value
        else:
            for header in workload_headers_to_copy:
                row[header] = ""
                row[_workload_source_field(header)] = ""

    log.info("Workload match enrichment done matched_rows=%s total_rows=%s source=%s", matched_rows, len(filtered_rows), workload_csv)


def enrich_filtered_rows_with_scope(filtered_rows: List[Dict[str, Any]]) -> None:
    for row in filtered_rows:
        row["F_Excluded"] = "N"
        network = _get_row_value_by_candidates(row, ["network"])
        iplist_name = _get_row_value_by_candidates(row, ["ILU_IPLIST", "IPLIST"])

        normalized_network = _normalize_lookup_value(network)
        normalized_iplist = _normalize_lookup_value(iplist_name)

        if (not normalized_network) or ("L1" in normalized_network):
            row["In scope"] = "TRUE"
            continue

        row["In scope"] = "TRUE" if normalized_network in normalized_iplist else "FALSE"


def read_servers_to_exclude(exclusions_file: str) -> List[str]:
    path = Path(exclusions_file)
    if not path.is_file():
        log.info("Servers-to-exclude file not found: %s", exclusions_file)
        return []

    delimiter = detect_csv_delimiter(str(path), default=",")
    values: List[str] = []
    seen: set[str] = set()
    with path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for row in reader:
            if not row:
                continue
            raw_value = str(row[0] or "").strip()
            if not raw_value:
                continue
            if normalize_header_name(raw_value) in {"hostname", "host", "server", "servers_to_exclude"}:
                continue
            normalized = _normalize_hostname_for_compare(raw_value)
            if not normalized or normalized in seen:
                continue
            seen.add(normalized)
            values.append(raw_value)

    log.info("Servers-to-exclude loaded entries=%s source=%s", len(values), exclusions_file)
    return values


def apply_manual_exclusions(filtered_rows: List[Dict[str, Any]], servers_to_exclude: List[str]) -> List[Dict[str, str]]:
    lookup_columns: List[Tuple[str, List[str], bool]] = [
        ("HOSTNAME", ["HOSTNAME", "hostname", "server_hostname", "host_name"], False),
        ("USUAL NAME", ["USUAL NAME", "usual_name", "server_usual_name"], False),
        ("FRIENDLY NAME", ["FRIENDLY NAME", "friendly_name", "server_friendly_name"], True),
        ("INV_ocs_name", ["INV_ocs_name"], False),
        ("INV_hostname", ["INV_hostname"], False),
    ]

    index_by_hostname: Dict[str, List[Tuple[str, Dict[str, Any]]]] = {}
    for row in filtered_rows:
        for label, candidates, skip_if_has_spaces in lookup_columns:
            value = _get_row_value_by_candidates(row, candidates)
            if not value:
                continue
            if skip_if_has_spaces and any(ch.isspace() for ch in value):
                continue
            normalized = _normalize_hostname_for_compare(value)
            if not normalized:
                continue
            index_by_hostname.setdefault(normalized, []).append((label, row))

    excluded_rows: List[Dict[str, str]] = []
    for input_value in servers_to_exclude:
        normalized_input = _normalize_hostname_for_compare(input_value)
        matches = index_by_hostname.get(normalized_input, [])
        selected_match = matches[0] if matches else ("", {})
        matched_by, row = selected_match

        if row:
            row["F_Excluded"] = "Y"
            row["F_FILTER_ALL"] = "N"
            row["In scope"] = "FALSE"

        excluded_rows.append(
            {
                "Server to exclude": input_value,
                "Retrived by": matched_by,
                "uid": _get_row_value_by_candidates(row, ["uid"]),
                "short_label": _get_row_value_by_candidates(row, ["short_label", "SHORT LABEL REL"]),
                "DSI REL": _get_row_value_by_candidates(row, ["DSI REL", "dsi"]),
                "DALI STATUS": _get_row_value_by_candidates(row, ["DALI STATUS", "usage"]),
                "USUAL NAME": _get_row_value_by_candidates(row, ["USUAL NAME", "usual_name"]),
                "cloud_type": _get_row_value_by_candidates(row, ["cloud_type", "CLOUD TYPE", "server_cloud_type"]),
                "OS NAME": _get_row_value_by_candidates(row, ["OS NAME", "os_name"]),
                "INV_Ocs_Name": _get_row_value_by_candidates(row, ["INV_Ocs_Name", "INV_ocs_name"]),
                "managed": _get_row_value_by_candidates(row, ["ILU_managed", "managed"]),
                "IPLIST": _get_row_value_by_candidates(row, ["ILU_IPLIST", "IPLIST"]),
                "enforcement": _get_row_value_by_candidates(row, ["ILU_enforcement", "enforcement"]),
                "role": _get_row_value_by_candidates(row, ["ILU_role", "role"]),
                "app": _get_row_value_by_candidates(row, ["ILU_app", "app"]),
                "env": _get_row_value_by_candidates(row, ["ILU_env", "env"]),
                "loc": _get_row_value_by_candidates(row, ["ILU_loc", "loc"]),
            }
        )

    return excluded_rows


def _filtered_row_dedup_key(row: Dict[str, Any], row_index: int) -> Tuple[str, str, str, str]:
    """Build a stable deduplication key for FILTRED rows.

    Rows are deduplicated at server granularity inside one monitored application/program scope.
    If no server identifier is available, keep row uniqueness by index.
    """
    app_uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"]))
    program = _normalize_lookup_value(_get_row_value_by_candidates(row, ["program"]))
    server_uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
    hostname = _normalize_lookup_value(
        _short_hostname(
            _get_row_value_by_candidates(
                row,
                ["HOSTNAME", "hostname", "INV_hostname", "INV_ocs_name", "server_hostname", "host_name"],
            )
        )
    )
    server_identity = server_uid or hostname
    if not server_identity:
        server_identity = f"ROW_{row_index}"
    return app_uid, program, server_identity, _normalize_lookup_value(_get_row_value_by_candidates(row, ["taken"]))


def deduplicate_filtered_rows_by_network_iplist(filtered_rows: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Keep one row per server using network/IPLIST consistency rules.

    Priority inside duplicates:
    1) keep row where network matches IPLIST
    2) else keep row where IPLIST is empty
    3) else keep first row
    """
    groups: Dict[Tuple[str, str, str, str], List[Dict[str, Any]]] = {}
    for idx, row in enumerate(filtered_rows):
        groups.setdefault(_filtered_row_dedup_key(row, idx), []).append(row)

    deduped_rows: List[Dict[str, Any]] = []
    for key, rows in groups.items():
        if len(rows) == 1:
            deduped_rows.append(rows[0])
            continue

        def _rank(candidate: Dict[str, Any]) -> int:
            network = _normalize_lookup_value(_get_row_value_by_candidates(candidate, ["network"]))
            iplist = _normalize_lookup_value(_get_row_value_by_candidates(candidate, ["ILU_IPLIST", "IPLIST"]))
            if network and iplist and network in iplist:
                return 0
            if not iplist:
                return 1
            return 2

        selected = min(rows, key=_rank)
        deduped_rows.append(selected)
        log.debug("FILTRED dedupe group=%s size=%s kept_network=%s kept_iplist=%s", key, len(rows), selected.get("network", ""), selected.get("IPLIST", ""))

    if len(deduped_rows) != len(filtered_rows):
        log.info("FILTRED dedupe applied before=%s after=%s removed=%s", len(filtered_rows), len(deduped_rows), len(filtered_rows) - len(deduped_rows))
    return deduped_rows


def _scope_trace_key(row: Dict[str, Any]) -> Tuple[str, str, str, str]:
    """Stable key for FILTRED/SCOPE traceability comparisons."""
    app_uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"]))
    program = _normalize_lookup_value(_get_row_value_by_candidates(row, ["program"]))
    server_uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
    hostname = _normalize_lookup_value(
        _short_hostname(
            _get_row_value_by_candidates(
                row,
                ["HOSTNAME", "hostname", "INV_hostname", "INV_ocs_name", "server_hostname", "host_name"],
            )
        )
    )
    server_identity = server_uid or hostname or "NO_SERVER"
    taken = _normalize_lookup_value(_get_row_value_by_candidates(row, ["taken"]))
    return app_uid, program, server_identity, taken


def build_enrich_rows_from_marley(
    marley_rows: List[Dict[str, Any]],
    inv_by_account_rows: List[Dict[str, Any]],
    dict_kear_account_rows: List[Dict[str, Any]],
    mappings: List[Tuple[str, str]],
    raw_extra_fieldnames: List[str],
) -> List[Dict[str, Any]]:
    """Build ENRICH rows from get_marley_gen2_by_uuid FOUND rows, shaped like RAW headers."""
    raw_fieldnames = ["uid", "Server UID"] + [display for display, _ in mappings] + list(raw_extra_fieldnames)
    dict_by_uid = {
        _normalize_lookup_value(row.get("uid", "")): row
        for row in dict_kear_account_rows
        if _normalize_lookup_value(row.get("uid", ""))
    }
    dict_by_beneficiary = {
        _normalize_lookup_value(row.get("INV_Beneficiary_Account", "")): row
        for row in dict_kear_account_rows
        if _normalize_lookup_value(row.get("INV_Beneficiary_Account", ""))
    }
    inv_by_ocs_name = _index_rows_by_ocs_name(inv_by_account_rows)
    inv_by_beneficiary: Dict[str, Dict[str, Any]] = {}
    for row in inv_by_account_rows:
        key = _normalize_lookup_value(row.get("beneficiary", ""))
        if key and key not in inv_by_beneficiary:
            inv_by_beneficiary[key] = row
    out: List[Dict[str, Any]] = []
    for marley in marley_rows:
        if _normalize_lookup_value(marley.get("lookup_status", "")) != "FOUND":
            continue
        enrich_row = {field: "" for field in raw_fieldnames}
        enrich_row["uid"] = _normalize_cell_value(marley.get("KEAR_OVERRIDE", ""))
        enrich_row["Server UID"] = _normalize_cell_value(_get_row_value_by_candidates(marley, ["uuid", "Server UID", "server_uid"]))
        beneficiary_key = _normalize_lookup_value(marley.get("beneficiary", ""))
        dict_row = dict_by_uid.get(_normalize_lookup_value(enrich_row.get("uid", "")), {}) or dict_by_beneficiary.get(beneficiary_key, {})
        inv_row = inv_by_ocs_name.get(_normalize_lookup_value(marley.get("ocs_name", "")), {}) or inv_by_beneficiary.get(beneficiary_key, {})
        overrides = {
            "main_application": _normalize_cell_value(marley.get("app_info.app_id", "")),
            "environment": _normalize_cell_value(dict_row.get("INV_Beneficiary_Account_ENV", "")),
            "hostname": _normalize_cell_value(inv_row.get("ocs_name", "")),
            "usage": _normalize_cell_value(marley.get("usage", "")),
            "status": "In production",
            "main_ip": _normalize_cell_value(inv_row.get("ip", "")),
            "usual_name": _normalize_cell_value(inv_row.get("ocs_name", "")),
            "friendly_name": _normalize_cell_value(inv_row.get("hostname", "")),
            "typology": _normalize_cell_value(marley.get("typologie", "")),
            "cloud_type": "Gen 2",
            "service_offer": _normalize_cell_value(inv_row.get("service_name", "")),
            "os_name": _normalize_cell_value(marley.get("os_name", "")),
            "os_release": _normalize_cell_value(marley.get("os_version", "")),
            "vrf_name": "",
            "silo": _normalize_cell_value(marley.get("silos", "")),
            "updated_by": "KEAR",
            "beneficiary_account_id": "",
            "owner_account_id": "",
            "server.status": _normalize_cell_value(marley.get("status", "")),
            "dns_name": _normalize_cell_value(marley.get("dns", "")),
        }
        for display, technical in mappings:
            marley_value = _normalize_cell_value(overrides.get(technical, marley.get(technical, "")))
            dict_value = _normalize_cell_value(dict_row.get(technical, ""))
            enrich_row[display] = marley_value or dict_value

        enrich_row["INV_ocs_name"] = _normalize_cell_value(inv_row.get("ocs_name", ""))
        enrich_row["INV_status"] = _normalize_cell_value(inv_row.get("status", ""))
        enrich_row["INV_hostname"] = _normalize_cell_value(inv_row.get("hostname", ""))
        enrich_row["INV_Owner_Account"] = _normalize_cell_value(marley.get("owner_app_name", "")) or _normalize_cell_value(enrich_row.get("INV_Owner_Account", ""))
        enrich_row["INV_Beneficiary_Account"] = _normalize_cell_value(inv_row.get("beneficiary", ""))
        enrich_row["INV_Beneficiary_Account_ENV"] = _normalize_cell_value(dict_row.get("INV_Beneficiary_Account_ENV", ""))
        enrich_row["ILU_managed"] = _normalize_cell_value(marley.get("ILU_managed", ""))
        enrich_row["ILU_IPLIST"] = _normalize_cell_value(marley.get("ILU_IPLIST", ""))
        enrich_row["ILU_SUBNET"] = _normalize_cell_value(marley.get("ILU_SUBNET", ""))
        enrich_row["ILU_enforcement"] = _normalize_cell_value(marley.get("ILU_enforcement", ""))
        enrich_row["ILU_role"] = _normalize_cell_value(marley.get("ILU_role", ""))
        enrich_row["ILU_app"] = _normalize_cell_value(marley.get("ILU_app", ""))
        enrich_row["ILU_env"] = _normalize_cell_value(marley.get("ILU_env", ""))
        enrich_row["ILU_loc"] = _normalize_cell_value(marley.get("ILU_loc", ""))
        enrich_row["ILU_OS"] = _normalize_cell_value(marley.get("ILU_OS", ""))
        enrich_row["ILU_hostname"] = _normalize_cell_value(marley.get("ILU_hostname", ""))
        enrich_row["ILU_short_hostname"] = _normalize_cell_value(marley.get("ILU_short_hostname", ""))
        enrich_row["ILU_ip_with_default_gw"] = _normalize_cell_value(marley.get("ILU_ip_with_default_gw", ""))
        enrich_row["ILU_ocs_nam_from_IP"] = _normalize_cell_value(marley.get("ILU_ocs_nam_from_IP", ""))
        out.append(enrich_row)

    log.info("ENRICH build from Marley FOUND rows=%s", len(out))
    return out


def _lookup_variants(value: str) -> List[str]:
    raw = str(value or "").strip()
    if not raw:
        return []
    variants: List[str] = []
    short_raw = _short_hostname(raw)
    for candidate in (raw, short_raw, raw.upper(), raw.lower(), short_raw.upper(), short_raw.lower()):
        if candidate and candidate not in variants:
            variants.append(candidate)
    return variants

def _pick_inventory_row(docs: List[Dict[str, Any]], retrieved_from: str) -> Dict[str, str]:
    if not docs:
        return {}
    first = docs[0]
    return {
        "INV_ocs_name": _normalize_cell_value(first.get("ocs_name")),
        "INV_status": _normalize_status(first.get("status")),
        "INV_hostname": _short_hostname(_normalize_cell_value(first.get("hostname"))),
        "Retrived from": retrieved_from,
        "INV_Owner_Account": _normalize_cell_value(first.get("owner_app_name")),
        "INV_Beneficiary_Account": _normalize_cell_value(first.get("beneficiary")),
    }


def _deduplicate_docs(docs: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    unique_docs: List[Dict[str, Any]] = []
    fingerprints = set()
    for doc in docs:
        fingerprint = json.dumps(doc, sort_keys=True, ensure_ascii=False)
        if fingerprint in fingerprints:
            continue
        fingerprints.add(fingerprint)
        unique_docs.append(doc)
    return unique_docs


def _inventory_search_by_field(
    client: Data4secClient,
    search_field: str,
    values: List[str],
    source_fields_override: Optional[List[str]] = None,
) -> Dict[str, List[Dict[str, Any]]]:
    """Centralized inventory search wrapper with consistent logging/filters."""
    cfg = QUERY_CONFIG["inventory"]
    log.info(
        "Inventory lookup start field=%s lookup_values=%s index=%s",
        search_field,
        len(values),
        cfg["index"],
    )
    result = client.bulk_search_multi(
        index_name=cfg["index"],
        search_field=search_field,
        values=values,
        source_fields=source_fields_override or cfg["source_fields"],
        scroll_timeout=QUERY_CONFIG.get("scroll_timeout", "10m"),
        size=QUERY_CONFIG.get("batch_size", 500),
        term_filters=cfg.get("term_filters", {}),
    )
    non_empty = sum(1 for docs in result.values() if docs)
    total_docs = sum(len(docs) for docs in result.values())
    log.info(
        "Inventory lookup done field=%s matched_values=%s total_docs=%s",
        search_field,
        non_empty,
        total_docs,
    )
    return result


def query_inventory_for_server_uids(client: Data4secClient, server_uids: List[str]) -> Dict[str, Dict[str, str]]:
    """Resolve inventory rows from Server UID with hostid->srn fallback strategy."""
    uid_to_hostid: Dict[str, str] = {}
    uid_to_srn: Dict[str, str] = {}
    for server_uid in server_uids:
        normalized_uid = _normalize_lookup_value(server_uid)
        if not normalized_uid:
            continue
        uid_to_hostid[normalized_uid] = _inventory_hostid_from_server_uid(normalized_uid)
        uid_to_srn[normalized_uid] = _inventory_srn_from_server_uid(normalized_uid)

    if not uid_to_hostid:
        log.info("Inventory enrichment skipped: no Server UID values to query")
        return {}

    hostid_lookup_values = sorted(set(uid_to_hostid.values()))
    srn_lookup_values = sorted(set(uid_to_srn.values()))
    log.info(
        "Inventory Server UID enrichment prepared server_uids=%s hostid_lookup_values=%s srn_lookup_values=%s",
        len(uid_to_hostid),
        len(hostid_lookup_values),
        len(srn_lookup_values),
    )

    aggregated: Dict[str, List[Dict[str, Any]]] = {uid: [] for uid in uid_to_hostid.keys()}

    hostid_results = _inventory_search_by_field(client=client, search_field="hostid", values=hostid_lookup_values)
    for input_value, docs in hostid_results.items():
        normalized_value = _normalize_lookup_value(input_value)
        if not normalized_value or not docs:
            continue
        for uid, expected_hostid in uid_to_hostid.items():
            if normalized_value == _normalize_lookup_value(expected_hostid):
                aggregated[uid].extend(docs)

    missing_uids = [uid for uid, docs in aggregated.items() if not docs]
    if missing_uids:
        log.info("Inventory Server UID hostid retry without status filter missing_uids=%s", len(missing_uids))
        cfg = QUERY_CONFIG["inventory"]
        hostid_values_for_missing = [uid_to_hostid[uid] for uid in missing_uids if uid in uid_to_hostid]
        hostid_no_status_results = client.bulk_search_multi(
            index_name=cfg["index"],
            search_field="hostid",
            values=hostid_values_for_missing,
            source_fields=cfg["source_fields"],
            scroll_timeout=QUERY_CONFIG.get("scroll_timeout", "10m"),
            size=QUERY_CONFIG.get("batch_size", 500),
            term_filters={},
        )
        for input_value, docs in hostid_no_status_results.items():
            normalized_value = _normalize_lookup_value(input_value)
            if not normalized_value or not docs:
                continue
            for uid in missing_uids:
                expected_hostid = uid_to_hostid.get(uid, "")
                if normalized_value == _normalize_lookup_value(expected_hostid):
                    aggregated[uid].extend(docs)

    missing_uids = [uid for uid, docs in aggregated.items() if not docs]
    if missing_uids:
        log.info("Inventory Server UID enrichment fallback on srn missing_uids=%s", len(missing_uids))
        cfg = QUERY_CONFIG["inventory"]
        for uid in missing_uids:
            uid_token = _normalize_lookup_value(uid)
            if not uid_token:
                continue
            wildcard_value = f"*server:{uid_token}*"
            docs = client.search_by_wildcard(
                index_name=cfg["index"],
                search_field="srn",
                wildcard_value=wildcard_value,
                source_fields=cfg["source_fields"],
                scroll_timeout=QUERY_CONFIG.get("scroll_timeout", "10m"),
                size=QUERY_CONFIG.get("batch_size", 500),
                term_filters=cfg.get("term_filters", {}),
            )
            if not docs:
                log.info(
                    "SRN wildcard with status filter returned 0 for uid=%s, retrying without status filter",
                    uid,
                )
                docs = client.search_by_wildcard(
                    index_name=cfg["index"],
                    search_field="srn",
                    wildcard_value=wildcard_value,
                    source_fields=cfg["source_fields"],
                    scroll_timeout=QUERY_CONFIG.get("scroll_timeout", "10m"),
                    size=QUERY_CONFIG.get("batch_size", 500),
                    term_filters={},
                )
            if docs:
                aggregated[uid].extend(docs)

    output: Dict[str, Dict[str, str]] = {}
    matched = 0
    for uid, docs in aggregated.items():
        dedup_docs = _deduplicate_docs(docs)
        if dedup_docs:
            matched += 1
        output[uid] = _pick_inventory_row(dedup_docs, retrieved_from="Dali Export")
    log.info(
        "Inventory Server UID enrichment done matched_server_uids=%s total_server_uids=%s",
        matched,
        len(uid_to_hostid),
    )
    return output


def query_inventory_for_beneficiaries(client: Data4secClient, beneficiaries: List[str]) -> Dict[str, List[Dict[str, Any]]]:
    """Fetch inventory documents by beneficiary account list."""
    cfg = QUERY_CONFIG["inventory"]
    search_field = cfg.get("beneficiary_search_field", "beneficiary")
    lookup_values = [_normalize_lookup_value(value) for value in beneficiaries if _normalize_lookup_value(value)]
    if not lookup_values:
        log.info("Inventory beneficiary discovery skipped: no beneficiaries")
        return {}
    log.info("Inventory beneficiary discovery prepared beneficiaries=%s", len(lookup_values))

    result_map = _inventory_search_by_field(client=client, search_field=search_field, values=lookup_values)
    out: Dict[str, List[Dict[str, Any]]] = {}
    for key, docs in result_map.items():
        normalized_key = _normalize_lookup_value(key)
        if not normalized_key or not docs:
            continue
        out.setdefault(normalized_key, []).extend(docs)

    deduped = {key: _deduplicate_docs(value) for key, value in out.items()}
    log.info(
        "Inventory beneficiary discovery done matched_beneficiaries=%s total_docs=%s",
        len(deduped),
        sum(len(v) for v in deduped.values()),
    )
    return deduped


def _extract_env_from_platform_tags(tags_value: Any) -> str:
    raw_items: List[str] = []
    if isinstance(tags_value, list):
        raw_items = [str(item or "").strip() for item in tags_value if str(item or "").strip()]
    else:
        text = str(tags_value or "").strip()
        if text.startswith("[") and text.endswith("]"):
            text = text[1:-1]
        raw_items = [part.strip() for part in text.split(",") if part.strip()]

    for item in raw_items:
        if ":" not in item:
            continue
        key, value = item.split(":", 1)
        if _normalize_lookup_value(key) == "ENV":
            return str(value or "").strip()
    return ""


def query_platform_accounts_env_by_names(client: Data4secClient, account_names: List[str]) -> Dict[str, str]:
    cfg = QUERY_CONFIG.get("platform_accounts", {})
    index_name = str(cfg.get("index", "platform_accounts"))
    search_field = str(cfg.get("search_field", "name"))
    source_fields = list(cfg.get("source_fields", ["name", "tags"]))
    term_filters = cfg.get("term_filters", {})

    normalized_names = sorted({
        _normalize_lookup_value(name)
        for name in account_names
        if _normalize_lookup_value(name)
    })
    if not normalized_names:
        return {}

    log.info(
        "Platform accounts lookup start index=%s search_field=%s beneficiaries=%s",
        index_name,
        search_field,
        len(normalized_names),
    )
    result_map = client.bulk_search_multi(
        index_name=index_name,
        search_field=search_field,
        values=normalized_names,
        source_fields=source_fields,
        scroll_timeout=QUERY_CONFIG.get("scroll_timeout", "10m"),
        size=QUERY_CONFIG.get("batch_size", 500),
        term_filters=term_filters,
    )

    output: Dict[str, str] = {}
    for name in normalized_names:
        docs = result_map.get(name, [])
        env_value = ""
        for doc in docs:
            env_value = _extract_env_from_platform_tags(doc.get("tags", ""))
            if env_value:
                break
        if env_value:
            output[name] = env_value
    log.info(
        "Platform accounts lookup done matched_env=%s total_names=%s",
        len(output),
        len(normalized_names),
    )
    return output


def query_marley_original_by_field(
    client: Data4secClient,
    lookup_values: List[str],
    search_field: str,
    lookup_label: str,
    source_fields_override: Optional[List[str]] = None,
) -> Dict[str, List[Dict[str, Any]]]:
    cfg = QUERY_CONFIG.get("marley_original", {})
    index_name = str(cfg.get("index", "marley_original"))
    source_fields = list(cfg.get("source_fields", []))
    term_filters = cfg.get("term_filters", {})

    normalized_lookup_values = [value for value in (_normalize_lookup_value(name) for name in lookup_values) if value]
    if not normalized_lookup_values:
        log.info("Marley lookup skipped: no %s values", lookup_label)
        return {}

    log.info(
        "Marley lookup start index=%s search_field=%s lookup_values=%s lookup_label=%s",
        index_name,
        search_field,
        len(normalized_lookup_values),
        lookup_label,
    )
    result_map = client.bulk_search_multi(
        index_name=index_name,
        search_field=search_field,
        values=sorted(set(normalized_lookup_values)),
        source_fields=source_fields_override or source_fields,
        scroll_timeout=QUERY_CONFIG.get("scroll_timeout", "10m"),
        size=QUERY_CONFIG.get("batch_size", 500),
        term_filters=term_filters,
    )

    out: Dict[str, List[Dict[str, Any]]] = {}
    for key, docs in result_map.items():
        normalized_key = _normalize_lookup_value(key)
        if not normalized_key or not docs:
            continue
        out.setdefault(normalized_key, []).extend(docs)

    missing_keys = [key for key in sorted(set(normalized_lookup_values)) if not out.get(key)]
    if missing_keys:
        log.info("Marley lookup fallback start missing_keys=%s (case-insensitive wildcard)", len(missing_keys))
        for key in missing_keys:
            docs = client.search_by_wildcard(
                index_name=index_name,
                search_field=search_field,
                wildcard_value=key,
                source_fields=source_fields_override or source_fields,
                scroll_timeout=QUERY_CONFIG.get("scroll_timeout", "10m"),
                size=QUERY_CONFIG.get("batch_size", 500),
                term_filters=term_filters,
            )
            if not docs:
                docs = client.search_by_wildcard(
                    index_name=index_name,
                    search_field=search_field,
                    wildcard_value=f"*{key}*",
                    source_fields=source_fields,
                    scroll_timeout=QUERY_CONFIG.get("scroll_timeout", "10m"),
                    size=QUERY_CONFIG.get("batch_size", 500),
                    term_filters=term_filters,
                )
            if docs:
                out.setdefault(key, []).extend(docs)

    deduped = {key: _deduplicate_docs(value) for key, value in out.items()}
    log.info(
        "Marley lookup done matched_%s=%s total_docs=%s",
        lookup_label,
        len(deduped),
        sum(len(v) for v in deduped.values()),
    )
    return deduped


def query_marley_original_by_ocs_names(client: Data4secClient, ocs_names: List[str]) -> Dict[str, List[Dict[str, Any]]]:
    cfg = QUERY_CONFIG.get("marley_original", {})
    search_field = str(cfg.get("search_field", "hostname"))
    return query_marley_original_by_field(
        client=client,
        lookup_values=ocs_names,
        search_field=search_field,
        lookup_label="ocs_names",
    )


def query_marley_original_by_uuids(client: Data4secClient, lookup_uuids: List[str]) -> Dict[str, List[Dict[str, Any]]]:
    return query_marley_original_by_field(
        client=client,
        lookup_values=lookup_uuids,
        search_field="uuid",
        lookup_label="uuids",
    )


def build_marley_sheet_rows(
    inventory_by_account_rows: List[Dict[str, Any]],
    marley_docs_by_lookup: Dict[str, List[Dict[str, Any]]],
    monitored_uids: set[str],
    lookup_source_field: str = "ocs_name",
    lookup_output_field: str = "lookup_hostname",
) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    normalized_scope_uids = {_normalize_lookup_value(uid) for uid in monitored_uids if _normalize_lookup_value(uid)}
    for source_row in inventory_by_account_rows:
        ocs_name = _normalize_cell_value(source_row.get("ocs_name", ""))
        lookup_value = _normalize_cell_value(source_row.get(lookup_source_field, ""))
        normalized_lookup_value = _normalize_lookup_value(lookup_value)
        docs = marley_docs_by_lookup.get(normalized_lookup_value, [])

        if not docs:
            rows.append(
                {
                    lookup_output_field: normalized_lookup_value,
                    "ocs_name": ocs_name,
                    "beneficiary": _normalize_cell_value(source_row.get("beneficiary", "")),
                    "owner_app_name": _normalize_cell_value(source_row.get("owner_app_name", "")),
                    "app_info.account_id": "",
                    "app_info.app_id": "",
                    "app_info.app_name": "",
                    "app_info.env": "",
                    "app_info.factor": "",
                    "app_info.kear_uuid": "",
                    "app_info.kear_library": "",
                    "app_info.ref_app": "",
                    "app_info.service_line_name": "",
                    "uuid": "",
                    "net_info.net_ipadress": "",
                    "os_name": "",
                    "os_version": "",
                    "typologie": "",
                    "silos": "",
                    "dns": "",
                    "status": "",
                    "usage": "",
                    "Kear in scope": "FALSE",
                    "lookup_status": "NOT_FOUND",
                }
            )
            continue

        for doc in docs:
            marley_kear_uuid = _normalize_cell_value(_nested_get(doc, "app_info.kear_uuid", ""))
            rows.append(
                {
                    lookup_output_field: normalized_lookup_value,
                    "ocs_name": ocs_name,
                    "beneficiary": _normalize_cell_value(source_row.get("beneficiary", "")),
                    "owner_app_name": _normalize_cell_value(source_row.get("owner_app_name", "")),
                    "app_info.account_id": _normalize_cell_value(_nested_get(doc, "app_info.account_id", "")),
                    "app_info.app_id": _normalize_cell_value(_nested_get(doc, "app_info.app_id", "")),
                    "app_info.app_name": _normalize_cell_value(_nested_get(doc, "app_info.app_name", "")),
                    "app_info.env": _normalize_cell_value(_nested_get(doc, "app_info.env", "")),
                    "app_info.factor": _normalize_cell_value(_nested_get(doc, "app_info.factor", "")),
                    "app_info.kear_uuid": marley_kear_uuid,
                    "app_info.kear_library": _normalize_cell_value(_nested_get(doc, "app_info.kear_library", "")),
                    "app_info.ref_app": _normalize_cell_value(_nested_get(doc, "app_info.ref_app", "")),
                    "app_info.service_line_name": _normalize_cell_value(_nested_get(doc, "app_info.service_line_name", "")),
                    "uuid": _normalize_cell_value(doc.get("uuid", "")),
                    "net_info.net_ipadress": _normalize_cell_value(_nested_get(doc, "net_info.net_ipadress", "")),
                    "os_name": _normalize_cell_value(doc.get("os_name", "")),
                    "os_version": _normalize_cell_value(doc.get("os_version", "")),
                    "typologie": _normalize_cell_value(doc.get("typologie", "")),
                    "silos": _normalize_cell_value(doc.get("silos", "")),
                    "dns": _normalize_cell_value(doc.get("dns", "")),
                    "status": _normalize_cell_value(doc.get("status", "")),
                    "usage": _normalize_cell_value(doc.get("usage", "")),
                    "Kear in scope": "TRUE" if _normalize_lookup_value(marley_kear_uuid) in normalized_scope_uids else "FALSE",
                    "lookup_status": "FOUND",
                }
            )
    return rows


def filter_marley_sheet_rows(
    marley_rows: List[Dict[str, Any]],
    filtered_rows: List[Dict[str, Any]],
    filters: Optional[Dict[str, str]] = None,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    take_only_scope_kears = _parse_filter_bool(
        filters,
        "FILTER_TAKE_ONLY_KEARS_IN_SCOPE_FROM_MARLEY",
        default=True,
    )
    owner_not_taken_tokens = _parse_filter_tokens(filters, "FILTER_OWNER_ACCOUNT_NOT_TAKEN")
    main_app_not_taken_tokens = _parse_filter_tokens(filters, "FILTER_MAIN_APP_NOT_TAKEN")
    env_tokens = _parse_filter_tokens(filters, "FILTER_PRD_ENV")
    os_tokens = _parse_filter_tokens(filters, "FILTER_OS_NAME")
    accounts_not_to_enrich_tokens = _parse_filter_tokens(filters, "FILTER_OWNER_ACCOUNTS_NOT_TO_ENRICH")

    existing_server_uids = {
        _normalize_lookup_value(row.get("Server UID", ""))
        for row in filtered_rows
        if _normalize_lookup_value(row.get("Server UID", ""))
    }

    annotated_rows: List[Dict[str, Any]] = []
    kept_rows: List[Dict[str, Any]] = []
    for row in marley_rows:
        lookup_found_ok = _normalize_lookup_value(row.get("lookup_status", "")) == "FOUND"
        kear_scope_ok = _normalize_lookup_value(row.get("Kear in scope", "")) == "TRUE"
        status_active_ok = _normalize_lookup_value(row.get("status", "")) == "ACTIVE"
        usage_in_use_ok = _normalize_lookup_value(row.get("usage", "")) == "IN USE"

        uuid_value = _normalize_lookup_value(row.get("uuid", ""))
        uuid_in_filtered = bool(uuid_value) and uuid_value in existing_server_uids
        uuid_not_in_filtered_ok = bool(uuid_value) and (not uuid_in_filtered)

        owner_value = _normalize_cell_value(row.get("owner_app_name", ""))
        owner_not_taken_ok = not (owner_not_taken_tokens and _contains_any_token(owner_value, owner_not_taken_tokens))

        app_id_value = _normalize_cell_value(row.get("app_info.app_id", ""))
        main_app_not_taken_ok = not (main_app_not_taken_tokens and _contains_any_token(app_id_value, main_app_not_taken_tokens))

        env_value = _normalize_cell_value(row.get("app_info.env", ""))
        env_filter_ok = True if not env_tokens else _contains_any_token(env_value, env_tokens)

        os_name_value = _normalize_cell_value(row.get("os_name", ""))
        os_filter_ok = True if not os_tokens else _matches_exact_token(os_name_value, os_tokens)

        account_enrich_allowed = not (accounts_not_to_enrich_tokens and _matches_exact_token(owner_value, accounts_not_to_enrich_tokens))

        kear_scope_filter_ok = kear_scope_ok if take_only_scope_kears else True

        final_keep = all(
            [
                lookup_found_ok,
                kear_scope_filter_ok,
                status_active_ok,
                usage_in_use_ok,
                uuid_not_in_filtered_ok,
                owner_not_taken_ok,
                main_app_not_taken_ok,
                env_filter_ok,
                os_filter_ok,
                account_enrich_allowed,
            ]
        )

        annotated_row = dict(row)
        annotated_row.update(
            {
                "F_lookup_status_FOUND": "Y" if lookup_found_ok else "N",
                "F_kear_in_scope_TRUE": "Y" if kear_scope_ok else "N",
                "F_take_only_kears_in_scope": "Y" if take_only_scope_kears else "N",
                "F_kear_scope_filter": "Y" if kear_scope_filter_ok else "N",
                "F_status_Active": "Y" if status_active_ok else "N",
                "F_usage_In_use": "Y" if usage_in_use_ok else "N",
                "F_uuid_in_filtered": "Y" if uuid_in_filtered else "N",
                "F_uuid_not_in_filtered": "Y" if uuid_not_in_filtered_ok else "N",
                "F_owner_account_allowed": "Y" if owner_not_taken_ok else "N",
                "F_main_app_allowed": "Y" if main_app_not_taken_ok else "N",
                "F_env_match": "Y" if env_filter_ok else "N",
                "F_os_match": "Y" if os_filter_ok else "N",
                "F_account_not_to_enrich": "N" if account_enrich_allowed else "Y",
                "F_account_enrich_allowed": "Y" if account_enrich_allowed else "N",
                "F_final_keep": "Y" if final_keep else "N",
            }
        )
        annotated_rows.append(annotated_row)

        if final_keep:
            kept_rows.append(dict(annotated_row))

    log.info(
        "Marley sheet filtering done input_rows=%s output_rows=%s take_only_scope_kears=%s owner_excluded=%s main_app_excluded=%s env_tokens=%s os_tokens=%s accounts_not_to_enrich=%s",
        len(annotated_rows),
        len(kept_rows),
        take_only_scope_kears,
        owner_not_taken_tokens,
        main_app_not_taken_tokens,
        env_tokens,
        os_tokens,
        accounts_not_to_enrich_tokens,
    )
    return annotated_rows, kept_rows


def enrich_marley_rows_with_workload(marley_rows: List[Dict[str, Any]], workload_csv: Path) -> None:
    workload_rows = _read_workload_derived_rows(workload_csv)
    marley_ilu_headers = [
        "ILU_managed",
        "ILU_IPLIST",
        "ILU_SUBNET",
        "ILU_enforcement",
        "ILU_role",
        "ILU_app",
        "ILU_env",
        "ILU_loc",
        "ILU_OS",
        "ILU_hostname",
        "ILU_short_hostname",
        "ILU_interfaces",
        "ILU_ip_with_default_gw",
        "ILU_ocs_name_from_IP",
        "ILU_ocs_nam_from_IP",
    ]
    if not workload_rows:
        for row in marley_rows:
            row["MAIN IP"] = ""
            for header in marley_ilu_headers:
                row[header] = row.get(header, "")
            row.pop("interfaces", None)
        return

    managed_true_short_idx, managed_true_ocs_idx, managed_false_short_idx, managed_false_ocs_idx = _build_workload_lookup_indexes(workload_rows)
    for row in marley_rows:
        lookup_candidates = [
            _normalize_cell_value(row.get("ocs_name", "")),
            _normalize_cell_value(row.get("hostname", "")),
            _normalize_cell_value(row.get("INV_hostname", "")),
        ]
        match = _find_workload_match_from_candidates(
            lookup_candidates=lookup_candidates,
            managed_true_short_idx=managed_true_short_idx,
            managed_true_ocs_idx=managed_true_ocs_idx,
            managed_false_short_idx=managed_false_short_idx,
            managed_false_ocs_idx=managed_false_ocs_idx,
        )
        if not match:
            row["MAIN IP"] = ""
            for header in marley_ilu_headers:
                row[header] = ""
            row.pop("interfaces", None)
            continue

        interfaces_raw = _normalize_cell_value(match.get("interfaces", ""))
        ipv4_list = _parse_ipv4_strings(interfaces_raw)
        main_ips = _pick_main_ips_for_subnet(ipv4_list, match.get("SUBNET", ""))
        row["ILU_interfaces"] = interfaces_raw
        row["MAIN IP"] = ", ".join(main_ips)
        for header in marley_ilu_headers:
            source_field = _workload_source_field(header)
            row[header] = _normalize_cell_value(_workload_value(match, source_field))
        row.pop("interfaces", None)


def _index_rows_by_ocs_name(rows: List[Dict[str, Any]]) -> Dict[str, Dict[str, Any]]:
    out: Dict[str, Dict[str, Any]] = {}
    for row in rows:
        key = _normalize_lookup_value(row.get("ocs_name", ""))
        if key and key not in out:
            out[key] = row
    return out


def build_dict_kear_account_rows(filtered_rows: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Build pivot rows from FILTRED with beneficiary/account dictionary fields."""
    out: List[Dict[str, Any]] = []
    seen: set[tuple[str, str]] = set()
    for row in filtered_rows:
        cloud_type = _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type"]))
        retrieved_from = _normalize_lookup_value(row.get("Retrived from", ""))
        if cloud_type != "GEN 2" or retrieved_from != "DALI EXPORT":
            continue

        beneficiary = _normalize_cell_value(row.get("INV_Beneficiary_Account", "")).strip()
        uid = _normalize_cell_value(row.get("uid", "")).strip()
        if not beneficiary or not uid:
            continue

        dedupe_key = (_normalize_lookup_value(beneficiary), _normalize_lookup_value(uid))
        if dedupe_key in seen:
            continue
        seen.add(dedupe_key)
        out.append(
            {
                "INV_Beneficiary_Account": beneficiary,
                "INV_Beneficiary_Account_ENV": _normalize_cell_value(row.get("INV_Beneficiary_Account_ENV", "")).strip(),
                "uid": uid,
                "name": _normalize_cell_value(_get_row_value_by_candidates(row, ["name", "NAME REL"])).strip(),
                "short_label": _normalize_cell_value(_get_row_value_by_candidates(row, ["short_label", "SHORT LABEL REL"])).strip(),
                "asa": _normalize_cell_value(_get_row_value_by_candidates(row, ["asa", "ASA REL"])).strip(),
                "irt_code": _normalize_cell_value(_get_row_value_by_candidates(row, ["irt_code", "IRT CODE REL"])).strip(),
                "iappli_code": _normalize_cell_value(_get_row_value_by_candidates(row, ["iappli_code", "IAPPLI CODE REL"])).strip(),
                "trigram": _normalize_cell_value(_get_row_value_by_candidates(row, ["trigram", "TRIGRAM REL"])).strip(),
                "dsi": _normalize_cell_value(_get_row_value_by_candidates(row, ["dsi", "DSI REL"])).strip(),
                "application_management_rc": _normalize_cell_value(
                    _get_row_value_by_candidates(row, ["application_management_rc", "APPLICATION MANAGEMENT RC REL"])
                ).strip(),
                "application_development_manager": _normalize_cell_value(
                    _get_row_value_by_candidates(
                        row,
                        ["application_development_manager", "APPLICATION DEVELOPMENT MANAGER REL"],
                    )
                ).strip(),
                "DALI [APP] UID": uid,
                "DALI [APP] NAME": _normalize_cell_value(_get_row_value_by_candidates(row, ["name", "NAME REL"])).strip(),
                "DALI [APP] SHORT LABEL": _normalize_cell_value(_get_row_value_by_candidates(row, ["short_label", "SHORT LABEL REL"])).strip(),
                "DALI [APP] ASA": _normalize_cell_value(_get_row_value_by_candidates(row, ["asa", "ASA REL"])).strip(),
                "DALI [APP] IRT CODE": _normalize_cell_value(_get_row_value_by_candidates(row, ["irt_code", "IRT CODE REL"])).strip(),
                "DALI [APP] IAPPLI CODE": _normalize_cell_value(_get_row_value_by_candidates(row, ["iappli_code", "IAPPLI CODE REL"])).strip(),
                "DALI [APP] TRIGRAM": _normalize_cell_value(_get_row_value_by_candidates(row, ["trigram", "TRIGRAM REL"])).strip(),
                "DALI [APP] DSI": _normalize_cell_value(_get_row_value_by_candidates(row, ["dsi", "DSI REL"])).strip(),
                "DALI [APP] APPLICATION MANAGEMENT RC": _normalize_cell_value(
                    _get_row_value_by_candidates(row, ["application_management_rc", "APPLICATION MANAGEMENT RC REL"])
                ).strip(),
                "DALI [APP] APPLICATION DEVELOPMENT MANAGER REL": _normalize_cell_value(
                    _get_row_value_by_candidates(
                        row,
                        ["application_development_manager", "APPLICATION DEVELOPMENT MANAGER REL"],
                    )
                ).strip(),
            }
        )

    log.info("Dict_Kear_Account built rows=%s", len(out))
    return out


def _index_uid_by_beneficiary(dict_kear_account_rows: List[Dict[str, Any]]) -> Dict[str, str]:
    out: Dict[str, str] = {}
    for row in dict_kear_account_rows:
        beneficiary = _normalize_lookup_value(row.get("INV_Beneficiary_Account", ""))
        uid = _normalize_cell_value(row.get("uid", "")).strip()
        if not beneficiary or not uid:
            continue
        out.setdefault(beneficiary, uid)
    return out


def apply_kear_override_from_beneficiary(
    marley_rows: List[Dict[str, Any]],
    dict_kear_account_rows: List[Dict[str, Any]],
) -> None:
    uid_by_beneficiary = _index_uid_by_beneficiary(dict_kear_account_rows)
    for row in marley_rows:
        in_scope_value = _normalize_lookup_value(row.get("Kear in scope", ""))
        if in_scope_value == "TRUE":
            row["KEAR_OVERRIDE"] = ""
            continue
        beneficiary = _normalize_lookup_value(row.get("beneficiary", ""))
        row["KEAR_OVERRIDE"] = _normalize_cell_value(uid_by_beneficiary.get(beneficiary, ""))


def _resolve_mapping_value(
    source_sheet: str,
    source_column: str,
    marley_row: Dict[str, Any],
    inv_row: Dict[str, Any],
) -> str:
    if source_sheet in {"get_marley_by_ocsname", "get_marley_gen2_by_uuid"}:
        if not source_column:
            return ""
        if source_column == "Gen 2":
            return "Gen 2"
        return _normalize_cell_value(marley_row.get(source_column, ""))
    if source_sheet == "get_inv_by_account":
        if not source_column:
            return ""
        return _normalize_cell_value(inv_row.get(source_column, ""))
    if source_sheet == "FILTRED" and source_column == "IPLIST":
        return _normalize_cell_value(marley_row.get("IPLIST", ""))
    return ""


def append_marley_rows_to_filtered(
    filtered_rows: List[Dict[str, Any]],
    marley_rows: List[Dict[str, Any]],
    monitored_rows: List[Dict[str, str]],
    inv_by_account_rows: Optional[List[Dict[str, Any]]] = None,
    dict_kear_account_rows: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    uid_to_template: Dict[str, Dict[str, Any]] = {}
    for row in filtered_rows:
        uid = _normalize_lookup_value(row.get("uid", ""))
        if uid and uid not in uid_to_template:
            uid_to_template[uid] = dict(row)

    monitored_by_uid: Dict[str, List[Dict[str, str]]] = {}
    for row in monitored_rows:
        uid = _normalize_lookup_value(row.get("uid", ""))
        if uid:
            monitored_by_uid.setdefault(uid, []).append(row)

    existing_keys = {
        (
            _normalize_lookup_value(row.get("uid", "")),
            _normalize_lookup_value(row.get("Server UID", "")),
        )
        for row in filtered_rows
    }

    inv_by_ocs_name = _index_rows_by_ocs_name(inv_by_account_rows or [])
    uid_by_beneficiary = _index_uid_by_beneficiary(dict_kear_account_rows or [])

    appended: List[Dict[str, Any]] = []
    for marley in marley_rows:
        marley_uid = _normalize_lookup_value(marley.get("app_info.kear_uuid", ""))
        beneficiary = _normalize_lookup_value(marley.get("beneficiary", ""))
        beneficiary_uid = _normalize_lookup_value(uid_by_beneficiary.get(beneficiary, ""))
        effective_uid = beneficiary_uid or marley_uid
        server_uid = _normalize_lookup_value(marley.get("uuid", ""))
        if not effective_uid or not server_uid:
            continue
        key = (effective_uid, server_uid)
        if key in existing_keys:
            continue

        candidates = monitored_by_uid.get(effective_uid, [])
        chosen = candidates[0] if candidates else {}
        marley_iplist = _normalize_lookup_value(marley.get("IPLIST", ""))
        for candidate in candidates:
            network = _normalize_lookup_value(candidate.get("network", ""))
            if network and marley_iplist and network in marley_iplist:
                chosen = candidate
                break

        template = dict(uid_to_template.get(effective_uid, {}))
        inv_row = inv_by_ocs_name.get(_normalize_lookup_value(marley.get("ocs_name", "")), {})

        # Mandatory enrichment backbone
        template.update(
            {
                "uid": _normalize_cell_value(beneficiary_uid or marley_uid),
                "program": _normalize_cell_value(chosen.get("program", "")),
                "taken": _normalize_cell_value(chosen.get("taken", "")),
                "Server UID": _normalize_cell_value(marley.get("uuid", "")),
                "HOSTNAME": _normalize_cell_value(marley.get("ocs_name", "")),
                "hostname": _normalize_cell_value(marley.get("ocs_name", "")),
                "STATUS": "In production",
                "MAIN IP": _normalize_cell_value(marley.get("MAIN IP", "")),
                "managed": _normalize_cell_value(marley.get("managed", "")),
                "IPLIST": _normalize_cell_value(marley.get("IPLIST", "")),
                "SUBNET": _normalize_cell_value(marley.get("SUBNET", "")),
                "enforcement": _normalize_cell_value(marley.get("enforcement", "")),
                "role": _normalize_cell_value(marley.get("role", "")),
                "app": _normalize_cell_value(marley.get("app", "")),
                "env": _normalize_cell_value(marley.get("env", "")),
                "loc": _normalize_cell_value(marley.get("loc", "")),
                "Retrived from": "Enriched from Marely",
                "lookup_status": "FOUND",
                "error": "",
            }
        )

        # Apply mapping table (source -> target display/technical keys)
        for source_sheet, source_column, target_display, target_technical, _rule in MARLEY_ENRICHMENT_MAPPING_TABLE:
            value = _resolve_mapping_value(source_sheet, source_column, marley, inv_row)
            if target_display:
                template[target_display] = value
            if target_technical:
                template[target_technical] = value

        if not str(template.get("network", "")).strip():
            template["network"] = _normalize_cell_value(chosen.get("network", ""))

        appended.append(template)
        existing_keys.add(key)

    log.info("Marley append to FILTRED done appended_rows=%s", len(appended))
    return filtered_rows + appended


def _extract_monitored_app_links_from_dali(response: Dict[str, Any], monitored_uids: set[str]) -> List[Dict[str, Any]]:
    """Keep DALI edges whose trailing application UID is in monitored scope."""
    rows: List[Dict[str, Any]] = []
    result = response.get("result") if isinstance(response, dict) else None
    edges = [edge for edge in (result or []) if isinstance(edge, dict)] if isinstance(result, list) else []

    log.info("DALI application-link extraction from additional lookup edges=%s", len(edges))
    for edge in edges:
        leading_props = node_properties_to_dict(edge.get("leading_node"))
        trailing_props = node_properties_to_dict(edge.get("trailing_node"))

        app_uid = str(trailing_props.get("uid") or "").strip()
        if not app_uid or app_uid not in monitored_uids:
            continue

        server_hostname = _normalize_cell_value(leading_props.get("hostname")) or _normalize_cell_value(trailing_props.get("hostname"))
        server_cloud_type = _normalize_cell_value(leading_props.get("cloud_type")) or _normalize_cell_value(trailing_props.get("cloud_type"))
        server_uid = _normalize_cell_value(leading_props.get("uid")) or _normalize_cell_value(trailing_props.get("uid"))

        rows.append(
            {
                "uid": app_uid,
                "hostname": server_hostname,
                "cloud_type": server_cloud_type,
                "server_uid": server_uid,
            }
        )

    log.info("DALI application-link extraction kept_rows=%s (monitored application uids)", len(rows))
    return rows


def discover_additional_servers_from_inventory_accounts(
    # Beneficiary-driven expansion: inventory -> DALI -> monitored UID filtering
    client: DaliImpactAnalysisClient,
    d4s_client: Data4secClient,
    filtered_rows: List[Dict[str, Any]],
    monitored_uids: set[str],
    impact_endpoint: str,
    limit: Optional[int],
    depth_until: Optional[int],
    filters: Optional[Dict[str, str]] = None,
    inventory_by_account_rows: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    accounts_not_to_enrich_tokens = _parse_filter_tokens(filters, "FILTER_OWNER_ACCOUNTS_NOT_TO_ENRICH")
    beneficiary_values = {
        _normalize_lookup_value(row.get("INV_Beneficiary_Account", ""))
        for row in filtered_rows
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type", "CLOUD TYPE"])) == "GEN 2"
        if str(row.get("INV_Beneficiary_Account", "")).strip() not in {"", "NOT_FOUND", "NOT_GEN2"}
    }
    if accounts_not_to_enrich_tokens:
        beneficiary_values = {
            beneficiary
            for beneficiary in beneficiary_values
            if beneficiary and not _matches_exact_token(beneficiary, accounts_not_to_enrich_tokens)
        }

    if not beneficiary_values:
        log.info(
            "Additional inventory-account discovery skipped: no eligible beneficiary account available excluded_accounts=%s",
            accounts_not_to_enrich_tokens,
        )
        return []
    log.info(
        "Additional inventory-account discovery start distinct_beneficiaries=%s excluded_accounts=%s",
        len(beneficiary_values),
        accounts_not_to_enrich_tokens,
    )

    inventory_by_beneficiary = query_inventory_for_beneficiaries(d4s_client, sorted(beneficiary_values))
    inventory_docs: List[Dict[str, Any]] = []
    for beneficiary, docs in inventory_by_beneficiary.items():
        for doc in docs:
            owner_account_value = _normalize_cell_value(doc.get("owner_app_name"))
            if accounts_not_to_enrich_tokens and _matches_exact_token(owner_account_value, accounts_not_to_enrich_tokens):
                continue
            if inventory_by_account_rows is not None:
                hostid_value = _normalize_cell_value(doc.get("hostid"))
                srn_value = _normalize_cell_value(doc.get("srn"))
                inventory_by_account_rows.append(
                    {
                        "input_INV_Beneficiary_Account": beneficiary,
                        "beneficiary": beneficiary,
                        "ocs_name": _normalize_cell_value(doc.get("ocs_name")),
                        "hostname": _short_hostname(_normalize_cell_value(doc.get("hostname"))),
                        "status": _normalize_status(doc.get("status")),
                        "hostid": hostid_value,
                        "Normalized_uuid_from_hostid": _normalize_uuid_from_hostid(hostid_value),
                        "srn": srn_value,
                        "Normalized_uuid_from_srn": _normalize_uuid_from_srn(srn_value),
                        "owner_app_name": owner_account_value,
                        "ip": _normalize_cell_value(doc.get("ip")),
                        "service_name": _normalize_cell_value(doc.get("service_name")),
                    }
                )
            inventory_docs.append(doc)
    inventory_docs = _deduplicate_docs(inventory_docs)
    log.info("Additional inventory-account discovery inventory_docs=%s", len(inventory_docs))

    log.info(
        "Additional inventory-account DALI lookup disabled: keeping inventory_by_account_rows only docs=%s",
        len(inventory_docs),
    )
    return []


def enrich_filtered_rows_with_inventory(
    # Main enrichment path for FILTRED: Gen2 inventory + optional discovered rows
    filtered_rows: List[Dict[str, Any]],
    client: DaliImpactAnalysisClient,
    impact_endpoint: str,
    limit: Optional[int],
    depth_until: Optional[int],
    monitored_uids: set[str],
    raw_server_uids: Optional[set[str]] = None,
    filters: Optional[Dict[str, str]] = None,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]], List[Dict[str, Any]], List[Dict[str, Any]], List[Dict[str, Any]]]:
    server_uids_to_query: List[str] = []
    row_contexts: List[Tuple[Dict[str, Any], str, str]] = []
    d4s_client = Data4secClient()
    log.info("Inventory enrichment start filtered_rows=%s monitored_uids=%s", len(filtered_rows), len(monitored_uids))

    for row in filtered_rows:
        cloud_type = _get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type"])
        server_uid = _get_row_value_by_candidates(row, ["server_uid", "server uid", "Server UID"])
        row_contexts.append((row, cloud_type, server_uid))

        is_gen2 = _normalize_lookup_value(cloud_type) == "GEN 2"
        if is_gen2 and _normalize_lookup_value(server_uid):
            server_uids_to_query.append(server_uid)

    unique_server_uids = sorted({
        _normalize_lookup_value(uid)
        for uid in server_uids_to_query
        if _normalize_lookup_value(uid)
    })
    log.info(
        "Inventory enrichment query optimization server_uid_rows=%s unique_server_uids=%s",
        len(server_uids_to_query),
        len(unique_server_uids),
    )
    inventory_map = query_inventory_for_server_uids(client=d4s_client, server_uids=unique_server_uids)

    for row, cloud_type, server_uid in row_contexts:
        is_gen2 = _normalize_lookup_value(cloud_type) == "GEN 2"
        if not is_gen2:
            for column in INVENTORY_HEADERS:
                row[column] = "NOT_GEN2"
            continue

        normalized_server_uid = _normalize_lookup_value(server_uid)
        inventory_row = inventory_map.get(normalized_server_uid, {}) if normalized_server_uid else {}

        if not inventory_row:
            row["INV_ocs_name"] = "NOT_FOUND"
            row["INV_status"] = "NOT_FOUND"
            row["INV_hostname"] = "NOT_FOUND"
            row["Retrived from"] = "NOT_FOUND"
            row["INV_Owner_Account"] = "NOT_FOUND"
            row["INV_Beneficiary_Account"] = "NOT_FOUND"
            continue

        for column in INVENTORY_HEADERS:
            row[column] = inventory_row.get(column, "")

    beneficiary_accounts = sorted(
        {
            _normalize_lookup_value(row.get("INV_Beneficiary_Account", ""))
            for row, cloud_type, _server_uid in row_contexts
            if _normalize_lookup_value(cloud_type) == "GEN 2"
            and str(row.get("INV_Beneficiary_Account", "")).strip() not in {"", "NOT_FOUND", "NOT_GEN2"}
        }
    )
    beneficiary_env_map = query_platform_accounts_env_by_names(d4s_client, beneficiary_accounts)
    for row, cloud_type, _server_uid in row_contexts:
        is_gen2 = _normalize_lookup_value(cloud_type) == "GEN 2"
        if not is_gen2:
            row["INV_Beneficiary_Account_ENV"] = "NOT_GEN2"
            continue
        beneficiary_account = _normalize_lookup_value(row.get("INV_Beneficiary_Account", ""))
        row["INV_Beneficiary_Account_ENV"] = beneficiary_env_map.get(beneficiary_account, "NOT_FOUND")

    inventory_by_account_rows: List[Dict[str, Any]] = []
    marley_by_ocsname_rows: List[Dict[str, Any]] = []
    marley_gen2_by_uuid_rows: List[Dict[str, Any]] = []

    discovered_rows = discover_additional_servers_from_inventory_accounts(
        client=client,
        d4s_client=d4s_client,
        filtered_rows=filtered_rows,
        monitored_uids=monitored_uids,
        impact_endpoint=impact_endpoint,
        limit=limit,
        depth_until=depth_until,
        filters=filters,
        inventory_by_account_rows=inventory_by_account_rows,
    )
    filtered_rows.extend(discovered_rows)

    normalized_raw_server_uids = {
        _normalize_lookup_value(value)
        for value in (raw_server_uids or set())
        if _normalize_lookup_value(value)
    }
    for row in inventory_by_account_rows:
        normalized_from_hostid = _normalize_lookup_value(row.get("Normalized_uuid_from_hostid", ""))
        row["asset_origin"] = "EXISTING_IN_RAW_IMPORT" if normalized_from_hostid and normalized_from_hostid in normalized_raw_server_uids else "ENRICHED_NEW_ASSET"

    beneficiary_not_taken_tokens = [
        _normalize_lookup_value(token)
        for token in _get_beneficiary_not_taken_tokens(filters)
        if _normalize_lookup_value(token)
    ]
    before_beneficiary_exclusion_count = len(filtered_rows)
    if beneficiary_not_taken_tokens:
        filtered_rows = [
            row
            for row in filtered_rows
            if not _contains_any_token(row.get("INV_Beneficiary_Account", ""), beneficiary_not_taken_tokens)
        ]
    removed_beneficiary_exclusion_count = before_beneficiary_exclusion_count - len(filtered_rows)
    log.info(
        "Inventory enrichment beneficiary exclusion tokens=%s removed_rows=%s kept_rows=%s",
        beneficiary_not_taken_tokens,
        removed_beneficiary_exclusion_count,
        len(filtered_rows),
    )

    _recompute_prd_env_flags(filtered_rows, filters)
    prod_tokens = _get_prod_beneficiary_tokens(filters)
    before_prod_filter_count = len(filtered_rows)
    filtered_rows = [
        row
        for row in filtered_rows
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type"])) != "GEN 2"
        or _contains_any_token(_effective_gen2_prd_env_value(row), prod_tokens)
    ]
    removed_non_prod_count = before_prod_filter_count - len(filtered_rows)
    log.info(
        "Inventory enrichment prod beneficiary filter tokens=%s scope=GEN2 removed_rows=%s kept_rows=%s",
        prod_tokens,
        removed_non_prod_count,
        len(filtered_rows),
    )

    log.info(
        "Inventory enrichment done base_rows=%s discovered_rows=%s total_rows=%s",
        len(row_contexts),
        len(discovered_rows),
        len(filtered_rows),
    )

    marley_rows_for_append: List[Dict[str, Any]] = []

    marley_source_rows: List[Dict[str, Any]] = []
    seen_marley_input_keys: set[tuple[str, str, str]] = set()
    for row in inventory_by_account_rows:
        if _normalize_lookup_value(row.get("asset_origin", "")) != "ENRICHED_NEW_ASSET":
            continue
        hostid_uuid = _normalize_lookup_value(row.get("Normalized_uuid_from_hostid", ""))
        srn_uuid = _normalize_lookup_value(row.get("Normalized_uuid_from_srn", ""))
        beneficiary = _normalize_lookup_value(row.get("beneficiary", ""))
        dedupe_key = (hostid_uuid, srn_uuid, beneficiary)
        if dedupe_key in seen_marley_input_keys:
            continue
        seen_marley_input_keys.add(dedupe_key)
        marley_source_rows.append(row)

    hostid_candidate_rows = [
        row for row in marley_source_rows if _normalize_cell_value(row.get("Normalized_uuid_from_hostid", ""))
    ]
    marley_hostid_uuids = sorted(
        {
            _normalize_cell_value(row.get("Normalized_uuid_from_hostid", ""))
            for row in hostid_candidate_rows
            if _normalize_cell_value(row.get("Normalized_uuid_from_hostid", ""))
        }
    )
    marley_docs_by_uuid = query_marley_original_by_uuids(d4s_client, marley_hostid_uuids)
    hostid_matched_rows = [
        row
        for row in hostid_candidate_rows
        if marley_docs_by_uuid.get(_normalize_lookup_value(row.get("Normalized_uuid_from_hostid", "")))
    ]
    marley_gen2_by_uuid_rows = build_marley_sheet_rows(
        inventory_by_account_rows=hostid_matched_rows,
        marley_docs_by_lookup=marley_docs_by_uuid,
        monitored_uids=monitored_uids,
        lookup_source_field="Normalized_uuid_from_hostid",
        lookup_output_field="lookup_uuid",
    )
    missing_hostid_rows = [
        row
        for row in marley_source_rows
        if not marley_docs_by_uuid.get(_normalize_lookup_value(row.get("Normalized_uuid_from_hostid", "")))
    ]
    marley_srn_uuids = sorted(
        {
            _normalize_cell_value(row.get("Normalized_uuid_from_srn", ""))
            for row in missing_hostid_rows
            if _normalize_cell_value(row.get("Normalized_uuid_from_srn", ""))
        }
    )
    if marley_srn_uuids:
        marley_docs_by_srn = query_marley_original_by_uuids(d4s_client, marley_srn_uuids)
        marley_gen2_by_uuid_rows.extend(
            build_marley_sheet_rows(
                inventory_by_account_rows=missing_hostid_rows,
                marley_docs_by_lookup=marley_docs_by_srn,
                monitored_uids=monitored_uids,
                lookup_source_field="Normalized_uuid_from_srn",
                lookup_output_field="lookup_uuid",
            )
        )
    marley_gen2_by_uuid_rows, marley_rows_for_append = filter_marley_sheet_rows(
        marley_rows=marley_gen2_by_uuid_rows,
        filtered_rows=filtered_rows,
        filters=filters,
    )
    log.info(
        "Marley UUID sheet build done source_new_assets=%s source_hostid_uuids=%s source_srn_uuids=%s output_rows=%s",
        len(marley_source_rows),
        len(marley_hostid_uuids),
        len(marley_srn_uuids),
        len(marley_gen2_by_uuid_rows),
    )

    return filtered_rows, inventory_by_account_rows, marley_by_ocsname_rows, marley_gen2_by_uuid_rows, marley_rows_for_append


def enrich_rows_with_inventory_for_gen2(rows: List[Dict[str, Any]], filters: Optional[Dict[str, str]] = None) -> None:
    """Populate inventory columns for every Gen2 row in the provided collection."""
    d4s_client = Data4secClient()
    server_uids = sorted(
        {
            _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
            for row in rows
            if _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type", "CLOUD TYPE"])) == "GEN 2"
            and _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
        }
    )
    if not server_uids:
        log.info("RAW inventory enrichment skipped: no GEN2 server uid found")
        return

    log.info("RAW inventory enrichment start gen2_servers=%s", len(server_uids))
    inventory_map = query_inventory_for_server_uids(client=d4s_client, server_uids=server_uids)

    for row in rows:
        is_gen2 = _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type", "CLOUD TYPE"])) == "GEN 2"
        if not is_gen2:
            continue
        server_uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
        inventory_row = inventory_map.get(server_uid, {}) if server_uid else {}
        if not inventory_row:
            row["INV_ocs_name"] = row.get("INV_ocs_name", "NOT_FOUND") or "NOT_FOUND"
            row["INV_status"] = row.get("INV_status", "NOT_FOUND") or "NOT_FOUND"
            row["INV_hostname"] = row.get("INV_hostname", "NOT_FOUND") or "NOT_FOUND"
            row["INV_Owner_Account"] = row.get("INV_Owner_Account", "NOT_FOUND") or "NOT_FOUND"
            row["INV_Beneficiary_Account"] = row.get("INV_Beneficiary_Account", "NOT_FOUND") or "NOT_FOUND"
            continue
        for column in INVENTORY_HEADERS:
            if not str(row.get(column, "")).strip():
                row[column] = inventory_row.get(column, "")

    beneficiary_accounts = sorted(
        {
            _normalize_lookup_value(row.get("INV_Beneficiary_Account", ""))
            for row in rows
            if _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type", "CLOUD TYPE"])) == "GEN 2"
            and str(row.get("INV_Beneficiary_Account", "")).strip() not in {"", "NOT_FOUND", "NOT_GEN2"}
        }
    )
    beneficiary_env_map = query_platform_accounts_env_by_names(d4s_client, beneficiary_accounts)
    for row in rows:
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type", "CLOUD TYPE"])) != "GEN 2":
            row["INV_Beneficiary_Account_ENV"] = row.get("INV_Beneficiary_Account_ENV", "NOT_GEN2") or "NOT_GEN2"
            continue
        beneficiary_account = _normalize_lookup_value(row.get("INV_Beneficiary_Account", ""))
        row["INV_Beneficiary_Account_ENV"] = beneficiary_env_map.get(beneficiary_account, row.get("INV_Beneficiary_Account_ENV", "NOT_FOUND") or "NOT_FOUND")

    _recompute_prd_env_flags(rows, filters)
    log.info("RAW inventory enrichment done")


def extract_rows_from_response(
    response: Dict[str, Any],
    base_row: Dict[str, Any],
    mappings: List[Tuple[str, str]],
    err_text: str,
    filters: Optional[Dict[str, str]] = None,
    apply_filters: bool = True,
) -> List[Dict[str, Any]]:
    def _mapping_alias_keys(dali_attr: str) -> List[str]:
        raw_attr = str(dali_attr or "").strip()
        if not raw_attr:
            return []
        aliases = [raw_attr]
        underscore_alias = raw_attr.replace(".", "_")
        if underscore_alias not in aliases:
            aliases.append(underscore_alias)
        return aliases

    raw_debug_fieldnames = [name for pair in RAW_FILTER_COLUMN_PAIRS for name in pair[:2] if name] + ["F_FILTER_ALL"]
    if err_text:
        row = dict(base_row)
        row.update({
            "lookup_status": "ERROR",
            "count": 0,
            "error": err_text,
        })
        for display_name, _ in mappings:
            row[display_name] = ""
        for _display_name, dali_attr in mappings:
            for alias in _mapping_alias_keys(dali_attr):
                row.setdefault(alias, "")
        for field in raw_debug_fieldnames:
            row[field] = ""
        return [row]

    result = response.get("result") if isinstance(response, dict) else None
    edges = [edge for edge in (result or []) if isinstance(edge, dict)] if isinstance(result, list) else []
    count_value = response.get("count", 0) if isinstance(response, dict) else 0

    if not edges:
        row = dict(base_row)
        row.update({
            "lookup_status": "NOT_FOUND",
            "count": count_value,
            "error": "",
        })
        for display_name, _ in mappings:
            row[display_name] = ""
        for _display_name, dali_attr in mappings:
            for alias in _mapping_alias_keys(dali_attr):
                row.setdefault(alias, "")
        for field in raw_debug_fieldnames:
            row[field] = ""
        return [row]

    out_rows: List[Dict[str, Any]] = []
    for edge in edges:
        leading_node = edge.get("leading_node") if isinstance(edge, dict) else None
        trailing_node = edge.get("trailing_node") if isinstance(edge, dict) else None
        lead = node_properties_to_dict(edge.get("leading_node"))
        trail = node_properties_to_dict(edge.get("trailing_node"))
        row = dict(base_row)
        row["Server UID"] = _extract_server_uid_from_edge(edge)
        row.update({
            "lookup_status": "FOUND",
            "count": count_value,
            "error": "",
        })
        for display_name, dali_attr in mappings:
            raw_value = _resolve_edge_mapping_value(edge=edge, dali_attr=dali_attr, base_row=base_row)
            normalized_value = _normalize_cell_value(raw_value)
            row[display_name] = normalized_value
            for alias in _mapping_alias_keys(dali_attr):
                if not str(row.get(alias, "")).strip():
                    row[alias] = normalized_value
        row.update(
            _raw_filter_debug_columns(
                lead=lead,
                trail=trail,
                filters=filters,
                leading_node=leading_node if isinstance(leading_node, dict) else None,
                trailing_node=trailing_node if isinstance(trailing_node, dict) else None,
                row=row,
            )
        )
        if (not apply_filters) or _edge_matches_filters(
            lead=lead,
            trail=trail,
            filters=filters,
            leading_node=leading_node if isinstance(leading_node, dict) else None,
            trailing_node=trailing_node if isinstance(trailing_node, dict) else None,
            row=row,
        ):
            out_rows.append(row)
    return out_rows


def write_output_csv(
    output_file: str,
    rows: List[Dict[str, Any]],
    mappings: List[Tuple[str, str]],
    extra_fieldnames: Optional[List[str]] = None,
    base_fieldnames: Optional[List[str]] = None,
) -> None:
    Path(output_file).parent.mkdir(parents=True, exist_ok=True)
    effective_base = ["uid", "program", "network", "taken", "Server UID"] if base_fieldnames is None else list(base_fieldnames)
    fieldnames = effective_base + [display for display, _ in mappings] + (extra_fieldnames or [])
    with open(output_file, "w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames, extrasaction="ignore")
        writer.writeheader()
        writer.writerows(rows)


def _xlsx_col_ref(index: int) -> str:
    ref = ""
    idx = index + 1
    while idx > 0:
        idx, rem = divmod(idx - 1, 26)
        ref = chr(65 + rem) + ref
    return ref


def _compute_col_widths(rows: List[List[str]], min_width: float = 10.0, max_width: float = 60.0) -> List[float]:
    if not rows:
        return []
    max_cols = max(len(row) for row in rows)
    widths: List[float] = []
    for col_idx in range(max_cols):
        max_len = 0
        for row in rows:
            value = row[col_idx] if col_idx < len(row) else ""
            max_len = max(max_len, len(str(value or "")))
        widths.append(min(max_width, max(min_width, float(max_len + 2))))
    return widths


def _xlsx_cols_xml(widths: List[float]) -> str:
    if not widths:
        return ""
    cols = []
    for idx, width in enumerate(widths, start=1):
        cols.append(f'<col min="{idx}" max="{idx}" width="{width:.2f}" customWidth="1"/>')
    return '<cols>' + ''.join(cols) + '</cols>'




def _xml_safe_text(value: Any) -> str:
    text = str(value or "")
    sanitized = "".join(
        ch
        for ch in text
        if ch in {"\t", "\n", "\r"}
        or (0x20 <= ord(ch) <= 0xD7FF)
        or (0xE000 <= ord(ch) <= 0xFFFD)
        or (0x10000 <= ord(ch) <= 0x10FFFF)
    )
    return escape(sanitized)


def _xlsx_autofilter_xml(
    row_count: int,
    col_count: int,
    fieldnames: Optional[List[str]] = None,
    filter_criteria: Optional[Dict[str, List[str]]] = None,
) -> str:
    if col_count <= 0:
        return ""
    start_ref = "A1"
    end_ref = f"{_xlsx_col_ref(col_count - 1)}{max(1, row_count + 1)}"
    if not filter_criteria or not fieldnames:
        return f'<autoFilter ref="{start_ref}:{end_ref}"/>'

    filter_parts: List[str] = []
    for fieldname, allowed_values in filter_criteria.items():
        if fieldname not in fieldnames:
            continue
        normalized_values = [str(value or "") for value in allowed_values if str(value or "") != ""]
        if not normalized_values:
            continue
        col_id = fieldnames.index(fieldname)
        values_xml = ''.join(f'<filter val="{_xml_safe_text(value)}"/>' for value in normalized_values)
        filter_parts.append(f'<filterColumn colId="{col_id}"><filters blank="0">{values_xml}</filters></filterColumn>')

    if not filter_parts:
        return f'<autoFilter ref="{start_ref}:{end_ref}"/>'
    return f'<autoFilter ref="{start_ref}:{end_ref}">' + ''.join(filter_parts) + '</autoFilter>'


def _coerce_excel_numeric(value: Any) -> Optional[str]:
    if isinstance(value, bool):
        return None
    if isinstance(value, int):
        return str(value)
    if isinstance(value, float):
        return str(value) if math.isfinite(value) else None

    raw = str(value or "").strip()
    if not raw:
        return None

    normalized = raw.replace(" ", "")
    if re.fullmatch(r"[+-]?\d+", normalized):
        stripped = normalized.lstrip("+-")
        if len(stripped) > 1 and stripped.startswith("0"):
            return None
        return normalized

    if re.fullmatch(r"[+-]?\d+[\.,]\d+", normalized):
        return normalized.replace(",", ".")

    return None


def _ratio_percent_from_label(value: Any) -> float:
    raw = str(value or "").strip()
    if not raw:
        return float("inf")

    # Expected format: "(x/y) nn,nn%". Keep only the percentage part after the last space.
    pct_candidate = raw.split()[-1].replace(",", ".")
    match = re.fullmatch(r"(-?\d+(?:\.\d+)?)%", pct_candidate)
    if not match:
        match = re.search(r"(-?\d+(?:\.\d+)?)\s*%", raw.replace(",", "."))
    if not match:
        return float("inf")
    try:
        return float(match.group(1))
    except ValueError:
        return float("inf")


def _variation_percent_from_label(value: Any) -> float:
    raw = str(value or "").strip().replace(",", ".")
    match = re.search(r"([+-]?\d+(?:\.\d+)?)\s*%", raw)
    if not match:
        return 0.0
    try:
        return float(match.group(1))
    except ValueError:
        return 0.0


def _append_stats_visual_columns(
    rows: List[Dict[str, Any]],
    headers: List[str],
    previous_stats_by_key: Optional[Dict[Tuple[str, str], Dict[str, str]]] = None,
) -> Tuple[List[Dict[str, Any]], List[str]]:
    pct_columns = [
        "% servers with illumio installed",
        "% servers with illumio installed (Enriched)",
        "% servers with illumio agent in blocking mode",
        "% servers with illumio agent in blocking mode (Enriched)",
    ]
    previous_stats_by_key = previous_stats_by_key or {}

    for row in rows:
        row_program = str(row.get("Program", "") or "").strip()
        row_uid = _normalize_lookup_value(row.get("Kear ID", ""))
        previous_row = previous_stats_by_key.get((row_program, row_uid), {})

        for pct_column in pct_columns:
            indicator_col = f"{pct_column} Indicator Icon"
            trend_col = f"{pct_column} Trend Icon"

            pct_value = _ratio_percent_from_label(row.get(pct_column, ""))
            row[indicator_col] = "" if pct_value == float("inf") else f"{pct_value:.2f}"

            previous_pct = _ratio_percent_from_label(previous_row.get(pct_column, "")) if previous_row else float("inf")
            if pct_value == float("inf") or previous_pct == float("inf"):
                row[trend_col] = ""
            else:
                row[trend_col] = f"{(pct_value - previous_pct):.2f}"

    extended_headers: List[str] = []
    for header in headers:
        extended_headers.append(header)
        if header in pct_columns:
            extended_headers.append(f"{header} Indicator Icon")
            extended_headers.append(f"{header} Trend Icon")

    return rows, extended_headers


def _append_total_directional_columns(rows: List[Dict[str, Any]], headers: List[str]) -> Tuple[List[Dict[str, Any]], List[str]]:
    trend_map = {
        "% servers with illumio installed": "Variation % servers with illumio installed",
        "% servers with illumio agent in blocking mode": "Variation % servers with illumio agent in blocking mode",
    }

    extended_headers: List[str] = []
    for header in headers:
        extended_headers.append(header)
        base_name = _base_field_name(header)
        if base_name in trend_map:
            extended_headers.append(f"{base_name} Trend Icon")

    for row in rows:
        for pct_col, var_base in trend_map.items():
            trend_col = f"{pct_col} Trend Icon"
            variation_key = next((k for k in row.keys() if _base_field_name(k) == var_base), "")
            row[trend_col] = f"{_variation_percent_from_label(row.get(variation_key, '')):.2f}" if variation_key else ""

    return rows, extended_headers


def _xlsx_conditional_formatting_xml(fieldnames: List[str], row_count: int) -> str:
    """
    FINAL VERSION — Indicator = legacy 3TrafficLights1
                    Trend    = x14 custom 3Triangles
    """

    if row_count <= 0:
        return ""

    legacy_rules: List[str] = []
    x14_rules: List[str] = []
    priority = 1
    end_row = row_count + 1

    for col_idx, field in enumerate(fieldnames):
        field_str = str(field).strip()
        col_ref = _xlsx_col_ref(col_idx)
        sqref = f"{col_ref}2:{col_ref}{end_row}"

        # ============================================================
        # ✅ INDICATOR ICONS (Traffic Lights)
        # EXACT legacy OOXML extracted from your working Excel file
        # ============================================================
        if "Indicator Icon" in field_str:
            legacy_rules.append(
                f'<conditionalFormatting sqref="{sqref}">'
                f'  <cfRule type="iconSet" priority="{priority}">'
                f'    <iconSet iconSet="3TrafficLights1" showValue="0">'
                f'      <cfvo type="percent" val="0"/>'
                f'      <cfvo type="percent" val="0" gte="0"/>'
                f'      <cfvo type="percent" val="100"/>'
                f'    </iconSet>'
                f'  </cfRule>'
                f'</conditionalFormatting>'
            )
            priority += 1
            continue

        # ============================================================
        # ✅ TREND ICONS (Triangles, x14 custom)
        # EXACT x14 structure extracted from your working Excel file
        # ============================================================
        if "Trend Icon" in field_str:
            rule_id = "{" + str(uuid.uuid4()).upper() + "}"
            x14_rules.append(
                f'<x14:conditionalFormatting xmlns:xm="http://schemas.microsoft.com/office/excel/2006/main">'
                f'  <x14:cfRule type="iconSet" priority="{priority}" id="{rule_id}">'
                f'    <x14:iconSet iconSet="3Triangles" custom="1" showValue="0">'
                f'      <x14:cfvo type="percent"><xm:f>0</xm:f></x14:cfvo>'
                f'      <x14:cfvo type="num"><xm:f>0</xm:f></x14:cfvo>'
                f'      <x14:cfvo type="num" gte="0"><xm:f>0</xm:f></x14:cfvo>'
                f'      <x14:cfIcon iconSet="3Triangles" iconId="0"/>'
                f'      <x14:cfIcon iconSet="3Triangles" iconId="1"/>'
                f'      <x14:cfIcon iconSet="3Triangles" iconId="2"/>'
                f'    </x14:iconSet>'
                f'  </x14:cfRule>'
                f'  <xm:sqref>{sqref}</xm:sqref>'
                f'</x14:conditionalFormatting>'
            )
            priority += 1
            continue

    # ============================================================
    # ✅ ASSEMBLAGE FINAL
    # Legacy rules go in worksheet root
    # Trend rules must be wrapped in extLst/x14:conditionalFormattings
    # ============================================================

    xml = "".join(legacy_rules)

    if x14_rules:
        xml += (
            '<extLst>'
            '  <ext uri="{78C0D931-6437-407d-A8EE-F0AAD7539E65}" '
            '       xmlns:x14="http://schemas.microsoft.com/office/spreadsheetml/2009/9/main">'
            '    <x14:conditionalFormattings>'
            + "".join(x14_rules)
            + '    </x14:conditionalFormattings>'
            '  </ext>'
            '</extLst>'
        )

    return xml


def _xlsx_sheet_xml_table(
    rows: List[Dict[str, Any]],
    fieldnames: List[str],
    shaded_columns: Optional[set[str]] = None,
    enriched_columns: Optional[set[str]] = None,
    header_multiline: bool = False,
    header_height: float = 40.0,
    fixed_widths: Optional[List[Optional[float]]] = None,
    hidden_header_columns: Optional[set[str]] = None,
    filter_criteria: Optional[Dict[str, List[str]]] = None,
) -> str:
    shaded_columns = shaded_columns or set()
    enriched_columns = enriched_columns or set()
    hidden_header_columns = hidden_header_columns or set()
    matrix: List[List[str]] = [fieldnames]
    for row in rows:
        matrix.append([str(row.get(field, "") or "") for field in fieldnames])

    sheet_rows: List[str] = []
    for row_idx, row_values in enumerate(matrix, start=1):
        cells: List[str] = []
        for col_idx, value in enumerate(row_values):
            col_ref = _xlsx_col_ref(col_idx)
            fieldname = fieldnames[col_idx] if col_idx < len(fieldnames) else ""
            is_shaded = fieldname in shaded_columns
            is_enriched = fieldname in enriched_columns

            if row_idx == 1:
                if fieldname in hidden_header_columns and header_multiline:
                    style_id = "22" if is_enriched else "21"
                elif is_enriched:
                    style_id = "17" if header_multiline else "10"
                elif is_shaded:
                    style_id = "16" if header_multiline else "4"
                else:
                    style_id = "15" if header_multiline else "1"
                header_text = "" if (fieldname in hidden_header_columns and header_multiline) else _xml_safe_text(value)
                cells.append(f'<c r="{col_ref}{row_idx}" s="{style_id}" t="inlineStr"><is><t>{header_text}</t></is></c>')
                continue

            numeric_value = _coerce_excel_numeric(value)
            right_aligned_column = fieldname in RECAP_RIGHT_ALIGNED_COLUMNS
            variation_column = _is_variation_column(fieldname)
            if numeric_value is not None:
                if variation_column:
                    style_id = "20" if is_enriched else ("19" if is_shaded else "18")
                    cells.append(f'<c r="{col_ref}{row_idx}" s="{style_id}" t="n"><v>{numeric_value}</v></c>')
                else:
                    style_id = "12" if is_enriched else ("6" if is_shaded else "5")
                    cells.append(f'<c r="{col_ref}{row_idx}" s="{style_id}" t="n"><v>{numeric_value}</v></c>')
            else:
                if variation_column:
                    style_id = "20" if is_enriched else ("19" if is_shaded else "18")
                elif right_aligned_column:
                    style_id = "14" if is_enriched else ("8" if is_shaded else "7")
                else:
                    style_id = "11" if is_enriched else ("3" if is_shaded else "0")
                header_text = "" if (fieldname in hidden_header_columns and header_multiline) else _xml_safe_text(value)
                cells.append(f'<c r="{col_ref}{row_idx}" s="{style_id}" t="inlineStr"><is><t>{header_text}</t></is></c>')
        row_attrs = f' r="{row_idx}" ht="{header_height:.0f}" customHeight="1"' if row_idx == 1 and header_multiline else f' r="{row_idx}"'
        sheet_rows.append(f'<row{row_attrs}>' + ''.join(cells) + '</row>')

    computed_widths = _compute_col_widths(matrix)
    if fixed_widths is None:
        effective_widths: List[float] = computed_widths
    else:
        effective_widths = [
            fixed_widths[idx] if idx < len(fixed_widths) and fixed_widths[idx] is not None else computed_widths[idx]
            for idx in range(len(computed_widths))
        ]

    freeze_header_xml = '<sheetViews><sheetView workbookViewId="0"><pane ySplit="1" topLeftCell="A2" activePane="bottomLeft" state="frozen"/></sheetView></sheetViews>'
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        + freeze_header_xml
        + _xlsx_cols_xml(effective_widths)
        + '<sheetData>' + ''.join(sheet_rows) + '</sheetData>'
        + _xlsx_autofilter_xml(row_count=len(rows), col_count=len(fieldnames), fieldnames=fieldnames, filter_criteria=filter_criteria)
        + _xlsx_conditional_formatting_xml(fieldnames, len(rows))
        + '</worksheet>'
    )


def _xlsx_sheet_xml_summary(summary_rows: List[Tuple[str, str]], fixed_widths: Optional[List[float]] = None) -> str:
    matrix = [[left, right] for left, right in summary_rows]
    sheet_rows: List[str] = []
    for row_idx, (left, right) in enumerate(matrix, start=1):
        is_section = str(left).strip().startswith("Section ")
        if is_section:
            cells = [
                f'<c r="A{row_idx}" s="2" t="inlineStr"><is><t>{_xml_safe_text(left)}</t></is></c>',
                f'<c r="B{row_idx}" s="2" t="inlineStr"><is><t>{_xml_safe_text(right)}</t></is></c>',
            ]
            sheet_rows.append(f'<row r="{row_idx}">' + ''.join(cells) + '</row>')
            continue

        left_num = _coerce_excel_numeric(left)
        right_num = _coerce_excel_numeric(right)
        left_cell = (
            f'<c r="A{row_idx}" s="5" t="n"><v>{left_num}</v></c>'
            if left_num is not None
            else f'<c r="A{row_idx}" s="0" t="inlineStr"><is><t>{_xml_safe_text(left)}</t></is></c>'
        )
        right_cell = (
            f'<c r="B{row_idx}" s="5" t="n"><v>{right_num}</v></c>'
            if right_num is not None
            else f'<c r="B{row_idx}" s="0" t="inlineStr"><is><t>{_xml_safe_text(right)}</t></is></c>'
        )
        sheet_rows.append(f'<row r="{row_idx}">' + left_cell + right_cell + '</row>')

    freeze_header_xml = '<sheetViews><sheetView workbookViewId="0"><pane ySplit="1" topLeftCell="A2" activePane="bottomLeft" state="frozen"/></sheetView></sheetViews>'
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        + freeze_header_xml
        + _xlsx_cols_xml(fixed_widths or _compute_col_widths(matrix))
        + '<sheetData>' + ''.join(sheet_rows) + '</sheetData>'
        + '</worksheet>'
    )


def _fieldnames_for_rows(rows: List[Dict[str, Any]]) -> List[str]:
    ordered: List[str] = []
    seen = set()
    for row in rows:
        for key in row.keys():
            if key not in seen:
                seen.add(key)
                ordered.append(key)
    return ordered


def _ordered_fieldnames_with_preferred(rows: List[Dict[str, Any]], preferred_columns: List[str]) -> List[str]:
    available = _fieldnames_for_rows(rows)
    ordered = list(preferred_columns)
    ordered.extend([column for column in available if column not in ordered])
    return ordered


def _ordered_fieldnames_with_filter_tail(rows: List[Dict[str, Any]], preferred_columns: List[str]) -> List[str]:
    available = _fieldnames_for_rows(rows)
    preferred = [column for column in preferred_columns if column in available]
    remaining = [column for column in available if column not in preferred]
    filter_columns = [column for column in remaining if str(column).startswith("F_")]
    non_filter_columns = [column for column in remaining if column not in filter_columns]
    return preferred + non_filter_columns + filter_columns


def build_filtered_output_fieldnames(mappings: List[Tuple[str, str]]) -> List[str]:
    base = ["uid", "program", "network", "taken", "Server UID"]
    mapping_columns = [display for display, _ in mappings]
    ordered: List[str] = []
    for name in base + mapping_columns + FILTERED_FIXED_EXTRA_HEADERS:
        if name not in ordered:
            ordered.append(name)
    return ordered


def enrich_raw_rows_with_scope_trace(raw_rows: List[Dict[str, Any]], scope_rows: List[Dict[str, Any]]) -> None:
    """Populate RAW with inventory/workload/exclusion trace columns from final scope when keys match."""
    by_uid_server: Dict[Tuple[str, str], Dict[str, Any]] = {}
    by_uid_hostname: Dict[Tuple[str, str], Dict[str, Any]] = {}
    rows_by_uid: Dict[str, List[Dict[str, Any]]] = {}

    for row in scope_rows:
        uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"]))
        server_uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
        hostname = _normalize_lookup_value(
            _short_hostname(
                _get_row_value_by_candidates(
                    row,
                    ["HOSTNAME", "hostname", "USUAL NAME", "usual_name", "INV_hostname", "INV_ocs_name"],
                )
            )
        )
        if uid:
            rows_by_uid.setdefault(uid, []).append(row)
        if uid and server_uid:
            by_uid_server[(uid, server_uid)] = row
        if uid and hostname:
            by_uid_hostname[(uid, hostname)] = row

    unique_row_by_uid: Dict[str, Dict[str, Any]] = {
        uid: uid_rows[0]
        for uid, uid_rows in rows_by_uid.items()
        if len(uid_rows) == 1
    }

    for row in raw_rows:
        uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"]))
        server_uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
        raw_hostname = _normalize_lookup_value(
            _short_hostname(
                _get_row_value_by_candidates(
                    row,
                    ["HOSTNAME", "hostname", "USUAL NAME", "usual_name", "INV_hostname", "INV_ocs_name"],
                )
            )
        )

        source = None
        if uid and server_uid:
            source = by_uid_server.get((uid, server_uid))
        if source is None and uid and raw_hostname:
            source = by_uid_hostname.get((uid, raw_hostname))
        if source is None and uid and not server_uid and not raw_hostname:
            source = unique_row_by_uid.get(uid)

        for header in RAW_SCOPE_TRACE_HEADERS:
            row[header] = source.get(header, row.get(header, "")) if source else row.get(header, "")
        if str(row.get("F_Excluded", "")).strip() == "":
            row["F_Excluded"] = "N"
        if _normalize_lookup_value(row.get("F_Excluded", "N")) == "Y":
            row["F_FILTER_ALL"] = "N"


def annotate_raw_scope_programs(raw_rows: List[Dict[str, Any]], monitored_rows: List[Dict[str, str]]) -> None:
    """Mark RAW rows with In Scope(s)/Program(s) based on (uid, IPLIST) vs (uid, network)."""
    monitored_by_uid: Dict[str, List[Dict[str, str]]] = {}
    for monitored_row in monitored_rows:
        uid = _normalize_lookup_value(monitored_row.get("uid", ""))
        if not uid:
            continue
        monitored_by_uid.setdefault(uid, []).append(monitored_row)

    for row in raw_rows:
        row["In Scope(s)"] = "N"
        row["Program(s)"] = ""

        if _normalize_lookup_value(row.get("F_FILTER_ALL", "")) != "Y":
            continue

        uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"]))
        iplist = _normalize_lookup_value(_get_row_value_by_candidates(row, ["ILU_IPLIST", "IPLIST"]))
        if not uid:
            continue

        matched_programs: set[str] = set()
        for monitored_row in monitored_by_uid.get(uid, []):
            network = _normalize_lookup_value(monitored_row.get("network", ""))
            is_match = False
            if network and iplist and network in iplist:
                is_match = True
            elif not network and not iplist:
                is_match = True
            if is_match:
                program = str(monitored_row.get("program", "")).strip()
                if program:
                    matched_programs.add(program)

        if matched_programs:
            row["In Scope(s)"] = "Y"
            row["Program(s)"] = ",".join(sorted(matched_programs))


def build_filtered_rows_from_raw(raw_rows: List[Dict[str, Any]], monitored_rows: List[Dict[str, str]]) -> List[Dict[str, Any]]:
    """Build FILTRED rows directly from RAW using F_FILTER_ALL/In Scope(s) conditions."""
    monitored_by_uid: Dict[str, List[Dict[str, str]]] = {}
    for monitored_row in monitored_rows:
        uid = _normalize_lookup_value(monitored_row.get("uid", ""))
        if not uid:
            continue
        monitored_by_uid.setdefault(uid, []).append(monitored_row)

    out_rows: List[Dict[str, Any]] = []
    for raw_row in raw_rows:
        if _normalize_lookup_value(raw_row.get("F_FILTER_ALL", "")) != "Y":
            continue
        if _normalize_lookup_value(raw_row.get("In Scope(s)", "")) != "Y":
            continue

        uid = _normalize_lookup_value(_get_row_value_by_candidates(raw_row, ["uid"]))
        iplist = _normalize_lookup_value(_get_row_value_by_candidates(raw_row, ["ILU_IPLIST", "IPLIST"]))
        if not uid:
            continue

        matches: List[Dict[str, str]] = []
        for monitored_row in monitored_by_uid.get(uid, []):
            network = _normalize_lookup_value(monitored_row.get("network", ""))
            if (network and iplist and network in iplist) or (not network and not iplist):
                matches.append(monitored_row)

        if not matches:
            row = dict(raw_row)
            row["program"] = str(row.get("program", "")).strip()
            row["network"] = str(row.get("network", "")).strip()
            row["taken"] = str(row.get("taken", "")).strip()
            out_rows.append(row)
            continue

        for monitored_row in matches:
            row = dict(raw_row)
            row["program"] = str(monitored_row.get("program", "")).strip()
            row["network"] = str(monitored_row.get("network", "")).strip()
            row["taken"] = str(monitored_row.get("taken", "")).strip()
            out_rows.append(row)

    deduped = _deduplicate_initial_filtered_rows(out_rows)
    log.info("FILTRED rebuilt from RAW rows=%s deduped=%s", len(out_rows), len(deduped))
    return deduped


def populate_enrich_scope_columns_from_monitored(enrich_rows: List[Dict[str, Any]], monitored_rows: List[Dict[str, str]]) -> None:
    """Backfill program/network/taken for ENRICH rows using monitored_kears matching rules."""
    monitored_by_uid: Dict[str, List[Dict[str, str]]] = {}
    for monitored_row in monitored_rows:
        uid = _normalize_lookup_value(monitored_row.get("uid", ""))
        if not uid:
            continue
        monitored_by_uid.setdefault(uid, []).append(monitored_row)

    for row in enrich_rows:
        uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"]))
        if not uid:
            continue

        iplist = _normalize_lookup_value(_get_row_value_by_candidates(row, ["ILU_IPLIST", "IPLIST"]))
        programs_raw = str(_get_row_value_by_candidates(row, ["Program(s)", "program"]) or "")
        program_tokens = {
            _normalize_lookup_value(token)
            for token in programs_raw.split(",")
            if _normalize_lookup_value(token)
        }

        for monitored_row in monitored_by_uid.get(uid, []):
            monitored_program = _normalize_lookup_value(monitored_row.get("program", ""))
            if program_tokens and monitored_program and monitored_program not in program_tokens:
                continue

            network = _normalize_lookup_value(monitored_row.get("network", ""))
            network_match = (network and iplist and network in iplist) or (not network and not iplist)
            if not network_match:
                continue

            row["program"] = str(monitored_row.get("program", "")).strip()
            row["network"] = str(monitored_row.get("network", "")).strip()
            row["taken"] = str(monitored_row.get("taken", "")).strip()
            break


def _raw_filter_fieldnames() -> List[str]:
    return [name for pair in RAW_FILTER_COLUMN_PAIRS for name in pair[:2] if name] + ["F_FILTER_ALL"]


def _insert_column_after(fieldnames: List[str], anchor: str, column_name: str) -> None:
    if column_name not in fieldnames:
        return
    if anchor not in fieldnames:
        return
    fieldnames.remove(column_name)
    anchor_index = fieldnames.index(anchor)
    fieldnames.insert(anchor_index + 1, column_name)


def _insert_column_after(fieldnames: List[str], anchor: str, column_name: str) -> None:
    if column_name not in fieldnames:
        return
    if anchor not in fieldnames:
        return
    fieldnames.remove(column_name)
    anchor_index = fieldnames.index(anchor)
    fieldnames.insert(anchor_index + 1, column_name)


def _is_truthy_flag(value: Any) -> bool:
    return _normalize_lookup_value(value) in {"TRUE", "1", "YES", "Y"}


def _is_not_in_scope_flag(value: Any) -> bool:
    return _normalize_lookup_value(value) in {"N", "NO", "FALSE", "0"}


def _sanitize_sheet_name(name: str) -> str:
    cleaned = "".join("_" if ch in {"\\", "/", "*", "?", ":", "[", "]"} else ch for ch in str(name or "").strip())
    return cleaned[:31] or "Program"


def _format_ratio_label(numerator: int, denominator: int) -> str:
    percent = (float(numerator) / float(denominator) * 100.0) if denominator > 0 else 0.0
    return f"({numerator}/{denominator}) {percent:.2f}%".replace(".", ",")


TOTAL_SHEET_COLUMN_WIDTHS = {
    "Program": 20.0,
    "Entity": 10.0,
    "Number of Applications": 10.0,
    "Total Assets in Dali (in scope)": 15.0,
    "Variation Total servers": 20.0,
    "% servers with illumio installed": 20.0,
    "Variation % servers with illumio installed": 20.0,
    "% servers with illumio agent in blocking mode": 20.0,
    "Variation % servers with illumio agent in blocking mode": 20.0,
    "% servers with illumio installed Trend Icon": 6.0,
    "% servers with illumio agent in blocking mode Trend Icon": 6.0,
}

EXCEL_TO_OOXML_WIDTH_RATIO = 0.939

STATS_SHEET_COLUMN_WIDTHS: Dict[str, Optional[float]] = {
    "Index": 5.0,
    "Program": 20.0,
    "Entity": 10.0,
    "Kear ID": 35.0,
    "Application Short Label": None,
    "Total Assets in Dali (in scope)": 10.0,
    "Total Assets in Dali (Enriched)": 10.0,
    "Assets in Dali not in illumio": 10.0,
    "Assets in Dali (Enriched) not in illumio": 10.0,
    "% servers with illumio installed": 15.0,
    "% servers with illumio installed Indicator Icon": 2.5,
    "% servers with illumio installed Trend Icon": 2.5,
    "% servers with illumio installed (Enriched)": 15.0,
    "% servers with illumio installed (Enriched) Indicator Icon": 2.5,
    "% servers with illumio installed (Enriched) Trend Icon": 2.5,
    "% servers with illumio agent in blocking mode": 15.0,
    "% servers with illumio agent in blocking mode Indicator Icon": 2.5,
    "% servers with illumio agent in blocking mode Trend Icon": 2.5,
    "% servers with illumio agent in blocking mode (Enriched)": 15.0,
    "% servers with illumio agent in blocking mode (Enriched) Indicator Icon": 2.5,
    "% servers with illumio agent in blocking mode (Enriched) Trend Icon": 2.5,
}

STATS_ICON_HEADER_COLUMNS = {
    "% servers with illumio installed Indicator Icon",
    "% servers with illumio installed Trend Icon",
    "% servers with illumio installed (Enriched) Indicator Icon",
    "% servers with illumio installed (Enriched) Trend Icon",
    "% servers with illumio agent in blocking mode Indicator Icon",
    "% servers with illumio agent in blocking mode Trend Icon",
    "% servers with illumio agent in blocking mode (Enriched) Indicator Icon",
    "% servers with illumio agent in blocking mode (Enriched) Trend Icon",
}


def _fixed_total_sheet_widths(fieldnames: List[str]) -> List[float]:
    widths: List[float] = []
    for fieldname in fieldnames:
        raw_name = str(fieldname or "").strip()
        base_name = raw_name if raw_name.endswith(" Trend Icon") else re.sub(r"\s*\([^)]*\)\s*$", "", raw_name)
        widths.append(TOTAL_SHEET_COLUMN_WIDTHS.get(base_name, 20.0))
    return widths


def _fixed_stats_sheet_widths(fieldnames: List[str]) -> List[Optional[float]]:
    widths: List[Optional[float]] = []
    for fieldname in fieldnames:
        requested = STATS_SHEET_COLUMN_WIDTHS.get(str(fieldname or "").strip())
        if requested is None:
            widths.append(None)
        else:
            widths.append(round(float(requested) / EXCEL_TO_OOXML_WIDTH_RATIO, 2))
    return widths


def _base_field_name(fieldname: Any) -> str:
    return re.sub(r"\s*\([^)]*\)\s*$", "", str(fieldname or "").strip())


def _is_variation_column(fieldname: Any) -> bool:
    return _base_field_name(fieldname).startswith("Variation ")


RECAP_RIGHT_ALIGNED_COLUMNS = {
    "% servers with illumio installed",
    "% servers with illumio installed (Enriched)",
    "% servers with illumio agent in blocking mode",
    "% servers with illumio agent in blocking mode (Enriched)",
    "Variation Total servers",
    "Variation % servers with illumio installed",
    "Variation % servers with illumio agent in blocking mode",
}

STATS_ENRICHED_COLUMNS = {
    "Total Assets in Dali (Enriched)",
    "Assets in Dali (Enriched) not in illumio",
    "% servers with illumio installed (Enriched)",
    "% servers with illumio agent in blocking mode (Enriched)",
}


def build_illumio_gap_sheets(
    scope_rows: List[Dict[str, Any]],
    excluded_rows: Optional[List[Dict[str, str]]] = None,
) -> List[Tuple[str, List[Dict[str, Any]], Optional[List[str]]]]:
    not_in_illumio_headers = [
        "program",
        "DSI REL",
        "UID REL",
        "SHORT LABEL REL",
        "IPLIST",
        "SUBNET",
        "HOSTNAME",
        "Server Status",
        "ENVIRONMENT",
        "Server UID",
        "DALI STATUS",
        "STATUS",
        "CLOUD TYPE",
        "Retrived from",
        "INV_Owner_Account",
        "INV_Beneficiary",
    ]
    in_illumio_not_blocking_headers = [
        "program",
        "DSI REL",
        "UID REL",
        "SHORT LABEL REL",
        "IPLIST",
        "SUBNET",
        "HOSTNAME",
        "Server Status",
        "ENVIRONMENT",
        "Server UID",
        "DALI STATUS",
        "STATUS",
        "CLOUD TYPE",
        "Retrived from",
        "INV_Owner_Account",
        "INV_Beneficiary",
        "enforcement",
        "role",
        "app",
        "env",
        "loc",
    ]

    not_in_illumio_rows: List[Dict[str, Any]] = []
    in_illumio_not_blocking_rows: List[Dict[str, Any]] = []

    excluded_hostnames = {
        _normalize_hostname_for_compare(item.get("Server to exclude", ""))
        for item in (excluded_rows or [])
        if _normalize_hostname_for_compare(item.get("Server to exclude", ""))
    }

    for row in scope_rows:
        lookup_candidates = [
            _get_row_value_by_candidates(row, ["HOSTNAME", "hostname", "INV_hostname"]),
            _get_row_value_by_candidates(row, ["USUAL NAME", "usual_name"]),
            _get_row_value_by_candidates(row, ["INV_ocs_name", "INV_Ocs_Name"]),
        ]
        if any(_normalize_hostname_for_compare(value) in excluded_hostnames for value in lookup_candidates if value):
            continue

        base_row = {
            "program": _get_row_value_by_candidates(row, ["Program(s)", "program"]),
            "HOSTNAME": _get_row_value_by_candidates(row, ["DALI [CI] HOSTNAME", "HOSTNAME", "hostname", "INV_hostname"]),
            "Server Status": _get_row_value_by_candidates(row, ["DALI [CI] Server Status", "Server Status", "server_status", "server.status"]),
            "Server UID": _get_row_value_by_candidates(row, ["Server UID", "server_uid"]),
            "UID REL": _get_row_value_by_candidates(row, ["DALI [APP] UID", "UID REL", "uid"]),
            "SHORT LABEL REL": _get_row_value_by_candidates(row, ["DALI [APP] SHORT LABEL", "SHORT LABEL REL", "short_label"]),
            "DSI REL": _get_row_value_by_candidates(row, ["DALI [APP] DSI", "DSI REL", "dsi"]),
            "ENVIRONMENT": _get_row_value_by_candidates(row, ["DALI [CI] ENVIRONMENT", "ENVIRONMENT", "environment"]),
            "DALI STATUS": _get_row_value_by_candidates(row, ["DALI [CI] USAGE", "DALI STATUS", "usage"]),
            "STATUS": _get_row_value_by_candidates(row, ["DALI [CI] STATUS", "STATUS", "status"]),
            "CLOUD TYPE": _get_row_value_by_candidates(row, ["DALI [CI] CLOUD TYPE", "CLOUD TYPE", "cloud_type", "server_cloud_type"]),
            "Retrived from": _get_row_value_by_candidates(row, ["Retrived from"]),
            "INV_Owner_Account": _get_row_value_by_candidates(row, ["INV_Owner_Account"]),
            "INV_Beneficiary": _get_row_value_by_candidates(row, ["INV_Beneficiary", "INV_Beneficiary_Account"]),
            "IPLIST": _get_row_value_by_candidates(row, ["ILU_IPLIST", "IPLIST"]),
            "SUBNET": _get_row_value_by_candidates(row, ["ILU_SUBNET", "SUBNET"]),
        }

        managed_value = _normalize_lookup_value(_get_row_value_by_candidates(row, ["ILU_managed", "managed"]))
        if managed_value != "TRUE":
            not_in_illumio_rows.append(base_row)
            continue

        enforcement = _get_row_value_by_candidates(row, ["ILU_enforcement", "enforcement"])
        if _normalize_lookup_value(enforcement) not in {"SELECTIVE", "FULL"}:
            in_illumio_not_blocking_rows.append(
                {
                    **base_row,
                    "enforcement": enforcement,
                    "role": _get_row_value_by_candidates(row, ["ILU_role", "role"]),
                    "app": _get_row_value_by_candidates(row, ["ILU_app", "app"]),
                    "env": _get_row_value_by_candidates(row, ["ILU_env", "env"]),
                    "loc": _get_row_value_by_candidates(row, ["ILU_loc", "loc"]),
                }
            )

    return [
        ("NOT_IN_ILLUMIO", not_in_illumio_rows, not_in_illumio_headers),
        ("IN_ILLUMIO_BUT_NOT_BLOCKING", in_illumio_not_blocking_rows, in_illumio_not_blocking_headers),
    ]


def build_out_of_scope_sheet(
    raw_rows: List[Dict[str, Any]],
    enrich_rows: List[Dict[str, Any]],
    dict_kear_account_rows: Optional[List[Dict[str, Any]]] = None,
) -> Tuple[str, List[Dict[str, Any]], Optional[List[str]]]:
    out_of_scope_headers = [
        "UID REL",
        "SHORT LABEL REL",
        "DSI REL",
        "HOSTNAME",
        "Server Status",
        "IPLIST",
        "SUBNET",
        "ENVIRONMENT",
        "Server UID",
        "DALI STATUS",
        "STATUS",
        "CLOUD TYPE",
        "Retrived from",
        "INV_Owner_Account",
        "INV_Beneficiary",
        "INV_Beneficiary_Account_ENV",
        "enforcement",
        "role",
        "app",
        "env",
        "loc",
    ]

    beneficiary_env_by_account: Dict[str, str] = {}
    for dict_row in dict_kear_account_rows or []:
        beneficiary_key = _normalize_lookup_value(dict_row.get("INV_Beneficiary_Account", ""))
        if not beneficiary_key:
            continue
        env_value = _normalize_cell_value(dict_row.get("INV_Beneficiary_Account_ENV", "")).strip()
        if env_value and beneficiary_key not in beneficiary_env_by_account:
            beneficiary_env_by_account[beneficiary_key] = env_value

    def _row_is_out_of_scope_exact(row: Dict[str, Any]) -> bool:
        # Intentionally mirror GLOBAL's counting logic exactly:
        # F_FILTER_ALL = Y and In Scope(s) = N
        return (
            _normalize_lookup_value(_get_row_value_by_candidates(row, ["F_FILTER_ALL"])) == "Y"
            and _normalize_lookup_value(_get_row_value_by_candidates(row, ["In Scope(s)"])) == "N"
        )

    out_of_scope_rows: List[Dict[str, Any]] = []
    for row in [*(raw_rows or []), *(enrich_rows or [])]:
        if not _row_is_out_of_scope_exact(row):
            continue

        inv_beneficiary = _get_first_non_empty_by_candidates(row, ["INV_Beneficiary", "INV_Beneficiary_Account"])
        inv_beneficiary_env = beneficiary_env_by_account.get(
            _normalize_lookup_value(inv_beneficiary),
            _get_first_non_empty_by_candidates(row, ["INV_Beneficiary_Account_ENV"]),
        )

        out_of_scope_rows.append(
            {
                "UID REL": _get_first_non_empty_by_candidates(row, ["DALI [APP] UID", "UID REL", "uid"]),
                "SHORT LABEL REL": _get_first_non_empty_by_candidates(row, ["DALI [APP] SHORT LABEL", "SHORT LABEL REL", "short_label"]),
                "DSI REL": _get_first_non_empty_by_candidates(row, ["DALI [APP] DSI", "DSI REL", "dsi"]),
                "HOSTNAME": _get_first_non_empty_by_candidates(row, ["DALI [CI] HOSTNAME", "HOSTNAME", "hostname", "INV_hostname"]),
                "Server Status": _get_first_non_empty_by_candidates(row, ["DALI [CI] Server Status", "Server Status", "server_status", "server.status"]),
                "IPLIST": _get_first_non_empty_by_candidates(row, ["ILU_IPLIST", "IPLIST"]),
                "SUBNET": _get_first_non_empty_by_candidates(row, ["ILU_SUBNET", "SUBNET"]),
                "ENVIRONMENT": _get_first_non_empty_by_candidates(row, ["DALI [CI] ENVIRONMENT", "ENVIRONMENT", "environment"]),
                "Server UID": _get_first_non_empty_by_candidates(row, ["Server UID", "server_uid"]),
                "DALI STATUS": _get_first_non_empty_by_candidates(row, ["DALI [CI] USAGE", "DALI STATUS", "usage"]),
                "STATUS": _get_first_non_empty_by_candidates(row, ["DALI [CI] STATUS", "STATUS", "status"]),
                "CLOUD TYPE": _get_first_non_empty_by_candidates(row, ["DALI [CI] CLOUD TYPE", "CLOUD TYPE", "cloud_type", "server_cloud_type"]),
                "Retrived from": _get_first_non_empty_by_candidates(row, ["Retrived from"]),
                "INV_Owner_Account": _get_first_non_empty_by_candidates(row, ["INV_Owner_Account"]),
                "INV_Beneficiary": inv_beneficiary,
                "INV_Beneficiary_Account_ENV": inv_beneficiary_env,
                "enforcement": _get_first_non_empty_by_candidates(row, ["ILU_enforcement", "enforcement"]),
                "role": _get_first_non_empty_by_candidates(row, ["ILU_role", "role"]),
                "app": _get_first_non_empty_by_candidates(row, ["ILU_app", "app"]),
                "env": _get_first_non_empty_by_candidates(row, ["ILU_env", "env"]),
                "loc": _get_first_non_empty_by_candidates(row, ["ILU_loc", "loc"]),
            }
        )

    return ("OUT_OF_SCOPE", out_of_scope_rows, out_of_scope_headers)


def build_program_recap_sheets(
    monitored_rows: List[Dict[str, str]],
    filtered_rows: List[Dict[str, Any]],
    scope_rows: List[Dict[str, Any]],
    raw_rows: Optional[List[Dict[str, Any]]],
    enrich_rows: Optional[List[Dict[str, Any]]],
    output_path: Path,
) -> List[Tuple[str, List[Dict[str, Any]], Optional[List[str]]]]:
    headers = [
        "Index",
        "Program",
        "Entity",
        "Sub-Entity",
        "Kear ID",
        "Application Short Label",
        "Total Assets in Dali (in scope)",
        "Total Assets (enriched)",
        "Total Assets (not in Scope)",
        "Assets in Dali not in illumio",
        "% servers with illumio installed",
        "% servers with illumio agent in blocking mode",
        "Total Assets in Dali (Enriched)",
        "Assets in Dali (Enriched) not in illumio",
        "% servers with illumio installed (Enriched)",
        "% servers with illumio agent in blocking mode (Enriched)",
    ]

    def _row_uid(row: Dict[str, Any]) -> str:
        # STATS matching must be done from the first-column business key (`uid`)
        # used in RAW/ENRICH/STATS "Kear ID", not from derived UID fields.
        return _normalize_lookup_value(_get_first_non_empty_by_candidates(row, ["uid"]))

    def _row_program(row: Dict[str, Any]) -> str:
        return str(_get_first_non_empty_by_candidates(row, ["program", "Program(s)"])).strip() or "Unknown"

    index_by_key_filtered: Dict[Tuple[str, str], List[Dict[str, Any]]] = {}
    for row in filtered_rows:
        uid = _row_uid(row)
        if uid:
            index_by_key_filtered.setdefault((uid, _row_program(row)), []).append(row)

    index_by_key_scope: Dict[Tuple[str, str], List[Dict[str, Any]]] = {}
    for row in scope_rows:
        uid = _row_uid(row)
        if uid:
            index_by_key_scope.setdefault((uid, _row_program(row)), []).append(row)

    recap_rows: List[Dict[str, Any]] = []
    seen_keys: set[Tuple[str, str]] = set()

    for monitored_row in monitored_rows:
        program = str(monitored_row.get("program", "")).strip() or "Unknown"
        uid = _normalize_lookup_value(monitored_row.get("uid", ""))
        if not uid:
            continue

        key = (program, uid)
        if key in seen_keys:
            continue
        seen_keys.add(key)

        base_rows = index_by_key_filtered.get((uid, program), [])
        enriched_rows = index_by_key_scope.get((uid, program), [])

        base_total = len(base_rows)
        enriched_total = len(enriched_rows)

        log.debug(
            "STATS trace uid=%s program=%s base_rows=%s enriched_rows=%s",
            uid,
            program,
            base_total,
            enriched_total,
        )

        managed_true_base = [row for row in base_rows if _parse_managed_flag(_get_row_value_by_candidates(row, ["ILU_managed", "managed"]))]
        managed_true_enriched = [row for row in enriched_rows if _parse_managed_flag(_get_row_value_by_candidates(row, ["ILU_managed", "managed"]))]

        # "not in illumio" includes FALSE and blank managed values (anything not TRUE).
        not_in_illumio_base = base_total - len(managed_true_base)
        not_in_illumio_enriched = enriched_total - len(managed_true_enriched)

        blocking_base = sum(
            1
            for row in managed_true_base
            if _normalize_lookup_value(_get_row_value_by_candidates(row, ["ILU_enforcement", "enforcement"])) in {"SELECTIVE", "FULL"}
        )
        blocking_enriched = sum(
            1
            for row in managed_true_enriched
            if _normalize_lookup_value(_get_row_value_by_candidates(row, ["ILU_enforcement", "enforcement"])) in {"SELECTIVE", "FULL"}
        )

        display_rows = enriched_rows or base_rows
        metadata_rows = display_rows
        entity = next(
            (
                _normalize_cell_value(
                    _get_row_value_by_candidates(
                        row,
                        ["DALI [APP] DSI", "DSI REL", "dsi", "application_management_rc"],
                    )
                ).strip()
                for row in metadata_rows
                if _normalize_cell_value(
                    _get_row_value_by_candidates(
                        row,
                        ["DALI [APP] DSI", "DSI REL", "dsi", "application_management_rc"],
                    )
                ).strip()
            ),
            "",
        )
        sub_entity_source = _normalize_cell_value(
            _get_row_value_by_candidates(
                next(iter(metadata_rows), {}),
                ["DALI [APP] APPLICATION MANAGEMENT RC", "application_management_rc", "APPLICATION MANAGEMENT RC REL"],
            )
        ).strip()
        sub_entity = sub_entity_source.split("-", 1)[0].strip() if sub_entity_source else ""

        short_label = next(
            (
                _normalize_cell_value(
                    _get_row_value_by_candidates(
                        row,
                        ["DALI [APP] SHORT LABEL", "SHORT LABEL REL", "short_label"],
                    )
                ).strip()
                for row in metadata_rows
                if _normalize_cell_value(
                    _get_row_value_by_candidates(
                        row,
                        ["DALI [APP] SHORT LABEL", "SHORT LABEL REL", "short_label"],
                    )
                ).strip()
            ),
            "",
        )

        recap_rows.append(
            {
                "Program": program,
                "Entity": entity,
                "Sub-Entity": sub_entity,
                "Kear ID": uid,
                "Application Short Label": short_label,
                "Total Assets in Dali (in scope)": str(base_total),
                "Total Assets in Dali (Enriched)": str(enriched_total),
                "Assets in Dali not in illumio": str(not_in_illumio_base),
                "Assets in Dali (Enriched) not in illumio": str(not_in_illumio_enriched),
                "% servers with illumio installed": _format_ratio_label(len(managed_true_base), base_total),
                "% servers with illumio installed (Enriched)": _format_ratio_label(len(managed_true_enriched), enriched_total),
                "% servers with illumio agent in blocking mode": _format_ratio_label(blocking_base, base_total),
                "% servers with illumio agent in blocking mode (Enriched)": _format_ratio_label(
                    blocking_enriched,
                    enriched_total,
                ),
            }
        )

    recap_rows.sort(key=lambda row: _ratio_percent_from_label(row.get("% servers with illumio installed", "")))
    recap_rows.sort(key=lambda row: _normalize_lookup_value(row.get("Entity", "")))
    recap_rows.sort(key=lambda row: _normalize_lookup_value(row.get("Program", "")), reverse=True)

    for index_value, recap_row in enumerate(recap_rows, start=1):
        recap_row["Index"] = str(index_value)

    previous_stats = _load_previous_stats_workbook(output_path)
    recap_rows, headers = _append_stats_visual_columns(recap_rows, headers, previous_stats)
    for recap_row in recap_rows:
        recap_row.pop("Total Assets (enriched)", None)
        recap_row.pop("Total Assets (not in Scope)", None)
    headers = [header for header in headers if header not in {"Total Assets (enriched)", "Total Assets (not in Scope)"}]

    default_last_month_label = _last_month_label_from_output(output_path)
    previous_totals, previous_totals_file = _load_previous_totals_workbook(output_path)
    totals_baseline_label = _label_from_workbook_path(previous_totals_file, default_last_month_label)
    total_program_sheet = _build_total_program_rows(recap_rows, totals_baseline_label, previous_totals.get("TOTAL.PROGRAM", {}))
    total_entity_sheet = _build_total_entity_rows(recap_rows, totals_baseline_label, previous_totals.get("TOTAL.ENTITY", {}))

    return [
        ("STATS", recap_rows, headers),
        total_program_sheet,
        total_entity_sheet,
    ]


def build_global_sheet(
    raw_rows: List[Dict[str, Any]],
    enrich_rows: List[Dict[str, Any]],
) -> Tuple[str, List[Dict[str, Any]], Optional[List[str]]]:
    """Build GLOBAL app-oriented stats keyed only by uid.

    Counts are intentionally aligned with Excel COUNTIFS logic used for validation:
    - Total Assets PRD (enriched) = COUNTIFS(uid, F_FILTER_ALL=Y) on RAW + ENRICH
    - Total Assets PRD (not in Scope) = COUNTIFS(uid, F_FILTER_ALL=Y, In Scope(s)=N) on RAW + ENRICH
    """
    headers = [
        "DALI [APP] UID",
        "DALI [APP] SHORT LABEL",
        "Total Assets PRD (enriched)",
        "Total Assets PRD (not in Scope)",
    ]

    uid_to_short_label: Dict[str, str] = {}
    for row in list(raw_rows or []) + list(enrich_rows or []):
        uid = _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"]))
        if not uid:
            continue
        short_label = _normalize_cell_value(
            _get_row_value_by_candidates(row, ["DALI [APP] SHORT LABEL", "SHORT LABEL REL", "short_label"])
        ).strip()
        if short_label and uid not in uid_to_short_label:
            uid_to_short_label[uid] = short_label

    all_uids = sorted(uid_to_short_label.keys())
    global_rows: List[Dict[str, Any]] = []

    for uid in all_uids:
        raw_enriched = sum(
            1
            for row in (raw_rows or [])
            if _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"])) == uid
            and _normalize_lookup_value(_get_row_value_by_candidates(row, ["F_FILTER_ALL"])) == "Y"
        )
        enrich_enriched = sum(
            1
            for row in (enrich_rows or [])
            if _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"])) == uid
            and _normalize_lookup_value(_get_row_value_by_candidates(row, ["F_FILTER_ALL"])) == "Y"
        )

        raw_not_in_scope = sum(
            1
            for row in (raw_rows or [])
            if _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"])) == uid
            and _normalize_lookup_value(_get_row_value_by_candidates(row, ["F_FILTER_ALL"])) == "Y"
            and _normalize_lookup_value(_get_row_value_by_candidates(row, ["In Scope(s)"])) == "N"
        )
        enrich_not_in_scope = sum(
            1
            for row in (enrich_rows or [])
            if _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"])) == uid
            and _normalize_lookup_value(_get_row_value_by_candidates(row, ["F_FILTER_ALL"])) == "Y"
            and _normalize_lookup_value(_get_row_value_by_candidates(row, ["In Scope(s)"])) == "N"
        )

        global_rows.append(
            {
                "DALI [APP] UID": uid,
                "DALI [APP] SHORT LABEL": uid_to_short_label.get(uid, ""),
                "Total Assets PRD (enriched)": str(raw_enriched + enrich_enriched),
                "Total Assets PRD (not in Scope)": str(raw_not_in_scope + enrich_not_in_scope),
            }
        )

    return ("GLOBAL", global_rows, headers)


DICT_SHEET_HEADERS = [
    "Sheet",
    "Sheet Purpose",
    "Column",
    "Kind",
    "Source",
    "Calculation / Logic",
    "Quick Explanation",
]


def _dict_sheet_purpose(sheet_name: str) -> str:
    purposes = {
        "Summary": "High-level execution summary of the run and key counters.",
        "RAW": "Detailed DALI impact dataset before the FILTRED flattening step.",
        "FILTRED": "Filtered and flattened in-scope dataset used as the base for program-aware STATS calculations.",
        "get_inv_by_account": "Raw inventory extract used as an external reference for host and account enrichment.",
        "get_marley_gen2_by_uuid": "Raw Marley / DALI enrichment extract collected during the run.",
        "DictKearAccount": "Dictionary linking beneficiary accounts to DALI application identifiers and reference labels.",
        "MONITORED_SCOPES": "Direct copy of user_inputs/monitored_kears.csv used as the business scope input for the run.",
        "CANDIDATE": "Workload candidates with non-empty IPLIST enriched from Marley and checked against monitored scopes.",
        "ENRICH": "Enriched server-level dataset combining DALI, Marley, inventory and Illumio attributes.",
        "SCOPE": "Subset of ENRICH limited to monitored and in-scope rows after the scope logic is applied.",
        "STATS": "Application and program KPI sheet aggregated by the business key {Program, Kear ID}.",
        "TOTAL.PROGRAM": "Program-level KPI aggregation built from the STATS sheet.",
        "TOTAL.ENTITY": "Program and entity KPI aggregation built from the STATS sheet.",
        "NOT_IN_ILLUMIO": "SCOPE rows where the asset is not found as protected in Illumio.",
        "IN_ILLUMIO_BUT_NOT_BLOCKING": "SCOPE rows found in Illumio but not enforcing in blocking mode.",
        "EXCLUDED": "Rows excluded from the operational scope by explicit exclusion logic.",
        "GLOBAL": "Application-level PRD counts keyed only by application UID across RAW and ENRICH.",
        "OUT_OF_SCOPE": "Detailed list of PRD assets that pass filtering but are outside the monitored scope.",
        "Dict": "Workbook dictionary that explains the purpose of each sheet and the meaning of each visible column.",
    }
    return purposes.get(sheet_name, "Workbook output sheet generated by the pipeline.")


def _dict_column_metadata(sheet_name: str, column_name: str) -> Tuple[str, str, str, str]:
    column = str(column_name or "").strip()
    base_name = _base_field_name(column)

    if sheet_name == "Summary":
        if column == "Metric":
            return ("attribute", "Pipeline summary builder", "Direct label", "Name of the summary metric displayed in the Summary sheet.")
        if column == "Value":
            return ("attribute", "Pipeline summary builder", "Direct value", "Value associated with the summary metric for the current run.")

    if sheet_name == "MONITORED_SCOPES":
        return (
            "attribute",
            "user_inputs/monitored_kears.csv",
            "Direct copy from the input CSV",
            "Business scope input provided by the user and reused by the pipeline as the monitored reference.",
        )

    if sheet_name == "CANDIDATE":
        if column == "scope_alread_monitored?":
            return (
                "computed",
                "MONITORED_SCOPES + Marley + export_wkld.derived.csv",
                "Y when one Marley app_info.kear_uuid token matches a monitored uid and the monitored network is contained in IPLIST, else N",
                "Indicates whether the candidate server already belongs to a monitored scope.",
            )
        if column == "Program":
            return (
                "computed",
                "MONITORED_SCOPES",
                "Program value copied from monitored_kears.csv when scope_alread_monitored? = Y",
                "Program associated with the matched monitored scope.",
            )
        if column == "lookup_status":
            return (
                "computed",
                "data4sec/marley_original",
                "FOUND when a Marley document is returned after short_hostname lookup with hostname and ocs_name_from_IP fallbacks, otherwise NOT_FOUND",
                "Status of the Marley enrichment lookup for the candidate server.",
            )
        if column in {"short_hostname", "hostname", "ocs_name_from_IP", "IPLIST", "SUBNET", "managed", "enforcement", "created_at", "ip_with_default_gw", "app", "env", "loc", "role"}:
            return (
                "attribute",
                "raw/export_wkld.derived.csv",
                "Direct copy from the workload-derived export",
                "Workload-derived candidate attribute extracted from export_wkld.derived.csv.",
            )
        if column in {"status", "uuid", "app_info.kear_uuid", "app_info.factor", "app_info.app_id", "app_info.app_name", "app_info.env", "typologie", "os_name", "dns"}:
            return (
                "attribute",
                "data4sec/marley_original",
                "First document returned by the Marley lookup on short_hostname, or on hostname / ocs_name_from_IP as fallback",
                "Marley enrichment attribute added to the workload candidate row.",
            )
        if column in {"owner_app_name", "beneficiary"}:
            return (
                "attribute",
                "data4sec/inventory",
                "First document returned by the inventory lookup on short_hostname, or on hostname / ocs_name_from_IP as fallback",
                "Inventory enrichment attribute added to the workload candidate row.",
            )
        if column == "beneficiary_account_env":
            return (
                "attribute",
                "data4sec/platform_accounts",
                "ENV tag looked up from platform_accounts using the non-empty beneficiary account returned by inventory",
                "Environment resolved from the beneficiary account when available.",
            )

    if sheet_name == "STATS":
        stats_map = {
            "Index": ("computed", "STATS builder", "Sequential numbering after final sorting", "Display index of the STATS line."),
            "Program": ("attribute", "MONITORED_SCOPES / FILTRED", "Business key copied from the monitored scope and the FILTRED rows", "Program used together with Kear ID as the business key of the line."),
            "Entity": ("attribute", "ENRICH / SCOPE metadata", "First non-empty DSI-style entity label found for the {Program, Kear ID} scope", "Main entity used for aggregation and reporting."),
            "Sub-Entity": ("attribute", "ENRICH / SCOPE metadata", "Derived from the application management RC label when available", "Secondary organizational breakdown used for readability."),
            "Kear ID": ("attribute", "Business application UID", "Business key copied from the monitored application UID", "Application identifier used as the primary key on the application axis."),
            "Application Short Label": ("attribute", "ENRICH / SCOPE metadata", "First non-empty short label found for the current {Program, Kear ID}", "Human-readable short application name."),
            "Total Assets in Dali (in scope)": ("computed", "FILTRED", "Count of FILTRED rows for the current {Program, Kear ID}", "Number of in-scope assets found in DALI for the current application and program."),
            "Assets in Dali not in illumio": ("computed", "FILTRED", "Total in-scope assets minus rows where the managed flag is TRUE", "In-scope assets not currently reported as protected in Illumio."),
            "% servers with illumio installed": ("computed", "FILTRED", "(managed TRUE rows / total FILTRED rows) formatted as (x/y) zz,zz%", "Coverage ratio of Illumio installation for in-scope assets."),
            "% servers with illumio installed Indicator Icon": ("computed", "STATS visual layer", "Helper numeric flag used by the Excel icon set formatting", "Visual indicator of the current installation coverage level."),
            "% servers with illumio installed Trend Icon": ("computed", "Current STATS row vs historical baseline workbook", "Current installation percentage minus the baseline installation percentage for the same {Program, Kear ID}", "Trend delta used by the icon set to show improvement, deterioration or stability over time."),
            "Total Assets in Dali (Enriched)": ("computed", "SCOPE", "Count of SCOPE rows for the current {Program, Kear ID}", "Number of enriched and in-scope assets retained in the SCOPE sheet."),
            "Assets in Dali (Enriched) not in illumio": ("computed", "SCOPE", "Total enriched assets minus rows where the managed flag is TRUE", "Enriched assets not currently reported as protected in Illumio."),
            "% servers with illumio installed (Enriched)": ("computed", "SCOPE", "(managed TRUE SCOPE rows / total SCOPE rows) formatted as (x/y) zz,zz%", "Coverage ratio computed on the enriched perimeter."),
            "% servers with illumio installed (Enriched) Indicator Icon": ("computed", "STATS visual layer", "Helper numeric flag used by the Excel icon set formatting", "Visual indicator of the current enriched installation coverage level."),
            "% servers with illumio installed (Enriched) Trend Icon": ("computed", "Current STATS row vs historical baseline workbook", "Current enriched installation percentage minus the baseline enriched installation percentage for the same {Program, Kear ID}", "Trend delta used by the icon set to show the time progression of enriched coverage."),
            "% servers with illumio agent in blocking mode": ("computed", "FILTRED", "(managed rows in SELECTIVE/FULL enforcement / total FILTRED rows) formatted as (x/y) zz,zz%", "Blocking-mode enforcement ratio on the in-scope perimeter."),
            "% servers with illumio agent in blocking mode Indicator Icon": ("computed", "STATS visual layer", "Helper numeric flag used by the Excel icon set formatting", "Visual indicator of the current blocking-mode enforcement level."),
            "% servers with illumio agent in blocking mode Trend Icon": ("computed", "Current STATS row vs historical baseline workbook", "Current blocking percentage minus the baseline blocking percentage for the same {Program, Kear ID}", "Trend delta used by the icon set to show the time progression of blocking mode."),
            "% servers with illumio agent in blocking mode (Enriched)": ("computed", "SCOPE", "(managed SCOPE rows in SELECTIVE/FULL enforcement / total SCOPE rows) formatted as (x/y) zz,zz%", "Blocking-mode enforcement ratio on the enriched perimeter."),
            "% servers with illumio agent in blocking mode (Enriched) Indicator Icon": ("computed", "STATS visual layer", "Helper numeric flag used by the Excel icon set formatting", "Visual indicator of the current enriched blocking-mode level."),
            "% servers with illumio agent in blocking mode (Enriched) Trend Icon": ("computed", "Current STATS row vs historical baseline workbook", "Current enriched blocking percentage minus the baseline enriched blocking percentage for the same {Program, Kear ID}", "Trend delta used by the icon set to show the time progression of enriched blocking mode."),
        }
        if column in stats_map:
            return stats_map[column]

    if sheet_name in {"TOTAL.PROGRAM", "TOTAL.ENTITY"}:
        total_map = {
            "Program": ("attribute", "STATS aggregation", "Grouping key", "Program aggregation key used by the sheet."),
            "Entity": ("attribute", "STATS aggregation", "Grouping key", "Entity aggregation key used by TOTAL.ENTITY."),
            "Number of Applications": ("computed", "STATS", "Count of STATS rows aggregated into the current total line", "Number of application rows rolled up in the aggregation."),
            "Total Assets in Dali (in scope)": ("computed", "STATS", "Sum of STATS 'Total Assets in Dali (in scope)' across the aggregation perimeter", "Current in-scope asset total for the aggregated perimeter."),
            "Variation Total servers": ("computed", "Current totals vs historical baseline workbook", "Current total assets minus baseline total assets for the same aggregation key", "Absolute variation versus the selected historical reference workbook."),
            "% servers with illumio installed": ("computed", "STATS", "Ratio rebuilt from summed STATS numerators and denominators", "Aggregated installation coverage ratio for the perimeter."),
            "Variation % servers with illumio installed": ("computed", "Current totals vs historical baseline workbook", "Current installation percentage minus baseline installation percentage", "Percentage-point variation versus the selected historical reference workbook."),
            "% servers with illumio agent in blocking mode": ("computed", "STATS", "Ratio rebuilt from summed STATS numerators and denominators", "Aggregated blocking-mode enforcement ratio for the perimeter."),
            "Variation % servers with illumio agent in blocking mode": ("computed", "Current totals vs historical baseline workbook", "Current blocking percentage minus baseline blocking percentage", "Percentage-point variation versus the selected historical reference workbook."),
            "% servers with illumio installed Trend Icon": ("computed", "TOTAL visual layer", "Numeric helper derived from the variation column for the Excel icon set", "Visual trend of the installation ratio at the aggregated level."),
            "% servers with illumio agent in blocking mode Trend Icon": ("computed", "TOTAL visual layer", "Numeric helper derived from the variation column for the Excel icon set", "Visual trend of the blocking-mode ratio at the aggregated level."),
        }
        if base_name in total_map:
            return total_map[base_name]

    if sheet_name == "GLOBAL":
        global_map = {
            "DALI [APP] UID": ("attribute", "RAW + ENRICH", "Direct application UID", "Application UID used as the key of the GLOBAL sheet."),
            "DALI [APP] SHORT LABEL": ("attribute", "RAW + ENRICH", "First available short label for the application UID", "Human-readable application label."),
            "Total Assets PRD (enriched)": ("computed", "RAW + ENRICH", "COUNTIFS on RAW and ENRICH where uid matches and F_FILTER_ALL = Y", "PRD asset count across the filtered RAW and ENRICH perimeter."),
            "Total Assets PRD (not in Scope)": ("computed", "RAW + ENRICH", "COUNTIFS on RAW and ENRICH where uid matches, F_FILTER_ALL = Y and In Scope(s) = N", "PRD asset count outside the monitored scope for the application."),
        }
        if column in global_map:
            return global_map[column]

    if sheet_name == "OUT_OF_SCOPE":
        if column == "INV_Beneficiary_Account_ENV":
            return ("attribute", "DictKearAccount + OUT_OF_SCOPE details", "Matched from DictKearAccount using INV_Beneficiary as the lookup key", "Environment associated with the beneficiary account." )
        if column in {"UID REL", "SHORT LABEL REL", "DSI REL", "HOSTNAME", "Server Status", "IPLIST", "SUBNET", "ENVIRONMENT", "Server UID", "DALI STATUS", "STATUS", "CLOUD TYPE", "Retrived from", "INV_Owner_Account", "INV_Beneficiary", "enforcement", "role", "app", "env", "loc"}:
            return ("attribute", "RAW + ENRICH out-of-scope rows", "First relevant value taken from the underlying RAW or ENRICH detail row", "Detail attribute kept to explain why the asset is counted as out of scope.")

    if sheet_name == "Dict":
        dict_map = {
            "Sheet": ("attribute", "Dict builder", "Direct label", "Workbook sheet name documented by the dictionary."),
            "Sheet Purpose": ("attribute", "Dict builder", "Direct narrative", "Short explanation of what the sheet is used for."),
            "Column": ("attribute", "Dict builder", "Direct label", "Column name documented by the dictionary row."),
            "Kind": ("attribute", "Dict builder", "Direct label", "Whether the documented column is an attribute, a computed KPI, a filter flag or an overview entry."),
            "Source": ("attribute", "Dict builder", "Direct narrative", "Main source dataset or system used to populate the column."),
            "Calculation / Logic": ("attribute", "Dict builder", "Direct narrative", "How the column is calculated or populated."),
            "Quick Explanation": ("attribute", "Dict builder", "Direct narrative", "Short business-friendly explanation of the column."),
        }
        if column in dict_map:
            return dict_map[column]

    if column.startswith("F_FILTER_") or column.startswith("FILTER_VALUE_"):
        return (
            "filter/debug",
            "Pipeline filtering engine",
            "Computed during the filter evaluation stage",
            "Diagnostic filter column used to explain why a row is kept or excluded.",
        )

    if column in {"uid", "UID REL", "DALI [APP] UID", "Kear ID"}:
        return ("attribute", "DALI / business scope input", "Direct identifier", "Application UID or business application identifier.")
    if column in {"program", "Program", "Program(s)"}:
        return ("attribute", "Monitored scope input", "Direct business scope label", "Program associated with the row or aggregation." )
    if column in {"network", "IPLIST", "SUBNET"}:
        return ("attribute", "Monitored scope input / Illumio enrichment", "Direct network scope value", "Network scope or network object associated with the row." )
    if column in {"HOSTNAME", "hostname", "INV_hostname", "INV_ocs_name", "Server UID", "Server Status", "DALI STATUS", "STATUS", "ENVIRONMENT", "CLOUD TYPE", "INV_Owner_Account", "INV_Beneficiary", "INV_Beneficiary_Account", "INV_Beneficiary_Account_ENV", "role", "app", "env", "loc", "enforcement"}:
        return ("attribute", "DALI / Marley / inventory / Illumio", "Direct attribute or first non-empty mapped value", "Operational attribute kept for analysis or drill-down." )
    if column.startswith("Variation "):
        return (
            "computed",
            "Current workbook vs historical baseline workbook",
            "Current value minus baseline value for the same business key",
            "Historical variation used to compare the current run with the selected reference file.",
        )
    if column.endswith("Trend Icon"):
        return (
            "computed",
            "Visual layer",
            "Numeric helper used by the Excel icon set conditional formatting",
            "Trend helper displayed as icons to show improvement, deterioration or stability.",
        )
    if column.endswith("Indicator Icon"):
        return (
            "computed",
            "Visual layer",
            "Numeric helper used by the Excel icon set conditional formatting",
            "Current-state helper displayed as icons to highlight KPI level.",
        )
    if column.startswith("% servers with illumio"):
        return (
            "computed",
            "FILTRED / SCOPE / STATS aggregation",
            "Coverage ratio formatted as (x/y) zz,zz%",
            "Illumio installation or blocking-mode KPI ratio.",
        )
    if column.startswith("Total Assets") or column.startswith("Assets in Dali") or column.startswith("Number of Applications"):
        return (
            "computed",
            "Pipeline aggregation",
            "Count or sum rebuilt from the source perimeter of the sheet",
            "Count-based KPI used for reporting and tracking.",
        )

    return (
        "attribute",
        "Source dataset of the sheet",
        "Direct value or mapped attribute",
        "Visible column kept in the workbook for analysis or traceability.",
    )


def build_dict_sheet(
    raw_headers: List[str],
    filtered_headers: List[str],
    extra_sheets: List[Tuple[str, List[Dict[str, Any]], Optional[List[str]]]],
) -> Tuple[str, List[Dict[str, Any]], Optional[List[str]]]:
    rows: List[Dict[str, str]] = []

    def add_sheet(sheet_name: str, headers: List[str]) -> None:
        purpose = _dict_sheet_purpose(sheet_name)
        rows.append(
            {
                "Sheet": sheet_name,
                "Sheet Purpose": purpose,
                "Column": "[Sheet overview]",
                "Kind": "overview",
                "Source": "Workbook pipeline",
                "Calculation / Logic": "Not applicable",
                "Quick Explanation": purpose,
            }
        )
        for header in headers:
            kind, source, calc, explanation = _dict_column_metadata(sheet_name, header)
            rows.append(
                {
                    "Sheet": sheet_name,
                    "Sheet Purpose": purpose,
                    "Column": str(header),
                    "Kind": kind,
                    "Source": source,
                    "Calculation / Logic": calc,
                    "Quick Explanation": explanation,
                }
            )

    add_sheet("Summary", ["Metric", "Value"])
    add_sheet("RAW", raw_headers)
    add_sheet("FILTRED", filtered_headers)

    for sheet_name, sheet_rows, sheet_headers in extra_sheets:
        effective_headers = list(sheet_headers or _fieldnames_for_rows(sheet_rows))
        add_sheet(sheet_name, effective_headers)

    add_sheet("Dict", DICT_SHEET_HEADERS)
    return ("Dict", rows, DICT_SHEET_HEADERS)


def _split_ratio_label(value: Any) -> Tuple[int, int, float]:
    raw = str(value or "")
    match = re.search(r"\((\d+)\s*/\s*(\d+)\)", raw)
    if not match:
        return (0, 0, 0.0)
    numerator = int(match.group(1))
    denominator = int(match.group(2))
    percent = (float(numerator) / float(denominator) * 100.0) if denominator else 0.0
    return numerator, denominator, percent


def _format_variation_count(current: int, previous: Optional[int]) -> str:
    if previous is None:
        return "data unavailable"
    diff = current - previous
    return f"({diff:+d})"


def _format_variation_percent(current: float, previous: Optional[float]) -> str:
    if previous is None:
        return "data unavailable"
    diff = int(round(current - previous))
    return f"({diff:+d}%)"


def _last_month_label_from_output(output_path: Path) -> str:
    current = datetime.utcnow()
    stem = output_path.stem
    match = re.search(r"(\d{8})_\d{6}$", stem)
    if match:
        try:
            current = datetime.strptime(match.group(1), "%Y%m%d")
        except ValueError:
            pass
    year = current.year
    month = current.month - 1
    if month == 0:
        month = 12
        year -= 1
    return f"{month:02d}/{year}"


def _label_from_workbook_path(workbook_path: Optional[Path], default_label: str) -> str:
    if workbook_path is None:
        return default_label
    match = re.search(r"(\d{8})_(\d{6})", workbook_path.stem)
    if not match:
        return default_label
    try:
        stamp = datetime.strptime("".join(match.groups()), "%Y%m%d%H%M%S")
    except ValueError:
        return default_label
    return f"{stamp.day:02d}/{stamp.month:02d}/{stamp.year}"


def _load_previous_totals_workbook(output_path: Path) -> Tuple[Dict[str, Dict[Tuple[str, ...], Dict[str, str]]], Optional[Path]]:
    previous_file = _find_totals_baseline_workbook(output_path)
    if previous_file is None:
        return {}, None

    out: Dict[str, Dict[Tuple[str, ...], Dict[str, str]]] = {"TOTAL.PROGRAM": {}, "TOTAL.ENTITY": {}}
    for sheet_name, key_fields in (("TOTAL.PROGRAM", ["Program"]), ("TOTAL.ENTITY", ["Program", "Entity"])):
        rows = _read_table_sheet_from_xlsx(previous_file, sheet_name)
        by_key: Dict[Tuple[str, ...], Dict[str, str]] = {}
        for row in rows:
            key = tuple(str(row.get(field, "")).strip() for field in key_fields)
            if all(key):
                by_key[key] = row
        out[sheet_name] = by_key
    return out, previous_file


def _find_totals_baseline_workbook(output_path: Path) -> Optional[Path]:
    return _find_stats_baseline_workbook(output_path)


def _find_stats_baseline_workbook(output_path: Path) -> Optional[Path]:
    runs_dir = output_path.parent.parent.parent
    if not runs_dir.is_dir():
        return None

    current_stamp = datetime.utcnow()
    match = re.search(r"(\d{8})_(\d{6})$", output_path.stem)
    if match:
        try:
            current_stamp = datetime.strptime("".join(match.groups()), "%Y%m%d%H%M%S")
        except ValueError:
            pass

    previous_month_candidates: List[Tuple[datetime, Path]] = []
    current_month_candidates: List[Tuple[datetime, Path]] = []

    for candidate in runs_dir.glob("*/raw/dali_impact_analysis_*.xlsx"):
        if candidate.resolve() == output_path.resolve():
            continue
        date_match = re.search(r"(\d{8})_(\d{6})", candidate.stem)
        if not date_match:
            continue
        try:
            stamp = datetime.strptime("".join(date_match.groups()), "%Y%m%d%H%M%S")
        except ValueError:
            continue
        if stamp >= current_stamp:
            continue

        same_month = stamp.year == current_stamp.year and stamp.month == current_stamp.month
        previous_month = (
            (current_stamp.month == 1 and stamp.year == current_stamp.year - 1 and stamp.month == 12)
            or (current_stamp.month != 1 and stamp.year == current_stamp.year and stamp.month == current_stamp.month - 1)
        )

        if previous_month:
            previous_month_candidates.append((stamp, candidate))
        elif same_month:
            current_month_candidates.append((stamp, candidate))

    if previous_month_candidates:
        previous_month_candidates.sort(key=lambda item: item[0], reverse=True)
        return previous_month_candidates[0][1]

    if current_month_candidates:
        current_month_candidates.sort(key=lambda item: item[0])
        return current_month_candidates[0][1]

    return None


def _load_previous_stats_workbook(output_path: Path) -> Dict[Tuple[str, str], Dict[str, str]]:
    previous_file = _find_stats_baseline_workbook(output_path)
    if previous_file is None:
        return {}

    rows = _read_table_sheet_from_xlsx(previous_file, "STATS")
    by_key: Dict[Tuple[str, str], Dict[str, str]] = {}
    for row in rows:
        program = str(row.get("Program", "") or "").strip()
        uid = _normalize_lookup_value(row.get("Kear ID", ""))
        if program and uid:
            by_key[(program, uid)] = row
    return by_key


def _read_table_sheet_from_xlsx(workbook_path: Path, sheet_name: str) -> List[Dict[str, str]]:
    ns_main = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
    ns_rel = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    ns_pkg = "{http://schemas.openxmlformats.org/package/2006/relationships}"

    with zipfile.ZipFile(workbook_path, "r") as zf:
        workbook_root = ET.fromstring(zf.read("xl/workbook.xml"))
        rels_root = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))

        rid = None
        for sheet in workbook_root.findall(f"{ns_main}sheets/{ns_main}sheet"):
            if sheet.get("name") == sheet_name:
                rid = sheet.get(f"{ns_rel}id")
                break
        if not rid:
            return []

        target = None
        for rel in rels_root.findall(f"{ns_pkg}Relationship"):
            if rel.get("Id") == rid:
                target = rel.get("Target")
                break
        if not target:
            return []

        sheet_root = ET.fromstring(zf.read(f"xl/{target}"))
        rows = sheet_root.findall(f"{ns_main}sheetData/{ns_main}row")
        if not rows:
            return []

        header_values = [_xlsx_value_from_cell(cell, ns_main) for cell in rows[0].findall(f"{ns_main}c")]
        data_rows: List[Dict[str, str]] = []
        for row in rows[1:]:
            values = [_xlsx_value_from_cell(cell, ns_main) for cell in row.findall(f"{ns_main}c")]
            item = {header_values[idx]: values[idx] if idx < len(values) else "" for idx in range(len(header_values)) if header_values[idx]}
            if any(str(v).strip() for v in item.values()):
                data_rows.append(item)
        return data_rows


def _xlsx_value_from_cell(cell: ET.Element, ns_main: str) -> str:
    inline = cell.find(f"{ns_main}is/{ns_main}t")
    if inline is not None and inline.text is not None:
        return inline.text
    value = cell.find(f"{ns_main}v")
    return "" if value is None or value.text is None else value.text


def _build_total_program_rows(
    recap_rows: List[Dict[str, Any]],
    last_month_label: str,
    previous_rows: Dict[Tuple[str, ...], Dict[str, str]],
) -> Tuple[str, List[Dict[str, Any]], List[str]]:
    headers = [
        "Program",
        "Number of Applications",
        "Total Assets in Dali (in scope)",
        f"Variation Total servers ({last_month_label})",
        "% servers with illumio installed",
        f"Variation % servers with illumio installed ({last_month_label})",
        "% servers with illumio agent in blocking mode",
        f"Variation % servers with illumio agent in blocking mode ({last_month_label})",
    ]

    grouped: Dict[str, Dict[str, Any]] = {}
    for row in recap_rows:
        program = str(row.get("Program", "")).strip()
        if not program:
            continue
        agg = grouped.setdefault(program, {"apps": 0, "assets": 0, "installed_num": 0, "installed_den": 0, "blocking_num": 0, "blocking_den": 0})
        agg["apps"] += 1
        agg["assets"] += int(_coerce_excel_numeric(row.get("Total Assets in Dali (in scope)", "0")) or "0")
        ins_n, ins_d, _ = _split_ratio_label(row.get("% servers with illumio installed", ""))
        blk_n, blk_d, _ = _split_ratio_label(row.get("% servers with illumio agent in blocking mode", ""))
        agg["installed_num"] += ins_n
        agg["installed_den"] += ins_d
        agg["blocking_num"] += blk_n
        agg["blocking_den"] += blk_d

    out_rows: List[Dict[str, Any]] = []
    for program in sorted(grouped, key=lambda value: str(value).upper(), reverse=True):
        agg = grouped[program]
        installed_label = _format_ratio_label(agg["installed_num"], agg["installed_den"])
        blocking_label = _format_ratio_label(agg["blocking_num"], agg["blocking_den"])
        _, _, installed_pct = _split_ratio_label(installed_label)
        _, _, blocking_pct = _split_ratio_label(blocking_label)

        previous = previous_rows.get((program,))
        previous_total = int(_coerce_excel_numeric((previous or {}).get("Total Assets in Dali (in scope)", "")) or "0") if previous else None
        previous_installed = _ratio_percent_from_label((previous or {}).get("% servers with illumio installed", "")) if previous else None
        previous_blocking = _ratio_percent_from_label((previous or {}).get("% servers with illumio agent in blocking mode", "")) if previous else None
        if previous_installed is not None and previous_installed == float("inf"):
            previous_installed = None
        if previous_blocking is not None and previous_blocking == float("inf"):
            previous_blocking = None

        out_rows.append(
            {
                "Program": program,
                "Number of Applications": str(agg["apps"]),
                "Total Assets in Dali (in scope)": str(agg["assets"]),
                f"Variation Total servers ({last_month_label})": _format_variation_count(agg["assets"], previous_total),
                "% servers with illumio installed": installed_label,
                f"Variation % servers with illumio installed ({last_month_label})": _format_variation_percent(installed_pct, previous_installed),
                "% servers with illumio agent in blocking mode": blocking_label,
                f"Variation % servers with illumio agent in blocking mode ({last_month_label})": _format_variation_percent(blocking_pct, previous_blocking),
            }
        )

    out_rows, headers = _append_total_directional_columns(out_rows, headers)
    return ("TOTAL.PROGRAM", out_rows, headers)


def _build_total_entity_rows(
    recap_rows: List[Dict[str, Any]],
    last_month_label: str,
    previous_rows: Dict[Tuple[str, ...], Dict[str, str]],
) -> Tuple[str, List[Dict[str, Any]], List[str]]:
    headers = [
        "Entity",
        "Program",
        "Number of Applications",
        "Total Assets in Dali (in scope)",
        f"Variation Total servers ({last_month_label})",
        "% servers with illumio installed",
        f"Variation % servers with illumio installed ({last_month_label})",
        "% servers with illumio agent in blocking mode",
        f"Variation % servers with illumio agent in blocking mode ({last_month_label})",
    ]

    grouped: Dict[Tuple[str, str], Dict[str, Any]] = {}
    for row in recap_rows:
        program = str(row.get("Program", "")).strip()
        entity = str(row.get("Entity", "")).strip()
        if not program:
            continue
        key = (program, entity)
        agg = grouped.setdefault(key, {"apps": 0, "assets": 0, "installed_num": 0, "installed_den": 0, "blocking_num": 0, "blocking_den": 0})
        agg["apps"] += 1
        agg["assets"] += int(_coerce_excel_numeric(row.get("Total Assets in Dali (in scope)", "0")) or "0")
        ins_n, ins_d, _ = _split_ratio_label(row.get("% servers with illumio installed", ""))
        blk_n, blk_d, _ = _split_ratio_label(row.get("% servers with illumio agent in blocking mode", ""))
        agg["installed_num"] += ins_n
        agg["installed_den"] += ins_d
        agg["blocking_num"] += blk_n
        agg["blocking_den"] += blk_d

    out_rows: List[Dict[str, Any]] = []
    for program, entity in sorted(grouped, key=lambda item: (str(item[1]).upper(), str(item[0]).upper())):
        agg = grouped[(program, entity)]
        installed_label = _format_ratio_label(agg["installed_num"], agg["installed_den"])
        blocking_label = _format_ratio_label(agg["blocking_num"], agg["blocking_den"])
        _, _, installed_pct = _split_ratio_label(installed_label)
        _, _, blocking_pct = _split_ratio_label(blocking_label)

        previous = previous_rows.get((program, entity))
        previous_total = int(_coerce_excel_numeric((previous or {}).get("Total Assets in Dali (in scope)", "")) or "0") if previous else None
        previous_installed = _ratio_percent_from_label((previous or {}).get("% servers with illumio installed", "")) if previous else None
        previous_blocking = _ratio_percent_from_label((previous or {}).get("% servers with illumio agent in blocking mode", "")) if previous else None
        if previous_installed is not None and previous_installed == float("inf"):
            previous_installed = None
        if previous_blocking is not None and previous_blocking == float("inf"):
            previous_blocking = None

        out_rows.append(
            {
                "Entity": entity,
                "Program": program,
                "Number of Applications": str(agg["apps"]),
                "Total Assets in Dali (in scope)": str(agg["assets"]),
                f"Variation Total servers ({last_month_label})": _format_variation_count(agg["assets"], previous_total),
                "% servers with illumio installed": installed_label,
                f"Variation % servers with illumio installed ({last_month_label})": _format_variation_percent(installed_pct, previous_installed),
                "% servers with illumio agent in blocking mode": blocking_label,
                f"Variation % servers with illumio agent in blocking mode ({last_month_label})": _format_variation_percent(blocking_pct, previous_blocking),
            }
        )

    out_rows, headers = _append_total_directional_columns(out_rows, headers)
    return ("TOTAL.ENTITY", out_rows, headers)



def _candidate_effective_env_value(
    marley_doc: Dict[str, Any],
    inventory_row: Dict[str, Any],
    beneficiary_env_map: Dict[str, str],
) -> str:
    beneficiary_value = _normalize_lookup_value(inventory_row.get("beneficiary", ""))
    if beneficiary_value:
        return beneficiary_env_map.get(beneficiary_value, "")
    return _normalize_cell_value(_nested_get(marley_doc, "app_info.env", ""))



def _candidate_filter_debug_columns(
    marley_doc: Dict[str, Any],
    inventory_row: Dict[str, Any],
    beneficiary_env_map: Dict[str, str],
    filters: Optional[Dict[str, str]],
) -> Dict[str, str]:
    env_value = _candidate_effective_env_value(marley_doc, inventory_row, beneficiary_env_map)
    os_value = _normalize_cell_value(marley_doc.get("os_name", ""))
    server_status_value = _normalize_cell_value(marley_doc.get("status", ""))
    domain_value = _normalize_cell_value(marley_doc.get("dns", ""))
    typology_value = _normalize_cell_value(marley_doc.get("typologie", ""))

    env_tokens = _parse_filter_tokens(filters, "FILTER_PRD_ENV")
    os_tokens = _parse_filter_tokens(filters, "FILTER_OS_NAME")
    server_status_tokens = _parse_filter_tokens(filters, "FILTER_SERVER_STATUS")
    domain_tokens = _parse_filter_tokens(filters, "FILTER_DOMAIN_NOT_TAKEN")
    typology_tokens = _parse_filter_tokens(filters, "FILTER_TYPOLOGY_NOT_TAKEN")

    env_ok = True if not env_tokens else _contains_any_token(env_value, env_tokens)
    os_ok = True if not os_tokens else _matches_exact_token(os_value, os_tokens)
    server_status_ok = True if not server_status_tokens else _matches_exact_token(server_status_value, server_status_tokens)
    domain_ok = not _contains_any_token(domain_value, domain_tokens) if domain_tokens else True
    typology_ok = not _contains_any_token(typology_value, typology_tokens) if typology_tokens else True

    return {
        "FILTER_VALUE_environment": env_value,
        "F_FILTER_PRD_ENV": "Y" if env_ok else "N",
        "FILTER_VALUE_os_name": os_value,
        "F_FILTER_OS_NAME": "Y" if os_ok else "N",
        "FILTER_VALUE_server.status": server_status_value,
        "F_FILTER_SERVER_STATUS": "Y" if server_status_ok else "N",
        "FILTER_VALUE_domain": domain_value,
        "F_FILTER_DOMAIN": "Y" if domain_ok else "N",
        "FILTER_VALUE_typology": typology_value,
        "F_FILTER_TYPOLOGY_NOT_TAKEN": "Y" if typology_ok else "N",
        "F_FILTER_ALL": "Y" if all([env_ok, os_ok, server_status_ok, domain_ok, typology_ok]) else "N",
    }


def build_candidate_sheet(
    workload_csv: Path,
    monitored_rows: List[Dict[str, str]],
    scope_rows: Optional[List[Dict[str, Any]]] = None,
    filters: Optional[Dict[str, str]] = None,
) -> Tuple[str, List[Dict[str, Any]], List[str]]:
    """Build candidate workload rows enriched with Marley, inventory and monitored-scope matching."""
    workload_rows = _read_workload_derived_rows(workload_csv)
    candidate_source_rows = [
        row
        for row in workload_rows
        if str(_get_row_value_by_candidates(row, ["IPLIST", "IPLISTS"])).strip()
    ]

    d4s_client = Data4secClient()

    candidate_marley_source_fields = sorted(
        set(list(QUERY_CONFIG.get("marley_original", {}).get("source_fields", [])) + [
            "uuid",
            "status",
            "app_info.kear_uuid",
            "app_info.factor",
            "app_info.app_id",
            "app_info.app_name",
            "app_info.env",
            "dns",
            "typologie",
            "os_name",
        ])
    )
    candidate_inventory_source_fields = sorted(
        set(list(QUERY_CONFIG.get("inventory", {}).get("source_fields", [])) + [
            "hostname",
            "ocs_name",
            "owner_app_name",
            "beneficiary",
        ])
    )

    short_hostnames = sorted({
        _normalize_lookup_value(_get_row_value_by_candidates(row, ["short_hostname"]))
        for row in candidate_source_rows
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["short_hostname"]))
    })
    hostnames = sorted({
        _normalize_lookup_value(_get_row_value_by_candidates(row, ["hostname"]))
        for row in candidate_source_rows
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["hostname"]))
    })
    ocs_names_from_ip = sorted({
        _normalize_lookup_value(_get_row_value_by_candidates(row, ["ocs_name_from_IP"]))
        for row in candidate_source_rows
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["ocs_name_from_IP"]))
    })

    marley_docs_by_short_hostname = query_marley_original_by_field(
        client=d4s_client,
        lookup_values=short_hostnames,
        search_field="hostname",
        lookup_label="candidate_short_hostnames",
        source_fields_override=candidate_marley_source_fields,
    ) if short_hostnames else {}
    marley_docs_by_hostname = query_marley_original_by_field(
        client=d4s_client,
        lookup_values=hostnames,
        search_field="hostname",
        lookup_label="candidate_hostnames",
        source_fields_override=candidate_marley_source_fields,
    ) if hostnames else {}
    marley_docs_by_ocs_name_from_ip = query_marley_original_by_field(
        client=d4s_client,
        lookup_values=ocs_names_from_ip,
        search_field="hostname",
        lookup_label="candidate_ocs_name_from_ip",
        source_fields_override=candidate_marley_source_fields,
    ) if ocs_names_from_ip else {}

    inventory_docs_by_short_hostname = _inventory_search_by_field(
        client=d4s_client,
        search_field="hostname",
        values=short_hostnames,
        source_fields_override=candidate_inventory_source_fields,
    ) if short_hostnames else {}
    inventory_docs_by_hostname = _inventory_search_by_field(
        client=d4s_client,
        search_field="hostname",
        values=hostnames,
        source_fields_override=candidate_inventory_source_fields,
    ) if hostnames else {}
    inventory_docs_by_ocs_name_from_ip = _inventory_search_by_field(
        client=d4s_client,
        search_field="ocs_name",
        values=ocs_names_from_ip,
        source_fields_override=candidate_inventory_source_fields,
    ) if ocs_names_from_ip else {}

    beneficiary_values = sorted({
        _normalize_lookup_value(doc.get("beneficiary", ""))
        for result_map in (inventory_docs_by_short_hostname, inventory_docs_by_hostname, inventory_docs_by_ocs_name_from_ip)
        for docs in result_map.values()
        for doc in docs
        if _normalize_lookup_value(doc.get("beneficiary", ""))
    })
    beneficiary_env_map = query_platform_accounts_env_by_names(d4s_client, beneficiary_values) if beneficiary_values else {}

    monitored_index: List[Tuple[str, str, str]] = []
    for monitored_row in monitored_rows:
        monitored_uid = _normalize_lookup_value(monitored_row.get("uid", "") or monitored_row.get("kear", ""))
        monitored_network = _normalize_lookup_value(monitored_row.get("network", ""))
        monitored_program = str(monitored_row.get("program", "") or "").strip()
        if monitored_uid and monitored_network:
            monitored_index.append((monitored_uid, monitored_network, monitored_program))

    scope_server_uids = {
        _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
        for row in (scope_rows or [])
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
    }

    candidate_rows: List[Dict[str, Any]] = []
    for source_row in candidate_source_rows:
        short_hostname = _get_row_value_by_candidates(source_row, ["short_hostname"]).strip()
        hostname = _get_row_value_by_candidates(source_row, ["hostname"]).strip()
        ocs_name_from_ip = _get_row_value_by_candidates(source_row, ["ocs_name_from_IP"]).strip()
        normalized_short_hostname = _normalize_lookup_value(short_hostname)
        normalized_hostname = _normalize_lookup_value(hostname)
        normalized_ocs_name_from_ip = _normalize_lookup_value(ocs_name_from_ip)

        marley_docs = marley_docs_by_short_hostname.get(normalized_short_hostname, [])
        if not marley_docs and normalized_hostname:
            marley_docs = marley_docs_by_hostname.get(normalized_hostname, [])
        if not marley_docs and normalized_ocs_name_from_ip:
            marley_docs = marley_docs_by_ocs_name_from_ip.get(normalized_ocs_name_from_ip, [])
        marley_doc = marley_docs[0] if marley_docs else {}
        lookup_status = "FOUND" if marley_doc else "NOT_FOUND"

        inventory_docs = inventory_docs_by_short_hostname.get(normalized_short_hostname, [])
        if not inventory_docs and normalized_hostname:
            inventory_docs = inventory_docs_by_hostname.get(normalized_hostname, [])
        if not inventory_docs and normalized_ocs_name_from_ip:
            inventory_docs = inventory_docs_by_ocs_name_from_ip.get(normalized_ocs_name_from_ip, [])
        inventory_doc = inventory_docs[0] if inventory_docs else {}

        filter_columns = _candidate_filter_debug_columns(
            marley_doc=marley_doc,
            inventory_row=inventory_doc,
            beneficiary_env_map=beneficiary_env_map,
            filters=filters,
        )

        kear_uuid_raw = _normalize_cell_value(_nested_get(marley_doc, "app_info.kear_uuid", ""))
        normalized_kear_tokens = {
            _normalize_lookup_value(token)
            for token in re.split(r"\s*,\s*", kear_uuid_raw)
            if _normalize_lookup_value(token)
        }
        iplist_value = _get_row_value_by_candidates(source_row, ["IPLIST", "IPLISTS"]).strip()
        normalized_iplist = _normalize_lookup_value(iplist_value)

        matched_programs: List[str] = []
        for monitored_uid, monitored_network, monitored_program in monitored_index:
            if monitored_uid not in normalized_kear_tokens:
                continue
            if not monitored_network or not normalized_iplist or monitored_network not in normalized_iplist:
                continue
            if monitored_program and monitored_program not in matched_programs:
                matched_programs.append(monitored_program)

        beneficiary_value = _normalize_cell_value(inventory_doc.get("beneficiary", ""))
        beneficiary_env_value = beneficiary_env_map.get(_normalize_lookup_value(beneficiary_value), "") if beneficiary_value else ""

        candidate_row = {
            "short_hostname": short_hostname,
            "status": _normalize_cell_value(marley_doc.get("status", "")),
            "uuid": _normalize_cell_value(marley_doc.get("uuid", "")),
            "IPLIST": iplist_value,
            "scope_alread_monitored?": "Y" if matched_programs else "N",
            "Program": ",".join(matched_programs),
            "owner_app_name": _normalize_cell_value(inventory_doc.get("owner_app_name", "")),
            "beneficiary": beneficiary_value,
            "beneficiary_account_env": beneficiary_env_value,
            "hostname": hostname,
            "ocs_name_from_IP": ocs_name_from_ip,
            "SUBNET": _get_row_value_by_candidates(source_row, ["SUBNET"]).strip(),
            "managed": _get_row_value_by_candidates(source_row, ["managed"]).strip(),
            "enforcement": _get_row_value_by_candidates(source_row, ["enforcement"]).strip(),
            "created_at": _get_row_value_by_candidates(source_row, ["created_at"]).strip(),
            "ip_with_default_gw": _get_row_value_by_candidates(source_row, ["ip_with_default_gw"]).strip(),
            "app": _get_row_value_by_candidates(source_row, ["app"]).strip(),
            "env": _get_row_value_by_candidates(source_row, ["env"]).strip(),
            "loc": _get_row_value_by_candidates(source_row, ["loc"]).strip(),
            "role": _get_row_value_by_candidates(source_row, ["role"]).strip(),
            "app_info.kear_uuid": kear_uuid_raw,
            "app_info.factor": _normalize_cell_value(_nested_get(marley_doc, "app_info.factor", "")),
            "app_info.app_id": _normalize_cell_value(_nested_get(marley_doc, "app_info.app_id", "")),
            "app_info.app_name": _normalize_cell_value(_nested_get(marley_doc, "app_info.app_name", "")),
            "app_info.env": _normalize_cell_value(_nested_get(marley_doc, "app_info.env", "")),
            "typologie": _normalize_cell_value(marley_doc.get("typologie", "")),
            "os_name": _normalize_cell_value(marley_doc.get("os_name", "")),
            "dns": _normalize_cell_value(marley_doc.get("dns", "")),
            "lookup_status": lookup_status,
        }
        candidate_row.update(filter_columns)

        candidate_uuid = _normalize_lookup_value(candidate_row.get("uuid", ""))
        candidate_status = str(candidate_row.get("status", "") or "").strip()
        if candidate_uuid and candidate_uuid in scope_server_uids:
            continue
        if not candidate_status:
            continue

        candidate_rows.append(candidate_row)

    return ("CANDIDATE", candidate_rows, CANDIDATE_SHEET_HEADERS)

def build_kear_labels_accounts_sheet(
    filtered_rows: List[Dict[str, Any]],
    workload_csv: Path,
) -> Tuple[str, List[Dict[str, Any]], List[str]]:
    headers = [
        "program",
        "uid",
        "UID REL",
        "SHORT LABEL REL",
        "network",
        "IPLIST",
        "DSI REL",
        "ENVIRONMENT",
        "CLOUD TYPE",
        "Retrieved from",
        "INV_Owner_Account",
        "INV_Beneficiary_Account",
        "managed",
        "role",
        "app",
        "env",
        "Count In scope",
        "Count in PCE",
    ]

    key_columns = [
        "program",
        "uid",
        "UID REL",
        "SHORT LABEL REL",
        "network",
        "IPLIST",
        "DSI REL",
        "ENVIRONMENT",
        "CLOUD TYPE",
        "Retrieved from",
        "INV_Owner_Account",
        "INV_Beneficiary_Account",
        "role",
        "app",
        "env",
    ]
    distinct_keys: set[Tuple[str, ...]] = set()
    in_scope_counts: Dict[Tuple[str, ...], int] = {}

    def _row_key(row: Dict[str, Any]) -> Tuple[str, ...]:
        return tuple(
            str(_get_row_value_by_candidates(row, [column if column != "Retrieved from" else "Retrived from"]) or "").strip()
            for column in key_columns
        )

    for row in filtered_rows:
        key = _row_key(row)
        distinct_keys.add(key)
        if _is_truthy_flag(row.get("In scope", "")):
            in_scope_counts[key] = in_scope_counts.get(key, 0) + 1

    workload_rows = _read_workload_derived_rows(workload_csv)
    pce_counts: Dict[Tuple[str, str, str, str, str], int] = {}
    managed_values_by_combo: Dict[Tuple[str, str, str, str], List[str]] = {}
    for row in workload_rows:
        combo_wo_managed = (
            str(row.get("role", "")).strip(),
            str(row.get("app", "")).strip(),
            str(row.get("env", "")).strip(),
            str(row.get("IPLIST", "")).strip(),
        )
        managed_value = str(_get_row_value_by_candidates(row, ["ILU_managed", "managed"])).strip()
        combo = combo_wo_managed + (managed_value,)
        pce_counts[combo] = pce_counts.get(combo, 0) + 1
        if combo_wo_managed not in managed_values_by_combo:
            managed_values_by_combo[combo_wo_managed] = []
        if managed_value not in managed_values_by_combo[combo_wo_managed]:
            managed_values_by_combo[combo_wo_managed].append(managed_value)

    out_rows: List[Dict[str, Any]] = []
    for key in sorted(distinct_keys):
        base_item = {key_columns[idx]: value for idx, value in enumerate(key)}
        base_item["Count In scope"] = str(in_scope_counts.get(key, 0))
        combo_wo_managed = (base_item.get("role", ""), base_item.get("app", ""), base_item.get("env", ""), base_item.get("IPLIST", ""))
        managed_values = managed_values_by_combo.get(combo_wo_managed, [""])

        for managed_value in managed_values:
            item = dict(base_item)
            item["managed"] = managed_value
            pce_key = combo_wo_managed + (managed_value,)
            item["Count in PCE"] = str(pce_counts.get(pce_key, 0))
            out_rows.append(item)

    return ("KearLabelsAccounts", out_rows, headers)


def _natural_slide_sort_key(value: Any) -> Tuple[int, str]:
    raw = str(value or "").strip()
    match = re.search(r"(\d+)", raw)
    if match:
        return (int(match.group(1)), raw.upper())
    return (10**9, raw.upper())


def _pptx_indicator_symbol(value: Any) -> Tuple[str, Optional[Tuple[int, int, int]]]:
    """Render indicator buckets using the validated Excel business rule.

    - red:    <= 0
    - yellow: > 0 and < 100
    - green:  >= 100

    The PPT rendering must follow the helper value carried by the
    ``Indicator Icon`` column and must not infer thresholds from the
    visible percentage label next to it.
    """
    numeric = _coerce_excel_numeric(value)
    if numeric is None:
        return ("", None)
    try:
        parsed = float(numeric)
    except ValueError:
        return (str(value or ""), None)
    if parsed <= 0.0:
        return ("■", (192, 0, 0))
    if parsed >= 100.0:
        return ("■", (0, 128, 0))
    return ("■", (191, 144, 0))


def _pptx_trend_symbol(value: Any) -> Tuple[str, Optional[Tuple[int, int, int]]]:
    """Render trend buckets with a stability window around zero.

    Trend helper values are percentage-point deltas. Small fluctuations should
    be treated as stable instead of positive/negative trend.
    - up:      >= +1.0 pt
    - down:    <= -1.0 pt
    - stable:  between -1.0 and +1.0 pt
    """
    numeric = _coerce_excel_numeric(value)
    if numeric is None:
        return ("", None)
    try:
        parsed = float(numeric)
    except ValueError:
        return (str(value or ""), None)
    if parsed >= 1.0:
        return ("↑", (0, 128, 0))
    if parsed <= -1.0:
        return ("↓", (192, 0, 0))
    return ("→", (96, 96, 96))


def _pptx_render_cell_value(column_name: str, value: Any) -> Tuple[str, Optional[Tuple[int, int, int]]]:
    column = str(column_name or "").strip()
    if column.endswith("Indicator Icon"):
        return _pptx_indicator_symbol(value)
    if column.endswith("Trend Icon"):
        return _pptx_trend_symbol(value)
    return (str(value or ""), None)


def _pptx_header_label(column_name: Any) -> str:
    raw = str(column_name or "").strip()
    if raw == "Index":
        return "#"
    return "" if ("Indicator" in raw or "Trend" in raw) else raw


def _pptx_column_alignment(column_name: Any, table_kind: str = "default") -> str:
    raw = str(column_name or "").strip()
    kind = str(table_kind or "default").strip().lower()
    if "Indicator" in raw or "Trend" in raw:
        return "center"
    if kind in {"total_program", "total_entity"} and "Variation" in raw:
        return "center"
    if raw in {"Index", "Program", "Entity", "Application Short Label"}:
        return "left"
    if "%" in raw:
        return "right"
    if raw in {
        "Number of Applications",
        "Total Assets in Dali (in scope)",
        "Variation Total servers",
        "Total Assets in Dali (Enriched)",
        "Assets in Dali not in illumio",
        "Assets in Dali (Enriched) not in illumio",
    }:
        return "center"
    return "left"


def _pptx_body_font_size(column_name: Any, default_font_size_pt: float, table_kind: str = "default") -> float:
    raw = str(column_name or "").strip()
    kind = str(table_kind or "default").strip().lower()
    if kind == "total_program":
        return default_font_size_pt
    if kind == "stats" and raw in {
        "% servers with illumio installed (Enriched) Indicator Icon",
        "% servers with illumio installed (Enriched) Trend Icon",
        "% servers with illumio agent in blocking mode (Enriched) Indicator Icon",
        "% servers with illumio agent in blocking mode (Enriched) Trend Icon",
    }:
        return 14.0
    if kind == "stats" and raw in {"Program", "Entity", "Sub-Entity", "Application Short Label"}:
        return 10.0
    if raw in {"Index", "Program", "Entity", "Application Short Label"}:
        return 10.0
    return default_font_size_pt


def _pptx_apply_cell_border(cell: Any, color_rgb: Tuple[int, int, int] = (31, 78, 121), width_emu: int = 3175) -> None:
    from pptx.oxml.xmlchemy import OxmlElement

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for child in list(tcPr):
        if child.tag.endswith(("lnL", "lnR", "lnT", "lnB")):
            tcPr.remove(child)

    for side in ("L", "R", "T", "B"):
        ln = OxmlElement(f"a:ln{side}")
        ln.set("w", str(width_emu))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        solid_fill = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", f"{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}")
        solid_fill.append(srgb)
        ln.append(solid_fill)

        prst_dash = OxmlElement("a:prstDash")
        prst_dash.set("val", "solid")
        ln.append(prst_dash)

        round_node = OxmlElement("a:round")
        ln.append(round_node)

        head_end = OxmlElement("a:headEnd")
        head_end.set("type", "none")
        head_end.set("w", "med")
        head_end.set("len", "med")
        ln.append(head_end)

        tail_end = OxmlElement("a:tailEnd")
        tail_end.set("type", "none")
        tail_end.set("w", "med")
        tail_end.set("len", "med")
        ln.append(tail_end)

        tcPr.append(ln)


def _pptx_clear_table_style(table: Any) -> None:
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        return
    for child in list(tblPr):
        if child.tag.endswith("tableStyleId"):
            tblPr.remove(child)
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")


def _pptx_set_cell_style(
    cell: Any,
    text_value: str,
    font_size_pt: float,
    bold: bool = False,
    fill_rgb: Optional[Tuple[int, int, int]] = None,
    font_rgb: Optional[Tuple[int, int, int]] = None,
    alignment: str = "left",
    border_rgb: Tuple[int, int, int] = (31, 78, 121),
    margin_left_pt: float = 2.5,
    margin_right_pt: float = 2.5,
    margin_top_pt: float = 1.4,
    margin_bottom_pt: float = 1.4,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    alignment_value = str(alignment or "left").strip().lower()
    if alignment_value == "center":
        paragraph_alignment = PP_ALIGN.CENTER
    elif alignment_value == "right":
        paragraph_alignment = PP_ALIGN.RIGHT
    else:
        paragraph_alignment = PP_ALIGN.LEFT

    cell.text = str(text_value or "")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(margin_left_pt)
    cell.margin_right = Pt(margin_right_pt)
    cell.margin_top = Pt(margin_top_pt)
    cell.margin_bottom = Pt(margin_bottom_pt)
    if fill_rgb is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*fill_rgb)
    else:
        cell.fill.background()

    text_frame = cell.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = paragraph_alignment
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.text = str(text_value or "")
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.name = "Calibri"
    if font_rgb is not None:
        run.font.color.rgb = RGBColor(*font_rgb)

    for paragraph in text_frame.paragraphs[1:]:
        paragraph.alignment = paragraph_alignment
        for run in paragraph.runs:
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold
            run.font.name = "Calibri"
            if font_rgb is not None:
                run.font.color.rgb = RGBColor(*font_rgb)

    _pptx_apply_cell_border(cell, color_rgb=border_rgb)


def _pptx_apply_column_widths(
    table_shape: Any,
    fieldnames: List[str],
    available_width_in: float,
    width_weights: Dict[str, float],
    table_kind: str = "default",
) -> None:
    from pptx.util import Inches

    normalized_kind = str(table_kind or "default").strip().lower()
    if normalized_kind == "stats":
        explicit_widths = {str(field): float(width_weights.get(str(field), 0.0) or 0.0) for field in fieldnames}
        remaining_columns = [field for field in fieldnames if explicit_widths.get(str(field), 0.0) <= 0.0]
        used_width = sum(explicit_widths.get(str(field), 0.0) for field in fieldnames if explicit_widths.get(str(field), 0.0) > 0.0)
        remaining_width = max(0.0, available_width_in - used_width)
        auto_width = (remaining_width / len(remaining_columns)) if remaining_columns else 0.0
        for idx, field in enumerate(fieldnames):
            width_in = explicit_widths.get(str(field), 0.0)
            if width_in <= 0.0:
                width_in = auto_width
            table_shape.table.columns[idx].width = Inches(width_in)
        return

    total_weight = sum(width_weights.get(str(field), 1.0) for field in fieldnames) or 1.0
    for idx, field in enumerate(fieldnames):
        table_shape.table.columns[idx].width = Inches(available_width_in * (width_weights.get(str(field), 1.0) / total_weight))


def _pptx_add_table_slide(
    prs: Any,
    title: str,
    rows: List[Dict[str, Any]],
    fieldnames: List[str],
    width_weights: Dict[str, float],
    rows_per_slide: int,
    body_font_size: float,
    table_kind: str = "default",
    header_font_size: float = 10.0,
    header_height_in: float = 0.24,
    body_height_in: float = 0.29,
    chunk_title_builder: Optional[Any] = None,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    slide_width_in = prs.slide_width / 914400.0
    slide_height_in = prs.slide_height / 914400.0
    normalized_table_kind = str(table_kind or "default").strip().lower()
    if normalized_table_kind == "stats":
        left_in = 0.14
        top_in = 0.14
        title_h_in = 0.26
        bottom_margin_in = 0.05
        table_top_in = top_in + title_h_in + 0.08
    else:
        left_in = 0.18
        top_in = 0.55
        title_h_in = 0.35
        bottom_margin_in = 0.2
        table_top_in = top_in + title_h_in + 0.16
    table_h_in = slide_height_in - table_top_in - bottom_margin_in
    table_w_in = slide_width_in - (left_in * 2)

    chunk_size = max(1, rows_per_slide)
    chunks = [rows[idx: idx + chunk_size] for idx in range(0, len(rows), chunk_size)] or [[]]

    for chunk_idx, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(table_w_in), Inches(title_h_in))
        title_tf = title_box.text_frame
        title_tf.clear()
        title_p = title_tf.paragraphs[0]
        if chunk_title_builder is not None:
            computed_title = str(chunk_title_builder(chunk) or "").strip()
            title_text = computed_title or title
        else:
            title_text = title if len(chunks) == 1 else f"{title} ({chunk_idx}/{len(chunks)})"

        if normalized_table_kind == "stats" and ":" in title_text:
            left_title, right_title = title_text.split(":", 1)
            left_run = title_p.add_run()
            left_run.text = left_title
            left_run.font.size = Pt(18)
            left_run.font.bold = True
            left_run.font.name = "Calibri"
            left_run.font.color.rgb = RGBColor(139, 0, 0)

            right_run = title_p.add_run()
            right_run.text = f":{right_title}"
            right_run.font.size = Pt(18)
            right_run.font.bold = True
            right_run.font.name = "Calibri"
        else:
            title_p.text = title_text
            title_p.runs[0].font.size = Pt(18)
            title_p.runs[0].font.bold = True
            title_p.runs[0].font.name = "Calibri"

        table_shape = slide.shapes.add_table(len(chunk) + 1, len(fieldnames), Inches(left_in), Inches(table_top_in), Inches(table_w_in), Inches(table_h_in))
        table = table_shape.table
        _pptx_clear_table_style(table)
        _pptx_apply_column_widths(table_shape, fieldnames, table_w_in, width_weights, table_kind=table_kind)

        header_fill = (31, 78, 121)
        alt_fill = (247, 249, 252)
        white_fill = (255, 255, 255)
        header_font_rgb = (255, 255, 255)
        border_rgb = (31, 78, 121)

        total_requested_h = header_height_in + (body_height_in * len(chunk))
        scale = min(1.0, table_h_in / total_requested_h) if total_requested_h > 0 else 1.0
        table.rows[0].height = Inches(header_height_in * scale)
        for row_idx in range(1, len(chunk) + 1):
            table.rows[row_idx].height = Inches(body_height_in * scale)

        for col_idx, fieldname in enumerate(fieldnames):
            _pptx_set_cell_style(
                table.cell(0, col_idx),
                _pptx_header_label(fieldname),
                font_size_pt=header_font_size,
                bold=True,
                fill_rgb=header_fill,
                font_rgb=header_font_rgb,
                alignment="center",
                border_rgb=border_rgb,
            )

        for row_idx, row in enumerate(chunk, start=1):
            row_fill = alt_fill if row_idx % 2 == 0 else white_fill
            for col_idx, fieldname in enumerate(fieldnames):
                rendered_text, font_rgb = _pptx_render_cell_value(str(fieldname), row.get(fieldname, ""))
                is_symbol_col = ("Indicator" in str(fieldname) or "Trend" in str(fieldname))
                is_program_bold = (str(fieldname or "").strip() == "Program" and normalized_table_kind in {"total_program", "stats"})
                symbol_font_size = 14.0 if normalized_table_kind == "stats" else 14.0
                _pptx_set_cell_style(
                    table.cell(row_idx, col_idx),
                    rendered_text,
                    font_size_pt=(symbol_font_size if is_symbol_col else _pptx_body_font_size(fieldname, body_font_size, table_kind=table_kind)),
                    bold=(is_symbol_col or is_program_bold),
                    fill_rgb=row_fill,
                    font_rgb=font_rgb,
                    alignment=_pptx_column_alignment(fieldname, table_kind=table_kind),
                    border_rgb=border_rgb,
                )


def _pptx_reindex_stats_rows(rows: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    by_program: Dict[str, List[Dict[str, Any]]] = {}
    for row in rows:
        program = str(row.get("Program", "") or "").strip()
        by_program.setdefault(program, []).append(dict(row))

    out_rows: List[Dict[str, Any]] = []
    for program in sorted(by_program.keys(), key=lambda value: value.upper()):
        program_rows = by_program[program]
        program_rows.sort(key=lambda item: int(_coerce_excel_numeric(item.get("Index", "0")) or "0"))
        for idx, item in enumerate(program_rows, start=1):
            item["Index"] = str(idx)
            out_rows.append(item)
    return out_rows


def _pptx_stats_slide_title(rows: List[Dict[str, Any]]) -> str:
    programs: List[str] = []
    entities: List[str] = []
    seen_programs: set[str] = set()
    seen_entities: set[str] = set()

    for row in rows:
        program = str(row.get("Program", "") or "").strip()
        entity = str(row.get("Entity", "") or "").strip()
        if program and program.upper() not in seen_programs:
            seen_programs.add(program.upper())
            programs.append(program)
        if entity and entity.upper() not in seen_entities:
            seen_entities.add(entity.upper())
            entities.append(entity)

    left = " / ".join(programs) if programs else "UNASSIGNED"
    right = " / ".join(entities) if entities else "NO ENTITY"
    return f"{left}:{right}"


def write_output_pptx(
    output_file: str,
    total_program_sheet: Tuple[str, List[Dict[str, Any]], Optional[List[str]]],
    total_entity_sheet: Tuple[str, List[Dict[str, Any]], Optional[List[str]]],
    stats_sheet: Tuple[str, List[Dict[str, Any]], Optional[List[str]]],
    monitored_rows: List[Dict[str, str]],
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        raise RuntimeError("python-pptx is required to generate the PPTX output") from exc

    output_path = Path(output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total_program_headers = list(total_program_sheet[2] or _fieldnames_for_rows(total_program_sheet[1]))
    total_entity_headers = list(total_entity_sheet[2] or _fieldnames_for_rows(total_entity_sheet[1]))
    stats_headers = [column for column in PPTX_STATS_COLUMNS if column in (stats_sheet[2] or [])]

    total_program_widths = {
        "Program": 1.7,
        "Entity": 1.3,
        "Number of Applications": 0.88,
        "Total Assets in Dali (in scope)": 1.02,
        "Variation Total servers": 1.2,
        "% servers with illumio installed": 1.25,
        "Variation % servers with illumio installed": 1.2,
        "% servers with illumio agent in blocking mode": 1.3,
        "Variation % servers with illumio agent in blocking mode": 1.25,
        "% servers with illumio installed Trend Icon": 0.24,
        "% servers with illumio agent in blocking mode Trend Icon": 0.24,
    }
    total_entity_widths = {
        "Program": 1.65,
        "Entity": 1.15,
        "Number of Applications": 0.88,
        "Total Assets in Dali (in scope)": 1.08,
        "Variation Total servers": 1.18,
        "% servers with illumio installed": 1.2,
        "Variation % servers with illumio installed": 1.16,
        "% servers with illumio agent in blocking mode": 1.25,
        "Variation % servers with illumio agent in blocking mode": 1.2,
        "% servers with illumio installed Trend Icon": 0.24,
        "% servers with illumio agent in blocking mode Trend Icon": 0.24,
    }
    stats_widths = {
        "Index": 0.236,
        "Program": 1.496,
        "Entity": 0.787,
        "Sub-Entity": 1.654,
        "Application Short Label": 0.0,
        "Total Assets in Dali (Enriched)": 0.846,
        "Assets in Dali (Enriched) not in illumio": 0.886,
        "% servers with illumio installed (Enriched)": 1.378,
        "% servers with illumio installed (Enriched) Indicator Icon": 0.276,
        "% servers with illumio installed (Enriched) Trend Icon": 0.276,
        "% servers with illumio agent in blocking mode (Enriched)": 1.358,
        "% servers with illumio agent in blocking mode (Enriched) Indicator Icon": 0.276,
        "% servers with illumio agent in blocking mode (Enriched) Trend Icon": 0.276,
    }

    _pptx_add_table_slide(
        prs,
        "TOTAL.PROGRAM",
        total_program_sheet[1],
        total_program_headers,
        total_program_widths,
        rows_per_slide=11,
        body_font_size=14.0,
        table_kind="total_program",
        header_font_size=12.0,
        header_height_in=0.28,
        body_height_in=0.33,
    )
    _pptx_add_table_slide(
        prs,
        "TOTAL.ENTITY",
        total_entity_sheet[1],
        total_entity_headers,
        total_entity_widths,
        rows_per_slide=11,
        body_font_size=12.0,
        table_kind="total_entity",
        header_font_size=10.0,
        header_height_in=0.25,
        body_height_in=0.31,
    )

    slide_by_key: Dict[Tuple[str, str], str] = {}
    for monitored_row in monitored_rows:
        program = str(monitored_row.get("program", "") or "").strip()
        uid = _normalize_lookup_value(monitored_row.get("uid", "") or monitored_row.get("kear", ""))
        slide_name = str(monitored_row.get("slide", "") or "").strip()
        if program and uid and slide_name:
            slide_by_key.setdefault((program, uid), slide_name)

    stats_groups: Dict[str, List[Dict[str, Any]]] = {}
    for row in stats_sheet[1]:
        program = str(row.get("Program", "") or "").strip()
        uid = _normalize_lookup_value(row.get("Kear ID", ""))
        slide_name = slide_by_key.get((program, uid), "UNASSIGNED")
        projected_row = {column: row.get(column, "") for column in stats_headers}
        stats_groups.setdefault(slide_name, []).append(projected_row)

    for slide_name in sorted(stats_groups.keys(), key=_natural_slide_sort_key):
        stats_rows = _pptx_reindex_stats_rows(stats_groups[slide_name])
        _pptx_add_table_slide(
            prs,
            slide_name,
            stats_rows,
            stats_headers,
            stats_widths,
            rows_per_slide=14,
            body_font_size=12.0,
            table_kind="stats",
            header_font_size=10.0,
            header_height_in=0.23,
            body_height_in=0.26,
            chunk_title_builder=_pptx_stats_slide_title,
        )

    prs.save(str(output_path))



def write_output_xlsx(
    # Dynamic XLSX writer supporting optional diagnostic sheets
    output_file: str,
    raw_rows: List[Dict[str, Any]],
    filtered_rows: List[Dict[str, Any]],
    mappings: List[Tuple[str, str]],
    summary_rows: List[Tuple[str, str]],
    filtered_extra_fieldnames: Optional[List[str]] = None,
    raw_extra_fieldnames: Optional[List[str]] = None,
    raw_base_fieldnames: Optional[List[str]] = None,
    filtered_base_fieldnames: Optional[List[str]] = None,
    filtered_fieldnames_override: Optional[List[str]] = None,
    extra_sheets: Optional[List[Tuple[str, List[Dict[str, Any]], Optional[List[str]]]]] = None,
) -> None:
    output_path = Path(output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    effective_raw_base = ["uid", "program", "network", "taken", "Server UID"] if raw_base_fieldnames is None else list(raw_base_fieldnames)
    effective_filtered_base = ["uid", "program", "network", "taken", "Server UID"] if filtered_base_fieldnames is None else list(filtered_base_fieldnames)
    raw_fieldnames = effective_raw_base + [display for display, _ in mappings] + (raw_extra_fieldnames or [])
    filtered_fieldnames = (
        list(filtered_fieldnames_override)
        if filtered_fieldnames_override is not None
        else (effective_filtered_base + [display for display, _ in mappings] + (filtered_extra_fieldnames or []))
    )

    sheets: List[Tuple[str, str, Optional[List[Dict[str, Any]]], Optional[List[str]], Optional[set[str]], Optional[set[str]], bool, Optional[float], Optional[List[Optional[float]]], Optional[set[str]], bool]] = [
        ("Summary", "summary", None, None, None, None, False, None, None, None, False),
        ("RAW", "table", raw_rows, raw_fieldnames, {name for name in raw_fieldnames if str(name).startswith("F_")}, None, False, None, None, None, False),
        ("FILTRED", "table", filtered_rows, filtered_fieldnames, None, None, False, None, None, None, False),
    ]
    for name, rows, fieldnames in (extra_sheets or []):
        effective_fields = fieldnames or _fieldnames_for_rows(rows)
        enriched_columns = STATS_ENRICHED_COLUMNS if name == "STATS" else None
        shaded_columns = ({column for column in effective_fields if str(column).startswith("F_FILTER_")} if name in {"ENRICH", "SCOPE"} else ({column for column in effective_fields if str(column).startswith("F_") or str(column).startswith("FILTER_VALUE_")} if name == "CANDIDATE" else None))
        is_total_sheet = name in {"TOTAL.PROGRAM", "TOTAL.ENTITY"}
        is_stats_sheet = name == "STATS"
        header_multiline = is_total_sheet or is_stats_sheet
        header_height = 40.0 if header_multiline else None
        if is_total_sheet:
            fixed_widths = _fixed_total_sheet_widths(effective_fields)
        elif is_stats_sheet:
            fixed_widths = _fixed_stats_sheet_widths(effective_fields)
        else:
            fixed_widths = None
        hidden_header_columns = STATS_ICON_HEADER_COLUMNS if is_stats_sheet else None
        sheets.append((name, "table", rows, effective_fields, shaded_columns, enriched_columns, header_multiline, header_height, fixed_widths, hidden_header_columns, name in HIDDEN_WORKSHEET_NAMES))

    content_types_parts = [
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">',
        '  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>',
        '  <Default Extension="xml" ContentType="application/xml"/>',
        '  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>',
        '  <Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/>',
    ]
    for idx in range(1, len(sheets) + 1):
        content_types_parts.append(
            f'  <Override PartName="/xl/worksheets/sheet{idx}.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>'
        )
    content_types_parts.append('</Types>')
    content_types = "\n".join(content_types_parts)

    rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>'''

    workbook_parts = [
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">',
        '  <sheets>',
    ]
    for idx, (sheet_name, _, _, _, _, _, _, _, _, _, hidden_sheet) in enumerate(sheets, start=1):
        hidden_attr = " state=\"hidden\"" if hidden_sheet else ""
        workbook_parts.append(f'    <sheet name="{escape(sheet_name)}" sheetId="{idx}" r:id="rId{idx}"{hidden_attr}/>')
    workbook_parts.extend([
        '  </sheets>',
        '  <calcPr calcId="1" fullCalcOnLoad="1"/>',
        '</workbook>',
    ])
    workbook = "\n".join(workbook_parts)

    workbook_rels_parts = [
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">',
    ]
    for idx in range(1, len(sheets) + 1):
        workbook_rels_parts.append(
            f'  <Relationship Id="rId{idx}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet{idx}.xml"/>'
        )
    workbook_rels_parts.append(
        f'  <Relationship Id="rId{len(sheets) + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>'
    )
    workbook_rels_parts.append('</Relationships>')
    workbook_rels = "\n".join(workbook_rels_parts)

    styles = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <fonts count="4">
    <font><sz val="11"/><name val="Calibri Light"/></font>
    <font><b/><sz val="11"/><name val="Calibri Light"/></font>
    <font><sz val="9"/><name val="Calibri Light"/></font>
    <font><b/><sz val="9"/><name val="Calibri Light"/></font>
  </fonts>
  <fills count="5">
    <fill><patternFill patternType="none"/></fill>
    <fill><patternFill patternType="gray125"/></fill>
    <fill><patternFill patternType="solid"><fgColor rgb="FFD9E1F2"/><bgColor indexed="64"/></patternFill></fill>
    <fill><patternFill patternType="solid"><fgColor rgb="FFE6E6E6"/><bgColor indexed="64"/></patternFill></fill>
    <fill><patternFill patternType="solid"><fgColor rgb="FFF2F7E6"/><bgColor indexed="64"/></patternFill></fill>
  </fills>
  <borders count="2">
    <border><left/><right/><top/><bottom/><diagonal/></border>
    <border><left style="thin"/><right style="thin"/><top style="thin"/><bottom style="thin"/><diagonal/></border>
  </borders>
  <cellStyleXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="1"/></cellStyleXfs>
  <cellXfs count="23">
    <xf numFmtId="0" fontId="0" fillId="0" borderId="1" xfId="0" applyBorder="1"/>
    <xf numFmtId="0" fontId="1" fillId="2" borderId="1" xfId="0" applyFont="1" applyFill="1" applyBorder="1"/>
    <xf numFmtId="0" fontId="1" fillId="2" borderId="1" xfId="0" applyFont="1" applyFill="1" applyBorder="1"/>
    <xf numFmtId="0" fontId="0" fillId="3" borderId="1" xfId="0" applyFill="1" applyBorder="1"/>
    <xf numFmtId="0" fontId="1" fillId="3" borderId="1" xfId="0" applyFont="1" applyFill="1" applyBorder="1"/>
    <xf numFmtId="0" fontId="0" fillId="0" borderId="1" xfId="0" applyAlignment="1" applyBorder="1"><alignment horizontal="right"/></xf>
    <xf numFmtId="0" fontId="0" fillId="3" borderId="1" xfId="0" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="right"/></xf>
    <xf numFmtId="0" fontId="0" fillId="0" borderId="1" xfId="0" applyAlignment="1" applyBorder="1"><alignment horizontal="right"/></xf>
    <xf numFmtId="0" fontId="0" fillId="3" borderId="1" xfId="0" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="right"/></xf>
    <xf numFmtId="0" fontId="1" fillId="4" borderId="1" xfId="0" applyFont="1" applyFill="1" applyBorder="1"/>
    <xf numFmtId="0" fontId="0" fillId="4" borderId="1" xfId="0" applyFill="1" applyBorder="1"/>
    <xf numFmtId="0" fontId="0" fillId="4" borderId="1" xfId="0" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="right"/></xf>
    <xf numFmtId="0" fontId="0" fillId="4" borderId="1" xfId="0" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="right"/></xf>
    <xf numFmtId="0" fontId="1" fillId="4" borderId="1" xfId="0" applyFont="1" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="right"/></xf>
    <xf numFmtId="0" fontId="0" fillId="4" borderId="1" xfId="0" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="right"/></xf>
    <xf numFmtId="0" fontId="3" fillId="2" borderId="1" xfId="0" applyFont="1" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="center" wrapText="1" vertical="center"/></xf>
    <xf numFmtId="0" fontId="3" fillId="3" borderId="1" xfId="0" applyFont="1" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="center" wrapText="1" vertical="center"/></xf>
    <xf numFmtId="0" fontId="3" fillId="4" borderId="1" xfId="0" applyFont="1" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="center" wrapText="1" vertical="center"/></xf>
    <xf numFmtId="0" fontId="0" fillId="0" borderId="1" xfId="0" applyAlignment="1" applyBorder="1"><alignment horizontal="center" vertical="center"/></xf>
    <xf numFmtId="0" fontId="0" fillId="3" borderId="1" xfId="0" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="center" vertical="center"/></xf>
    <xf numFmtId="0" fontId="0" fillId="4" borderId="1" xfId="0" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="center" vertical="center"/></xf>
    <xf numFmtId="0" fontId="3" fillId="2" borderId="1" xfId="0" applyFont="1" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="center" vertical="center"/></xf>
    <xf numFmtId="0" fontId="3" fillId="4" borderId="1" xfId="0" applyFont="1" applyFill="1" applyAlignment="1" applyBorder="1"><alignment horizontal="center" vertical="center"/></xf>
  </cellXfs>
  <cellStyles count="1"><cellStyle name="Normal" xfId="0" builtinId="0"/></cellStyles>
</styleSheet>'''

    with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("xl/workbook.xml", workbook)
        zf.writestr("xl/_rels/workbook.xml.rels", workbook_rels)
        zf.writestr("xl/styles.xml", styles)

        for idx, (sheet_name, sheet_kind, rows, fieldnames, shaded_columns, enriched_columns, header_multiline, header_height, fixed_widths, hidden_header_columns, _hidden_sheet) in enumerate(sheets, start=1):
            if sheet_kind == "summary":
                xml = _xlsx_sheet_xml_summary(summary_rows)
            else:
                filter_criteria = {"F_FILTER_ALL": ["Y"], "scope_alread_monitored?": ["N"]} if sheet_name == "CANDIDATE" else None
                xml = _xlsx_sheet_xml_table(
                    rows or [],
                    fieldnames or [],
                    shaded_columns=shaded_columns,
                    enriched_columns=enriched_columns,
                    header_multiline=header_multiline,
                    header_height=header_height or 40.0,
                    fixed_widths=fixed_widths,
                    hidden_header_columns=hidden_header_columns,
                    filter_criteria=filter_criteria,
                )
            zf.writestr(f"xl/worksheets/sheet{idx}.xml", xml)


def write_output_json(output_file: str, payload: Dict[str, Any]) -> str:
    output_path = Path(output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    json_text = json.dumps(payload, ensure_ascii=False, indent=2) + "\n"

    with open(output_path, "w", encoding="utf-8") as handle:
        handle.write(json_text)

    gz_path = output_path if output_path.suffix == ".gz" else output_path.with_suffix(output_path.suffix + ".gz")
    with gzip.open(gz_path, "wt", encoding="utf-8") as handle:
        handle.write(json_text)

    output_path.unlink(missing_ok=True)
    return str(gz_path)


def build_impact_params(uid: str, limit: Optional[int], depth_until: Optional[int]) -> Dict[str, Any]:
    params = dict(IMPACT_DEFAULT_PARAMS)
    params["attributeValue"] = uid
    params["limit"] = str(limit if limit is not None else 10000)
    params["depthUntil"] = str(depth_until if depth_until is not None else 8)
    return params


def _filter_rows_from_debug_flags(rows: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Apply filters after enrichment using precomputed debug flags.

    Keep NOT_FOUND/ERROR rows for traceability, and keep FOUND rows only when
    all configured filters passed (`F_FILTER_ALL=Y`).
    """
    out: List[Dict[str, Any]] = []
    for row in rows:
        lookup_status = _normalize_lookup_value(row.get("lookup_status", ""))
        if lookup_status in {"ERROR", "NOT_FOUND"}:
            out.append(row)
            continue
        if str(row.get("F_FILTER_ALL", "")).strip().upper() == "Y":
            out.append(row)
    return out


def _deduplicate_initial_filtered_rows(rows: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Remove exact FILTRED duplicates while preserving business granularity."""
    seen: set[Tuple[str, str, str, str, str, str]] = set()
    deduped: List[Dict[str, Any]] = []
    for row in rows:
        key = (
            _normalize_lookup_value(_get_row_value_by_candidates(row, ["uid"])),
            _normalize_lookup_value(_get_row_value_by_candidates(row, ["program"])),
            _normalize_lookup_value(_get_row_value_by_candidates(row, ["network"])),
            _normalize_lookup_value(_get_row_value_by_candidates(row, ["taken"])),
            _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"])),
            _normalize_lookup_value(_get_row_value_by_candidates(row, ["lookup_status"])),
        )
        if key in seen:
            continue
        seen.add(key)
        deduped.append(row)
    if len(deduped) != len(rows):
        log.info(
            "Initial FILTRED dedupe applied before=%s after=%s removed=%s",
            len(rows),
            len(deduped),
            len(rows) - len(deduped),
        )
    return deduped


def run_impact_analysis(
    # Batch DALI extraction for monitored application UIDs/KEARs
    client: DaliImpactAnalysisClient,
    monitored_rows: List[Dict[str, str]],
    mappings: List[Tuple[str, str]],
    impact_endpoint: str,
    limit: Optional[int],
    depth_until: Optional[int],
    sleep_ms: int,
    dry_run: bool,
    filters: Optional[Dict[str, str]] = None,
    workload_csv: Optional[Path] = None,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]], Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    errors: List[Dict[str, str]] = []
    raw_rows: List[Dict[str, Any]] = []
    filtered_rows_candidates: List[Dict[str, Any]] = []

    deduped_monitored_rows: List[Dict[str, str]] = []
    seen_monitored_contexts: set[Tuple[str, str, str, str]] = set()
    for row in monitored_rows:
        key = (
            _normalize_lookup_value(row.get("uid", "")),
            _normalize_lookup_value(row.get("program", "")),
            _normalize_lookup_value(row.get("network", "")),
            _normalize_lookup_value(row.get("taken", "")),
        )
        if not key[0] or key in seen_monitored_contexts:
            continue
        seen_monitored_contexts.add(key)
        deduped_monitored_rows.append(row)

    total = len(deduped_monitored_rows)
    unique_uids = {str(row.get("uid", "")).strip() for row in deduped_monitored_rows if str(row.get("uid", "")).strip()}
    if total != len(monitored_rows):
        log.info("Monitored rows dedup applied before=%s after=%s", len(monitored_rows), total)
    job_started_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    log.info("Impact analysis batch prepared rows=%s unique_uids=%s", total, len(unique_uids))

    dali_response_cache: Dict[str, Dict[str, Any]] = {}
    dali_error_cache: Dict[str, str] = {}
    raw_uid_seen: set[str] = set()

    for idx, row in enumerate(deduped_monitored_rows, start=1):
        uid = row["uid"]
        log.info("[%s/%s] uid=%s", idx, total, uid)
        response: Dict[str, Any] = {}
        err_text = ""

        if uid in dali_response_cache:
            response = dali_response_cache[uid]
            err_text = dali_error_cache.get(uid, "")
            log.debug("Impact analysis cache hit uid=%s", uid)
        else:
            if dry_run:
                response = {"count": 0, "result": []}
            else:
                try:
                    params = build_impact_params(uid=uid, limit=limit, depth_until=depth_until)
                    response = client.get_json(endpoint=impact_endpoint, params=params)
                except Exception as exc:  # continue batch on error
                    err_text = str(exc)
                    log.warning("Impact analysis failed for uid=%s: %s", uid, err_text)
                    errors.append({"uid": uid, "error": err_text})
                    response = {}

            dali_response_cache[uid] = response
            if err_text:
                dali_error_cache[uid] = err_text

        items.append({"uid": uid, "response": response})

        filtered_base_row = {
            "uid": uid,
            "kear": row.get("kear", uid),
            "program": row.get("program", ""),
            "network": row.get("network", ""),
            "taken": row.get("taken", ""),
            "Server UID": "",
        }
        if uid not in raw_uid_seen:
            raw_base_row = {
                "uid": uid,
                "kear": row.get("kear", uid),
                "Server UID": "",
            }
            raw_rows.extend(
                extract_rows_from_response(response=response, base_row=raw_base_row, mappings=mappings, err_text=err_text, filters=filters, apply_filters=False)
            )
            raw_uid_seen.add(uid)
        filtered_rows_candidates.extend(
            extract_rows_from_response(response=response, base_row=filtered_base_row, mappings=mappings, err_text=err_text, filters=filters, apply_filters=False)
        )

        if sleep_ms > 0:
            time.sleep(sleep_ms / 1000.0)

    if workload_csv is not None:
        enrich_filtered_rows_with_workload_matches(raw_rows, workload_csv)
        enrich_filtered_rows_with_workload_matches(filtered_rows_candidates, workload_csv)

    filtered_rows = _filter_rows_from_debug_flags(filtered_rows_candidates)
    filtered_rows = _deduplicate_initial_filtered_rows(filtered_rows)

    success_count = len(deduped_monitored_rows) - len(errors)
    found_count = sum(1 for item in items if isinstance(item.get("response"), dict) and int(item.get("response", {}).get("count", 0) or 0) > 0)
    job_end_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    payload = {
        "meta": {
            "generated_at": job_end_at,
            "job_started_at": job_started_at,
            "job_end_at": job_end_at,
            "dali_base_url": client.base_url,
            "endpoint": impact_endpoint,
            "uid_count": len(deduped_monitored_rows),
            "success_count": success_count,
            "found_count": found_count,
            "error_count": len(errors),
            "depth_until": depth_until,
            "limit": limit,
            "dry_run": dry_run,
        },
        "items": items,
        "errors": errors,
    }
    return raw_rows, filtered_rows, payload


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Batch DALI impact analysis for monitored UIDs/KEARs.")
    parser.add_argument("--monitored-file", default="user_inputs/monitored_kears.csv", help="Path to monitored_kears.csv")
    parser.add_argument("--input", dest="monitored_file", help="Compatibility alias for --monitored-file")
    parser.add_argument("--headers-file", default="user_inputs/headers.csv", help="Path to headers.csv")
    parser.add_argument("--headers-xlsx", dest="headers_file", help="Compatibility alias for --headers-file")
    parser.add_argument("--headers-sheet", help="Compatibility option for legacy command (ignored in CSV mode)")
    parser.add_argument("--excel", action="store_true", help="Compatibility option for legacy command (ignored in CSV mode)")
    parser.add_argument("--filters-file", default="user_inputs/filters.conf", help="Path to filters.conf (key,value)")
    parser.add_argument(
        "--servers-to-exclude-file",
        default="user_inputs/servers_to_exclude.csv",
        help="Path to servers_to_exclude.csv",
    )
    parser.add_argument("--output", default="RUNS/dali_impact_analysis.xlsx", help="Output XLSX path (sheets RAW and FILTRED)")
    parser.add_argument("--json-out", default="RUNS/dali_impact_analysis.json", help="Output JSON path")
    parser.add_argument("--pptx-output", help="Optional output PPTX path (default: same folder/name as XLSX with .pptx)")
    parser.add_argument(
        "--impact-endpoint",
        default=os.getenv("DALI_IMPACT_ENDPOINT") or "/api/v1/impactAnalysis",
        help="DALI impact analysis endpoint path",
    )
    parser.add_argument("--depth-until", type=int, default=parse_positive_int("DALI_DEPTH_UNTIL", os.getenv("DALI_DEPTH_UNTIL")))
    parser.add_argument("--limit", type=int, default=parse_positive_int("DALI_LIMIT", os.getenv("DALI_LIMIT")))
    parser.add_argument("--sleep-ms", type=int, default=0, help="Sleep between UID calls in milliseconds")
    parser.add_argument("--dry-run", action="store_true", help="Do not call DALI API; generate empty responses")
    parser.add_argument("-v", "--verbose", action="store_true", help="Verbose logging")

    args = parser.parse_args()
    args.depth_until = parse_positive_int("--depth-until", str(args.depth_until) if args.depth_until is not None else None)
    args.limit = parse_positive_int("--limit", str(args.limit) if args.limit is not None else None)
    if args.sleep_ms < 0:
        raise ValueError("--sleep-ms must be >= 0")
    return args


def setup_logging(verbose: bool) -> None:
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(level=level, format="%(asctime)s | %(levelname)s | %(name)s | %(message)s")


def main() -> None:
    """CLI entrypoint: extract, enrich, then write CSV/XLSX/JSON outputs."""
    load_env_file()
    args = parse_args()
    setup_logging(args.verbose)

    if args.excel:
        log.info("--excel compatibility flag detected (CSV mode used).")
    if args.headers_sheet:
        log.info("--headers-sheet=%s ignored in CSV mode.", args.headers_sheet)

    mappings = read_headers_mapping(args.headers_file)
    monitored_rows = read_monitored_kears(args.monitored_file)
    monitored_sheet_headers, monitored_sheet_rows = read_csv_with_original_headers(args.monitored_file)
    filters = read_filters_conf(args.filters_file) if Path(args.filters_file).is_file() else {}
    output_xlsx = Path(args.output)
    output_pptx = Path(args.pptx_output) if args.pptx_output else output_xlsx.with_suffix(".pptx")
    workload_derived_csv = output_xlsx.parent / "export_wkld.derived.csv"

    client = DaliImpactAnalysisClient()
    raw_rows, filtered_rows, json_payload = run_impact_analysis(
        client=client,
        monitored_rows=monitored_rows,
        mappings=mappings,
        impact_endpoint=args.impact_endpoint,
        limit=args.limit,
        depth_until=args.depth_until,
        sleep_ms=args.sleep_ms,
        dry_run=args.dry_run,
        filters=filters,
        workload_csv=workload_derived_csv,
    )
    enrich_rows_with_inventory_for_gen2(raw_rows, filters=filters)
    annotate_raw_scope_programs(raw_rows=raw_rows, monitored_rows=monitored_rows)
    filtered_rows = build_filtered_rows_from_raw(raw_rows=raw_rows, monitored_rows=monitored_rows)
    filtered_rows_for_sheet = [dict(row) for row in filtered_rows]

    monitored_uids = {str(row.get("uid", "")).strip() for row in monitored_rows if str(row.get("uid", "")).strip()}
    raw_server_uids = {
        _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
        for row in raw_rows
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["Server UID", "server_uid", "serveruid"]))
    }
    scope_rows, inv_by_account_rows, _marley_by_ocsname_rows, marley_gen2_by_uuid_rows, marley_rows_for_append = enrich_filtered_rows_with_inventory(
        filtered_rows=filtered_rows,
        client=client,
        impact_endpoint=args.impact_endpoint,
        limit=args.limit,
        depth_until=args.depth_until,
        monitored_uids=monitored_uids,
        raw_server_uids=raw_server_uids,
        filters=filters,
    )

    enrich_filtered_rows_with_workload_matches(scope_rows, workload_derived_csv)
    enrich_marley_rows_with_workload(marley_gen2_by_uuid_rows, workload_derived_csv)
    enrich_marley_rows_with_workload(marley_rows_for_append, workload_derived_csv)
    dict_kear_account_rows = build_dict_kear_account_rows(scope_rows)
    apply_kear_override_from_beneficiary(marley_gen2_by_uuid_rows, dict_kear_account_rows)
    scope_rows = append_marley_rows_to_filtered(
        filtered_rows=scope_rows,
        marley_rows=marley_rows_for_append,
        monitored_rows=monitored_rows,
        inv_by_account_rows=inv_by_account_rows,
        dict_kear_account_rows=dict_kear_account_rows,
    )
    enrich_filtered_rows_with_scope(scope_rows)
    scope_rows = deduplicate_filtered_rows_by_network_iplist(scope_rows)
    servers_to_exclude = read_servers_to_exclude(args.servers_to_exclude_file)
    excluded_rows = apply_manual_exclusions(scope_rows, servers_to_exclude)
    enrich_raw_rows_with_scope_trace(raw_rows=raw_rows, scope_rows=scope_rows)
    annotate_raw_scope_programs(raw_rows=raw_rows, monitored_rows=monitored_rows)
    filtered_rows_for_sheet = build_filtered_rows_from_raw(raw_rows=raw_rows, monitored_rows=monitored_rows)
    raw_csv_path = output_xlsx.with_name(output_xlsx.stem + "_RAW.csv")
    filtered_csv_path = output_xlsx.with_name(output_xlsx.stem + "_FILTRED.csv")
    raw_filter_fieldnames = _raw_filter_fieldnames()
    raw_trace_headers = [name for name in RAW_SCOPE_TRACE_HEADERS if name not in RAW_FILTER_TAIL_HEADERS]
    raw_filter_tail = [name for name in RAW_FILTER_TAIL_HEADERS if name in (raw_filter_fieldnames + RAW_SCOPE_PROGRAM_HEADERS + ["F_Excluded"])]
    raw_extra_fieldnames = (
        raw_trace_headers
        + raw_filter_tail
    )
    enrich_rows = build_enrich_rows_from_marley(
        marley_rows=marley_gen2_by_uuid_rows,
        inv_by_account_rows=inv_by_account_rows,
        dict_kear_account_rows=dict_kear_account_rows,
        mappings=mappings,
        raw_extra_fieldnames=raw_extra_fieldnames,
    )
    for row in enrich_rows:
        row.update(_enrich_filter_columns_from_enrich_row(row=row, filters=filters, servers_to_exclude=servers_to_exclude))
    enrich_filtered_rows_with_scope(enrich_rows)
    annotate_raw_scope_programs(raw_rows=enrich_rows, monitored_rows=monitored_rows)
    populate_enrich_scope_columns_from_monitored(enrich_rows=enrich_rows, monitored_rows=monitored_rows)
    enrich_rows_in_scope = [
        row
        for row in enrich_rows
        if _is_truthy_flag(_get_row_value_by_candidates(row, ["In Scope(s)", "In scope"]))
    ]
    scope_rows_for_sheet = filtered_rows_for_sheet + enrich_rows_in_scope
    expected_scope_rows = len(filtered_rows_for_sheet) + len(enrich_rows_in_scope)
    if len(scope_rows_for_sheet) != expected_scope_rows:
        raise RuntimeError(
            f"SCOPE row count mismatch expected={expected_scope_rows} actual={len(scope_rows_for_sheet)}"
        )
    log.info(
        "SCOPE append count check filtered=%s enrich_total=%s enrich_in_scope=%s scope=%s",
        len(filtered_rows_for_sheet),
        len(enrich_rows),
        len(enrich_rows_in_scope),
        len(scope_rows_for_sheet),
    )
    write_output_csv(str(raw_csv_path), raw_rows, mappings, extra_fieldnames=raw_extra_fieldnames, base_fieldnames=["uid", "Server UID"])
    write_output_csv(str(filtered_csv_path), filtered_rows_for_sheet, mappings, extra_fieldnames=None)

    now_utc = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    started_at = json_payload.get("meta", {}).get("job_started_at", now_utc)
    json_gz_path = write_output_json(args.json_out, json_payload)
    ended_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    error_count = int(json_payload.get("meta", {}).get("error_count", 0) or 0)
    execution_status = "SUCCESS" if error_count == 0 else "FAIL"
    summary_rows: List[Tuple[str, str]] = [
        ("Section 1 : Execution Report", ""),
        ("Report date", now_utc),
        ("Job Started at", started_at),
        ("Job End at", ended_at),
        ("Execution Report", execution_status),
        ("", ""),
        ("Section 2 : Applied filters", ""),
    ]
    if filters:
        for key, value in sorted(filters.items()):
            summary_rows.append((key, value))
    else:
        summary_rows.append(("No filter", "<none>"))

    gen2_rows = [row for row in scope_rows_for_sheet if _normalize_lookup_value(_get_row_value_by_candidates(row, ["cloud_type", "server_cloud_type"])) == "GEN 2"]
    inventory_found_rows = [
        row
        for row in gen2_rows
        if str(row.get("INV_ocs_name", "")).strip() not in {"", "NOT_FOUND", "NOT_GEN2"}
    ]

    summary_rows.extend(
        [
            ("", ""),
            ("Section 3 : Dali Report", ""),
            ("Number of processed kears", str(len(monitored_rows))),
            ("Total assets get from Dali", str(len(raw_rows))),
            ("Total assets after filtering", str(len(filtered_rows_for_sheet))),
            ("", ""),
            ("Section 4 : Data4sec/inventory report", ""),
            ("Number of processed GEN 2 servers", str(len(gen2_rows))),
            ("Number of assets found in inventory", str(len(inventory_found_rows))),
        ]
    )

    candidate_sheet = build_candidate_sheet(workload_csv=workload_derived_csv, monitored_rows=monitored_rows, scope_rows=scope_rows_for_sheet, filters=filters)

    recap_program_sheets = build_program_recap_sheets(
        monitored_rows=monitored_rows,
        filtered_rows=filtered_rows_for_sheet,
        scope_rows=scope_rows_for_sheet,
        raw_rows=raw_rows,
        enrich_rows=enrich_rows,
        output_path=output_xlsx,
    )
    recap_by_name = {name: (name, rows, headers) for name, rows, headers in recap_program_sheets}
    illumio_gap_sheets = build_illumio_gap_sheets(scope_rows=scope_rows_for_sheet, excluded_rows=excluded_rows)
    illumio_by_name = {name: (name, rows, headers) for name, rows, headers in illumio_gap_sheets}
    scope_fieldnames = build_filtered_output_fieldnames(mappings)
    enrich_fieldnames = ["uid", "Server UID"] + [display for display, _ in mappings] + raw_extra_fieldnames
    marley_sheet_preferred = [
        "lookup_uuid",
        "lookup_status",
        "ocs_name",
        "uuid",
        "owner_app_name",
        "beneficiary",
        "app_info.kear_uuid",
        "Kear in scope",
        "KEAR_OVERRIDE",
        "app_info.app_id",
        "app_info.app_name",
        "app_info.env",
        "app_info.kear_library",
        "status",
        "usage",
        "app_info.account_id",
        "app_info.factor",
        "app_info.ref_app",
        "app_info.service_line_name",
        "net_info.net_ipadress",
        "typologie",
        "os_name",
        "os_version",
        "silos",
        "dns",
        "MAIN IP",
        "ILU_managed",
        "ILU_IPLIST",
        "ILU_SUBNET",
        "ILU_enforcement",
        "ILU_role",
        "ILU_app",
        "ILU_env",
        "ILU_loc",
        "ILU_OS",
        "ILU_hostname",
        "ILU_short_hostname",
        "ILU_interfaces",
        "ILU_ip_with_default_gw",
        "ILU_ocs_name_from_IP",
        "ILU_ocs_nam_from_IP",
    ]
    marley_fieldnames = [
        name
        for name in _ordered_fieldnames_with_preferred(marley_gen2_by_uuid_rows, marley_sheet_preferred)
        if not str(name).startswith("F_")
    ]
    filtered_sheet_fieldnames = build_filtered_output_fieldnames(mappings)
    raw_sheet_fieldnames = ["uid", "Server UID"] + [display for display, _ in mappings] + raw_extra_fieldnames
    ordered_sheets: List[Tuple[str, List[Dict[str, Any]], Optional[List[str]]]] = [
        ("get_inv_by_account", inv_by_account_rows, None),
        ("get_marley_gen2_by_uuid", marley_gen2_by_uuid_rows, marley_fieldnames),
        (
            "DictKearAccount",
            dict_kear_account_rows,
                [
                    "INV_Beneficiary_Account",
                    "INV_Beneficiary_Account_ENV",
                    "DALI [APP] UID",
                "DALI [APP] NAME",
                "DALI [APP] SHORT LABEL",
                "DALI [APP] ASA",
                "DALI [APP] IRT CODE",
                "DALI [APP] IAPPLI CODE",
                "DALI [APP] TRIGRAM",
                "DALI [APP] DSI",
                "DALI [APP] APPLICATION MANAGEMENT RC",
                "DALI [APP] APPLICATION DEVELOPMENT MANAGER REL",
            ],
        ),
        ("MONITORED_SCOPES", monitored_sheet_rows, monitored_sheet_headers),
        candidate_sheet,
        ("ENRICH", enrich_rows, enrich_fieldnames),
        ("SCOPE", scope_rows_for_sheet, scope_fieldnames),
    ]
    for recap_name in ("STATS", "TOTAL.PROGRAM", "TOTAL.ENTITY"):
        if recap_name in recap_by_name:
            ordered_sheets.append(recap_by_name[recap_name])
    for gap_name in ("NOT_IN_ILLUMIO", "IN_ILLUMIO_BUT_NOT_BLOCKING"):
        if gap_name in illumio_by_name:
            ordered_sheets.append(illumio_by_name[gap_name])
    ordered_sheets.append(("EXCLUDED", excluded_rows, EXCLUDED_SHEET_HEADERS))

    # Final safety pass on the exact rows that will be written to the XLSX.
    annotate_raw_scope_programs(raw_rows=raw_rows, monitored_rows=monitored_rows)
    annotate_raw_scope_programs(raw_rows=enrich_rows, monitored_rows=monitored_rows)

    raw_blank_in_scope = sum(
        1
        for row in raw_rows
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["F_FILTER_ALL"])) == "Y"
        and str(_get_row_value_by_candidates(row, ["In Scope(s)"])).strip() == ""
    )
    enrich_blank_in_scope = sum(
        1
        for row in enrich_rows
        if _normalize_lookup_value(_get_row_value_by_candidates(row, ["F_FILTER_ALL"])) == "Y"
        and str(_get_row_value_by_candidates(row, ["In Scope(s)"])).strip() == ""
    )
    if raw_blank_in_scope or enrich_blank_in_scope:
        raise RuntimeError(
            f"Final scope annotation incomplete before XLSX write: "
            f"RAW blank In Scope(s)={raw_blank_in_scope}, "
            f"ENRICH blank In Scope(s)={enrich_blank_in_scope}"
        )

    ordered_sheets.append(build_global_sheet(raw_rows=raw_rows, enrich_rows=enrich_rows))
    ordered_sheets.append(build_out_of_scope_sheet(raw_rows=raw_rows, enrich_rows=enrich_rows, dict_kear_account_rows=dict_kear_account_rows))
    ordered_sheets.insert(0, build_dict_sheet(raw_sheet_fieldnames, filtered_sheet_fieldnames, ordered_sheets))

    write_output_xlsx(
        str(output_xlsx),
        raw_rows,
        filtered_rows_for_sheet,
        mappings,
        summary_rows,
        raw_base_fieldnames=["uid", "Server UID"],
        raw_extra_fieldnames=raw_extra_fieldnames,
        filtered_fieldnames_override=filtered_sheet_fieldnames,
        extra_sheets=ordered_sheets,
    )

    if all(name in recap_by_name for name in ("TOTAL.PROGRAM", "TOTAL.ENTITY", "STATS")):
        write_output_pptx(
            str(output_pptx),
            total_program_sheet=recap_by_name["TOTAL.PROGRAM"],
            total_entity_sheet=recap_by_name["TOTAL.ENTITY"],
            stats_sheet=recap_by_name["STATS"],
            monitored_rows=monitored_rows,
        )

    print(f"Monitored rows: {len(monitored_rows)}")
    print(f"Header mappings: {len(mappings)}")
    print(f"Custom filters loaded: {len(filters)}")
    print(f"Impact endpoint: {args.impact_endpoint}")
    print(f"Depth until: {args.depth_until}")
    print(f"Limit: {args.limit}")
    print(f"Errors: {len(json_payload.get('errors', []))}")
    print(f"RAW CSV written to: {raw_csv_path}")
    print(f"FILTRED CSV written to: {filtered_csv_path}")
    print(f"XLSX written to: {output_xlsx}")
    print(f"PPTX written to: {output_pptx}")
    print(f"JSON.GZ written to: {json_gz_path}")


if __name__ == "__main__":
    main()
